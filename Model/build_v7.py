import sys, os
from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import campus_model as M, index_model as I, indiastat as S

THEME = sys.argv[1] if len(sys.argv)>1 else "light"
OUT   = sys.argv[2] if len(sys.argv)>2 else "Flipkart_Minutes_WiRED_Round1"
TPL   = "/sessions/bold-modest-albattani/mnt/Flipkart/Case/Presentation_template.pptx"
IMG   = f"/sessions/bold-modest-albattani/mnt/Flipkart/Model/charts/{THEME}"
FONT  = "Calibri"

if THEME=="light":
    CONTENT_LAYOUT, APPENDIX_LAYOUT = "CUSTOM_3","CUSTOM_2"
    FG=C(0x0D,0x1F,0x5C); MUTE=C(0x5A,0x67,0x85); RULE=C(0xD3,0xDC,0xF0)
    BANNER_BG=C(0x0D,0x1F,0x5C); BANNER_FG=C(0xFF,0xFF,0xFF)
    CHIP=C(0xF2,0xF6,0xFF); CHIP2=C(0xFF,0xF8,0xE3)
    POS=C(0x1E,0x7A,0x46); NEG=C(0xC0,0x39,0x2B); HI=C(0xB5,0x7F,0x00)
    NAVBG=C(0xE8,0xEE,0xFB); NAVON=C(0x0D,0x1F,0x5C); NAVONFG=C(0xFF,0xFF,0xFF)
    BAND=C(0x0D,0x1F,0x5C); BANDFG=C(0xFF,0xFF,0xFF); BANDHI=C(0xFF,0xC2,0x20)
    A_FG=C(0xFF,0xFF,0xFF); A_MUTE=C(0xC3,0xD2,0xF5); A_LINK=C(0xFF,0xC2,0x20)
    A_BANNER_BG=C(0xFF,0xC2,0x20); A_BANNER_FG=C(0x0A,0x1A,0x4E); A_RULE=C(0x3E,0x5C,0xA8)
    A_CHIP=C(0x14,0x2C,0x77); A_CHIP2=C(0x1B,0x36,0x86); A_ZEBRA=C(0x13,0x2A,0x72)
    A_POS=C(0x6B,0xE0,0xA0); A_NEG=C(0xFF,0x9A,0x8B); A_HI=C(0xFF,0xC2,0x20)
    A_HDR=C(0x9F,0xB6,0xE8); A_LINK2=C(0x8F,0xB0,0xFF)
else:
    CONTENT_LAYOUT, APPENDIX_LAYOUT = "CUSTOM_2","CUSTOM_3"
    FG=C(0xFF,0xFF,0xFF); MUTE=C(0xBF,0xCE,0xF0); RULE=C(0x3E,0x5C,0xA8)
    BANNER_BG=C(0xFF,0xC2,0x20); BANNER_FG=C(0x0A,0x1A,0x4E)
    CHIP=C(0x12,0x2B,0x74); CHIP2=C(0x1B,0x36,0x86)
    POS=C(0x6B,0xE0,0xA0); NEG=C(0xFF,0x9A,0x8B); HI=C(0xFF,0xC2,0x20)
    NAVBG=C(0x14,0x2C,0x77); NAVON=C(0xFF,0xC2,0x20); NAVONFG=C(0x0A,0x1A,0x4E)
    BAND=C(0xFF,0xC2,0x20); BANDFG=C(0x0A,0x1A,0x4E); BANDHI=C(0x0A,0x1A,0x4E)
    A_FG=C(0x0D,0x1F,0x5C); A_MUTE=C(0x5A,0x67,0x85); A_LINK=C(0x12,0x4A,0xB8)
    A_BANNER_BG=C(0x0D,0x1F,0x5C); A_BANNER_FG=C(0xFF,0xFF,0xFF); A_RULE=C(0xD3,0xDC,0xF0)
    A_CHIP=C(0xF2,0xF6,0xFF); A_CHIP2=C(0xFF,0xF8,0xE3); A_ZEBRA=C(0xF5,0xF8,0xFF)
    A_POS=C(0x1E,0x7A,0x46); A_NEG=C(0xC0,0x39,0x2B); A_HI=C(0xB5,0x7F,0x00)
    A_HDR=C(0x7C,0x8B,0xAD); A_LINK2=C(0x12,0x4A,0xB8)

prs=Presentation(TPL)
L={l.name:l for l in prs.slide_masters[0].slide_layouts}
_ids=prs.slides._sldIdLst
for sid in list(_ids)[1:]:
    prs.part.drop_rel(sid.rId); _ids.remove(sid)

# ---------------- primitives ----------------
def strip(s):
    for sh in list(s.shapes):
        if sh.is_placeholder: sh._element.getparent().remove(sh._element)

def tb(s,l,t,w,h,anchor=MSO_ANCHOR.TOP):
    x=s.shapes.add_textbox(In(l),In(t),In(w),In(h)); f=x.text_frame
    f.word_wrap=True; f.vertical_anchor=anchor
    f.margin_left=f.margin_right=f.margin_top=f.margin_bottom=0
    return f

def par(f,first=False,before=0,space_after=0,lsp=None):
    p=f.paragraphs[0] if (first and not f.paragraphs[0].runs) else f.add_paragraph()
    if before: p.space_before=Pt(before)
    if space_after: p.space_after=Pt(space_after)
    if lsp: p.line_spacing=lsp
    return p

def run(p,txt,sz,col,bold=False,ital=False):
    r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=bold
    r.font.italic=ital; r.font.name=FONT; r.font.color.rgb=col; return r

def text(s,l,t,w,h,txt,sz,col,bold=False,ital=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,lsp=None):
    f=tb(s,l,t,w,h,anchor); p=par(f,True,lsp=lsp); p.alignment=align
    run(p,txt,sz,col,bold,ital); return f

def box(s,l,t,w,h,fill,radius=0.10,line=None,lw=0.75):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,In(l),In(t),In(w),In(h))
    try: sh.adjustments[0]=radius
    except Exception: pass
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=Pt(lw)
    sh.shadow.inherit=False
    sh.text_frame.word_wrap=True
    sh.text_frame.margin_left=In(0.07); sh.text_frame.margin_right=In(0.07)
    sh.text_frame.margin_top=In(0.03); sh.text_frame.margin_bottom=In(0.03)
    return sh

def rule(s,l,t,w,h,col=None):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,In(l),In(t),In(w),In(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=col or RULE
    sh.line.fill.background(); sh.shadow.inherit=False; return sh

from PIL import Image
def pic(s,name,l,t,maxw,maxh,center=True):
    p=f"{IMG}/{name}.png"
    iw,ih=Image.open(p).size; a=ih/iw
    w=maxw; h=w*a
    if h>maxh: h=maxh; w=h/a
    x=l+(maxw-w)/2 if center else l
    return s.shapes.add_picture(p,In(x),In(t),width=In(w),height=In(h))

STEPS=["01  THE FILTER","02  METRIC SYSTEM  (D1)","03  SEGMENTATION ENGINE  (D2)"]
def nav(s,active):
    x=0.55
    for i,lab in enumerate(STEPS):
        w=[1.30,2.05,2.55][i]
        on=(i==active)
        b=box(s,x,0.26,w,0.255,NAVON if on else NAVBG,radius=0.42)
        f=b.text_frame; p=par(f,True); p.alignment=PP_ALIGN.CENTER
        run(p,lab,7.4,NAVONFG if on else MUTE,bold=on)
        f.vertical_anchor=MSO_ANCHOR.MIDDLE
        if i<2: text(s,x+w+0.055,0.30,0.16,0.2,"›",8.5,MUTE)
        x+=w+0.30

def head(s,title,sub):
    text(s,0.55,0.60,8.12,0.40,title,19,FG,bold=True,lsp=0.92)
    text(s,0.55,1.06,9.0,0.30,sub,8.2,MUTE,ital=True,lsp=0.95)

def banner(s,l,t,w,label):
    b=box(s,l,t,w,0.235,BANNER_BG,radius=0.16)
    f=b.text_frame; f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True); run(p,label,7.6,BANNER_FG,bold=True)
    return t+0.235

def divider(s,x,t,h): rule(s,x,t,0.011,h)

def band(s,t,parts,h=0.40):
    b=box(s,0.55,t,8.90,h,BAND,radius=0.10)
    f=b.text_frame; f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True); p.line_spacing=0.98
    for txt,bold in parts: run(p,txt,8.4,BANDHI if bold else BANDFG,bold=bold)
    return b

def footer(s,txt): text(s,0.55,5.37,8.90,0.16,txt,5.5,MUTE,ital=True)

# =========================================================
# SLIDE 1
# =========================================================
s1=prs.slides.add_slide(L[CONTENT_LAYOUT]); strip(s1); nav(s1,0)
head(s1,"Only 11% of India's 4.46 crore students sit inside a dark store's reach",
     "AISHE 2023-24 census. The remaining 89% are day-scholars and distance-mode students. Campus rules then cut the reachable 11% again.")
DIV=4.98
divider(s1,DIV,1.40,3.42)
y=banner(s1,0.55,1.40,4.15,"WHO IS REACHABLE   ·   AISHE 2023-24 CENSUS")
pic(s1,"funnel",0.55,y+0.12,4.15,2.30)
text(s1,0.55,4.30,4.15,0.46,
     "National hostel occupancy is 56.3%. Of built capacity, 44% stands empty, so residential "
     "density and enrolment are different quantities.",7.4,MUTE,lsp=0.95)

y=banner(s1,5.30,1.40,4.15,"WHAT GATES THEM   ·   CODED PRIMARY CORPUS")
gates=[("PERMITTED 24x7","doorstep delivery observed to 3 AM",POS),
       ("GATE-RESTRICTED","walk-to-gate is the norm; 8 PM cut-offs observed",HI),
       ("BANNED OR REVOCABLE","access withdrawn after misuse; one case gender-differentiated",NEG)]
gy=y+0.12
for lab,desc,col in gates:
    b=box(s1,5.30,gy,4.15,0.335,CHIP,radius=0.14)
    rule(s1,5.30,gy,0.052,0.335,col)
    f=b.text_frame; f.margin_left=In(0.13); f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True); p.line_spacing=0.92
    run(p,lab+"   ",7.3,col,bold=True); run(p,desc,7.1,MUTE)
    gy+=0.395
text(s1,5.30,gy+0.02,4.15,0.18,"WHEN THEY CAN ORDER  ·  ILLUSTRATIVE DEMAND SHAPE",6.9,FG,bold=True)
pic(s1,"curfew",5.30,gy+0.20,4.15,1.28)

band(s1,4.86,[("Minutes runs the highest AOV in Indian quick commerce at Rs750-800. ",False),
   ("Students are its lowest-AOV cohort, so a campus store cannot be underwritten on AOV.",True)],h=0.36)
footer(s1,"AISHE 2023-24 (MoE) Tables 31, 42  ·  AOV benchmarks: Inc42, Aug 2026 (Minutes Rs750-800; Instamart Rs746; Blinkit Rs547; Zepto Rs410)  ·  Access regime coded from a keyword-retrieved public corpus; thematic coding establishes theme salience, and population frequencies cannot be inferred from it.")

# =========================================================
# SLIDE 2
# =========================================================
s2=prs.slides.add_slide(L[CONTENT_LAYOUT]); strip(s2); nav(s2,1)
BE=M.full_breakeven_aov(3.0)
text(s2,0.55,0.60,8.12,0.40,"Breakeven needs Rs528 an order, which a student basket does not reach",19,FG,bold=True,lsp=0.92)
pb=box(s2,0.55,1.03,9.00,0.30,None,radius=0.10,line=POS,lw=1.0)
f=pb.text_frame; f.margin_left=In(0.14); f.vertical_anchor=MSO_ANCHOR.MIDDLE
p=par(f,True)
run(p,"CALIBRATION   ",6.6,POS,bold=True)
run(p,"Rs134.7 revenue  −  Rs93.0 stack  −  Rs12.3 residual  =  ",7.0,FG)
run(p,"Rs29.4 model",7.0,POS,bold=True)
run(p,"      vs      ",7.0,MUTE)
run(p,"Rs29.4 reported",7.0,POS,bold=True)
run(p,"   (Eternal, FY26E)",6.6,MUTE)
banner(s2,0.55,1.42,8.90,"ONE ORDER, THREE STORES   ·   CONTRIBUTION WATERFALL, RUPEES PER ORDER")
pic(s2,"waterfall",0.55,1.70,8.90,1.66)

DIV2=5.06
divider(s2,DIV2,3.48,1.34)
text(s2,0.55,3.46,4.30,0.18,"THE METRIC SYSTEM  ·  EACH GUARDS ONE BAR ABOVE",7.2,FG,bold=True)
mets=[("Peak-to-Trough Demand Ratio",">4.0x","store ops"),
      ("Minimum Viable Hour","staffing rule","store ops"),
      ("Gate-Drop Consolidation Ratio","≥3.0x","last mile"),
      ("Tiered Wastage Rate","<2x ambient","net revenue"),
      ("Term-Weighted CM + Break Runway","≥1.0x","calendar"),
      ("Cohort-Bounded CAC Payback","≤40% of tenure","calendar"),
      ("Committed Revenue Ratio","≥25%","net revenue"),
      ("Calendar-Linked Labour Share","≥50%","store ops")]
cy=3.70
for i,(name,bandv,guards) in enumerate(mets):
    col = 0 if i<4 else 1
    x = 0.55 + col*2.24
    yy = cy + (i%4)*0.272
    b=box(s2,x,yy,2.16,0.255,CHIP if col==0 else CHIP2,radius=0.16)
    f=b.text_frame; f.margin_left=In(0.09); f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True); p.line_spacing=0.88
    run(p,name+"  ",6.4,FG,bold=True); run(p,bandv,6.4,HI,bold=True)
text(s2,0.55,4.76,4.30,0.26,"Capital payback on a dark store runs 56 months. A student is on campus for 36 to 48. Metric 6 exists because the store outlives its own customer base.",6.2,MUTE,lsp=0.94)

text(s2,5.38,3.46,4.07,0.18,"WHERE VIABILITY BREAKS  ·  AND WHAT MOVES IT",7.2,FG,bold=True)
pic(s2,"cliff",5.36,3.66,2.20,1.14)
by2=3.66
for c in (1.0,2.0,3.0,4.0):
    fb2=M.full_breakeven_aov(c)
    onbeat = (c==3.0)
    rb=box(s2,7.60,by2,1.85,0.235,CHIP2 if not onbeat else BANNER_BG,radius=0.18)
    f=rb.text_frame; f.margin_left=In(0.10); f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True)
    run(p,f"{c:.0f}x   ",6.3,(BANNER_FG if onbeat else FG),bold=True)
    run(p,f"Rs{fb2:.0f}",6.7,(BANNER_FG if onbeat else HI),bold=True)
    by2+=0.258
text(s2,7.60,by2+0.02,1.85,0.14,"consolidation → full breakeven",5.6,MUTE,ital=True)
band(s2,4.96,[("Rs528 demands 4.8% contribution margin. Top-cohort stores in mature markets run 4%. ",True),
  ("A campus store has to beat the best of the existing network, which is what D1 is built to detect early.",False)],h=0.34)
footer(s2,"Take rate 19.41%, AOV Rs694 and contribution Rs29.4/order: JM Financial Eternal model (FY26E), built on company data  ·  Throughput ceiling 1,400/day sits inside Blinkit's observed 1,334-1,487 range across nine quarters  ·  Per-order cost stack: published quick-commerce economics, 2026; Rs12.3 of variable cost is unallocated and carried openly.  Store contribution-margin benchmarks (2% mature, 4% top cohort) and the 56-month franchised-store payback from analyst channel checks.  A 3-month pilot instruments the store; industry stores need 6-12 months to mature.  Full derivation on A2.")

# =========================================================
# SLIDE 3
# =========================================================
s3=prs.slides.add_slide(L[CONTENT_LAYOUT]); strip(s3); nav(s3,2)
head(s3,f"The minimum viable catchment is {I.CLUSTER:,.0f} hostel residents in one radius",
     f"Orders per resident is anchored on Blinkit's disclosed 3.6 orders per month, taken at 1.5x for a captive campus. No single Indian campus reaches {I.CLUSTER:,.0f}, so the engine scores clusters. Run on {len(I.df)} states, {len(I.g)} clear the gate.")
DIV3=4.62
divider(s3,DIV3,1.36,3.78)
y=banner(s3,0.55,1.36,3.80,"THE ENGINE   ·   THREE STAGES")

g0=box(s3,0.55,y+0.08,3.80,0.62,CHIP2,radius=0.13,line=NEG,lw=1.0)
f=g0.text_frame; f.margin_left=In(0.10); f.vertical_anchor=MSO_ANCHOR.MIDDLE
p=par(f,True); p.line_spacing=0.90
run(p,"GATE 0  ACCESS REGIME  ",7.0,NEG,bold=True)
run(p,"screens, and corrects\n",6.3,MUTE,ital=True)
run(p,"A curfewed campus generates a truncated order history, so observed data under-states it. Gate 0 re-weights Layer 2 inputs before scoring.\n",6.3,MUTE)
run(p,"Digital access is at parity (1.3pp male-female gap). Physical access is not. The constraint is institutional.",6.3,HI)

def layer(s,t,title,note,rows,tint):
    text(s,0.55,t,3.80,0.16,title,7.0,FG,bold=True)
    text(s,0.55,t+0.155,3.80,0.14,note,6.2,MUTE,ital=True)
    yy=t+0.305
    for lab,w in rows:
        box(s,0.55,yy,3.80,0.205,tint,radius=0.20)
        wmax=max(r[1] for r in rows)
        rule(s,0.55,yy+0.055,0.94+2.10*(w/wmax),0.095,
             C(0x9F,0xB6,0xE8) if THEME=="light" else C(0x37,0x5F,0xC6))
        f=tb(s,0.65,yy+0.030,3.60,0.16); p=par(f,True)
        run(p,lab,6.3,FG,bold=True)
        f2=tb(s,0.62,yy+0.030,3.60,0.16); p2=par(f2,True); p2.alignment=PP_ALIGN.RIGHT
        run(p2,f"{w:.0%}",6.4,FG,bold=True)
        yy+=0.192
    return yy

y2=layer(s3,y+0.78,"LAYER 1  ·  MATHEMATICAL","scored, then tiered x1.0 / x0.7 / x0.4",
   [("Campus concentration  (avg enrolment/college)",0.30),
    ("Addressable base  (hostel residents)",0.30),
    ("Purchasing power  (urban MPCE)",0.25),
    ("Residential intensity  (occupancy %)",0.15)],CHIP)
y3=layer(s3,y2+0.055,"LAYER 2  ·  BEHAVIOURAL","coded occasion index",
   [("Occasion density  (mess-off nights, exams)",0.30),
    ("Habit formation  (weekly order frequency)",0.30),
    ("Category mix  (snacks + term-start durables)",0.25),
    ("Group ordering  (multi-account, one gate)",0.15)],CHIP2)
b=box(s3,0.55,y3+0.045,3.80,0.235,BANNER_BG,radius=0.14)
f=b.text_frame; f.vertical_anchor=MSO_ANCHOR.MIDDLE
p=par(f,True); p.alignment=PP_ALIGN.CENTER
run(p,"PRIORITY  =  BEHAVIOURAL INDEX  x  VIABILITY MULTIPLIER",6.9,BANNER_FG,bold=True)

y=banner(s3,4.94,1.36,4.51,f"THE ENGINE, RUN LIVE   ·   {len(I.g)} STATES CLEARING THE SCALE GATE")
pic(s3,"scatter",4.94,y+0.05,4.51,2.30)
text(s3,4.94,4.00,2.26,0.16,"RANK STABILITY  ·  ±10pp ON EVERY WEIGHT",6.3,FG,bold=True)
pic(s3,"rankrange",4.94,4.15,2.26,0.82)
text(s3,7.28,4.00,2.17,0.16,"VARIABLES THAT DISCRIMINATE  ·  ONE WAS DROPPED",6.3,FG,bold=True)
pic(s3,"discriminate",7.28,4.13,2.17,0.86)

band(s3,4.96,[("Posture by quadrant:  ",False),("build standalone",True),(" (priority)  ·  ",False),
  ("bounded ecosystem bet, kill trigger",True),(" (dense, low spend)  ·  ",False),
  ("serve from the city store",True),(" (affluent, fragmented)  ·  ",False),
  ("do not serve",True),(" (sub-scale)",False)],h=0.30)
footer(s3,f"AISHE 2023-24 (MoE) Tables 42 and 31  ·  HCES 2023-24 (MoSPI) urban MPCE  ·  Scale gate derived from the model on slide 2 at {I.ORD_RES} orders per resident per day (1.5x Blinkit's disclosed 3.6/month)  ·  Digital penetration scored and dropped: urban 15-24 daily internet use runs 90.5-97.3% across candidate states (NSS)  ·  Ranks 3-5 are one tied band.")

# =========================================================
# APPENDIX  ·  three slides
# =========================================================
def ap_slide(title, sub):
    sl=prs.slides.add_slide(L[APPENDIX_LAYOUT]); strip(sl)
    text(sl,0.55,0.40,8.90,0.30,title,17,A_FG,bold=True)
    text(sl,0.55,0.76,8.90,0.20,sub,7.8,A_MUTE,ital=True)
    return sl

def ap_banner(sl,l,t,w,label):
    b=box(sl,l,t,w,0.215,A_BANNER_BG,radius=0.16)
    f=b.text_frame; f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True); run(p,label,7.0,A_BANNER_FG,bold=True)
    return t+0.215

def ap_table(sl,l,t,w,cols,rows,fs=6.3,rh=0.176,hdr=True):
    """cols = [(label, frac_width, align)]  rows = list of tuples (str or (str,bold,colour))"""
    y=t
    if hdr:
        x=l
        for lab,fr,al in cols:
            f=tb(sl,x,y,w*fr,0.15); p=par(f,True)
            p.alignment=PP_ALIGN.RIGHT if al=="r" else PP_ALIGN.LEFT
            run(p,lab,5.9,A_HDR,bold=True)
            x+=w*fr
        y+=0.175
        rule(sl,l,y-0.035,w,0.008,A_RULE)
    for i,rw in enumerate(rows):
        if i%2==0: box(sl,l-0.05,y-0.018,w+0.10,rh,A_ZEBRA,radius=0.04)
        x=l
        for (lab,fr,al),cell in zip(cols,rw):
            txt,bold,col = cell if isinstance(cell,tuple) else (cell,False,A_FG)
            f=tb(sl,x,y,w*fr-0.04,rh); p=par(f,True); p.line_spacing=0.86
            p.alignment=PP_ALIGN.RIGHT if al=="r" else PP_ALIGN.LEFT
            run(p,txt,fs,col,bold=bold)
            x+=w*fr
        y+=rh
    return y

TIER={"T1":A_POS,"T2":A_HI,"D":A_LINK2,"A":A_NEG}
def tg(t): return (t,True,TIER.get(t,A_MUTE))

# ---------------- A1 · METHOD ----------------
a1=ap_slide("Appendix A1  ·  Method, parameter sourcing and limitations",
   "How each number entered the model, what has no source, and what the analysis cannot claim.")
DIVA=4.86
rule(a1,DIVA,1.16,0.011,3.90,A_RULE)

y=ap_banner(a1,0.55,1.16,4.10,"THE PARAMETER LADDER")
steps=[("1  SEARCH",A_POS,"Every document in the corpus is full-text indexed and grepped for the parameter before anything else. Analyst models state values that filings only imply."),
       ("2  DERIVE",A_HI,"Only where the search returns nothing. Derived values are computed from disclosed figures and labelled as derived, never presented as reported."),
       ("3  ASSUME",A_NEG,"Last resort. An assumed value must be sensitivity-tested and its swing reported on-slide.")]
for lab,col,desc in steps:
    b=box(a1,0.55,y+0.10,4.10,0.52,A_CHIP,radius=0.10); rule(a1,0.55,y+0.10,0.05,0.52,col)
    f=b.text_frame; f.margin_left=In(0.12); f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True); p.line_spacing=0.88
    run(p,lab+"   ",6.6,col,bold=True); run(p,desc,6.1,A_MUTE)
    y+=0.60
b=box(a1,0.55,y+0.22,4.10,0.72,A_CHIP2,radius=0.10,line=A_NEG,lw=0.9)
f=b.text_frame; f.margin_left=In(0.12); f.vertical_anchor=MSO_ANCHOR.MIDDLE
p=par(f,True); p.line_spacing=0.90
run(p,"TWO INPUTS HAVE NO SOURCE ANYWHERE\n",6.7,A_NEG,bold=True)
run(p,"Campus AOV has no reported value anywhere, so the model solves for the level at which the store breaks even. "
      "Orders per resident was an assumption until Blinkit's disclosed frequency (3.6 orders per month per user) allowed it to be "
      "anchored: the deck uses 1.5x that rate at full penetration, and reports the 0.12 to 0.25 range it sits inside.",6.2,A_MUTE)
text(a1,0.55,y+1.06,4.10,0.72,
 "LIMITATIONS.  Access-regime and behavioural inputs come from thematic coding of a keyword-retrieved public corpus, which "
 "establishes theme salience and not population frequencies. Cost inputs carry mixed vintages: AISHE 2023-24, HCES 2023-24, "
 "analyst estimates FY26E. The HCES household-type table reads about Rs558 below the headline urban series, a level shift "
 "rather than a ranking difference (rho 0.98), and the index is unchanged under either. State level is the demonstration; "
 "the production model runs at campus level.",
 6.4,A_MUTE,lsp=0.98)

y=ap_banner(a1,5.20,1.16,4.25,"ASSUMPTION TABLE   ·   T1 DISCLOSED   T2 INDUSTRY   D DERIVED   A ASSUMED")
cols=[("Input",0.50,"l"),("Value",0.28,"r"),("Tier",0.10,"r"),("",0.12,"r")]
rows=[
 ("Blinkit revenue take rate FY26E","19.41%",tg("T1"),""),
 ("Blinkit AOV FY26E","Rs694",tg("T1"),""),
 ("Blinkit contribution profit / order","Rs29.4",tg("T1"),""),
 ("Orders/day/store, 9-quarter range","1,334-1,487",tg("T1"),""),
 ("Net working capital","18 days",tg("T1"),""),
 ("ROCE target, management","> 40%",tg("T1"),""),
 ("Higher-ed enrolment / hostel residents","4.46 cr / 49.4 L",tg("T1"),""),
 ("Urban MPCE by state","HCES 2023-24",tg("T1"),""),
 ("Active academic months","8.5",tg("T1"),""),
 ("Last mile per order","Rs42",tg("T2"),""),
 ("Store operations per order","Rs39",tg("T2"),""),
 ("Packaging and support per order","Rs12",tg("T2"),""),
 ("Store fixed cost per month","Rs 5-10 L",tg("T2"),""),
 ("Capex per store","Rs 2.2-2.5 cr",tg("T2"),""),
 ("Minutes AOV / orders per day","Rs750-800 / 1,000-1,100",tg("T2"),""),
 ("Unallocated variable cost","Rs12.3",tg("D"),""),
 ("Calendar surcharge  (12 / 8.5 months)","1.412x",tg("D"),""),
 ("Blinkit orders/month per user","3.6",tg("T1"),""),
 ("Store CM: mature / top cohort","2% / 4% of GOV",tg("T1"),""),
 ("Franchised store payback","56 months",tg("T1"),""),
 ("Orders per resident per day","0.18",tg("D"),""),
]
ap_table(a1,5.20,y+0.14,4.25,cols,rows,fs=6.5,rh=0.205)

# ---------------- A2 · THE MODEL ----------------
a2=ap_slide("Appendix A2  ·  Full derivation, reference only",
   "Slide 2 states the headline and the calibration proof. Everything below is the worked computation behind it, reproducible from the workbook.")
rule(a2,DIVA,1.16,0.011,3.90,A_RULE)

y=ap_banner(a2,0.55,1.16,4.10,"CALIBRATION  ·  REPRODUCE A KNOWN ANSWER FIRST")
b=box(a2,0.55,y+0.12,4.10,0.86,A_CHIP,radius=0.10,line=A_POS,lw=1.0)
f=b.text_frame; f.margin_left=In(0.12); f.vertical_anchor=MSO_ANCHOR.MIDDLE
p=par(f,True); p.line_spacing=0.95
run(p,"Revenue per order  =  19.41% x Rs694  =  Rs134.7\n",6.5,A_FG)
run(p,"Less itemised variable stack  Rs42 + Rs39 + Rs12  =  Rs93.0\n",6.5,A_FG)
run(p,"Less unallocated residual  Rs12.3\n",6.5,A_FG)
run(p,"Model contribution / order  =  Rs29.4      Reported  =  Rs29.4",6.9,A_POS,bold=True)
text(a2,0.55,y+1.06,4.10,0.34,
 "The residual is the gap between the itemised stack and the total variable cost implied by the reported contribution. "
 "It is carried openly rather than distributed across the named lines.",6.1,A_MUTE,lsp=0.94)

yy=ap_banner(a2,0.55,y+1.52,4.10,"BREAKEVEN SURFACE   ·   RUPEES PER ORDER")
cols2=[("Gate-drop consolidation",0.44,"l"),("Contribution breakeven",0.30,"r"),("Full breakeven",0.26,"r")]
rows2=[(f"{c:.0f}x  " + ("one order, one trip" if c==1 else f"{c:.0f} orders, one trip"),
        f"Rs{M.cm_breakeven_aov(c):.0f}",
        (f"Rs{M.full_breakeven_aov(c):.0f}",True,A_HI if c!=3 else A_POS))
       for c in (1.0,2.0,3.0,4.0)]
ap_table(a2,0.55,yy+0.14,4.10,cols2,rows2,fs=6.6,rh=0.215)
text(a2,0.55,yy+1.16,4.10,0.34,
 "Full breakeven is evaluated at 1,400 orders per day, inside the observed range, and carries the 1.412x calendar surcharge "
 "for the 3.5 months the campus is empty.",6.1,A_MUTE,lsp=0.94)

y=ap_banner(a2,5.20,1.16,4.25,"WHAT THE MODEL RETURNS")
cols3=[("Campus AOV",0.24,"l"),("CM / order at 3x",0.30,"r"),("Orders/day needed",0.28,"r"),("Feasible",0.18,"r")]
rows3=[]
for a in (450,500,528,550,600,650):
    cmv=M.cm_per_order(a,3.0); od=M.orders_per_day(a,3.0)
    ok = od<=M.CEILING
    rows3.append((f"Rs{a}",f"Rs{cmv:.1f}",f"{od:,.0f}",("yes",True,A_POS) if ok else ("no",True,A_NEG)))
yv=ap_table(a2,5.20,y+0.13,4.25,cols3,rows3,fs=6.5,rh=0.200)

yv=ap_banner(a2,5.20,yv+0.20,4.25,"ASSET TURN  ·  THE CALENDAR PENALTY")
cols4=[("Store type",0.56,"l"),("NOV (Rs cr)",0.22,"r"),("Asset turn",0.22,"r")]
rows4=[("Minutes city store, 12 months","29.7",("12.6x",True,A_POS)),
       ("Campus store, 8.5 months","19.1",("8.1x",True,A_NEG)),
       ("Industry published band","-",("10.5-12x",False,A_MUTE))]
cols4b=[("Contribution margin benchmark",0.62,"l"),("% of order value",0.38,"r")]
rows4b=[("Campus store at Rs528 breakeven",("4.8%",True,A_NEG)),
        ("Top-cohort stores, mature network",("4.0%",True,A_HI)),
        ("Mature-market stores, typical",("2.0%",False,A_MUTE))]
yv=ap_table(a2,5.20,yv+0.12,4.25,cols4,rows4,fs=6.4,rh=0.188)
yv=ap_table(a2,5.20,yv+0.13,4.25,cols4b,rows4b,fs=6.4,rh=0.188)
text(a2,5.20,yv+0.12,4.25,0.56,
 "SENSITIVITY.  Full breakeven AOV at 3x falls from Rs528 to Rs495 as the take rate rises to 20.71%, and sits "
 "between Rs521 and Rs535 across the observed 1,487 to 1,334 throughput range. Gate-drop consolidation is the "
 "largest single lever: Rs672 at door-drop against Rs510 at 4x, a Rs162 span.",
 6.4,A_MUTE,lsp=0.98)

# ---------------- A3 · ENGINE + SOURCES ----------------
a3=ap_slide("Appendix A3  ·  Index specification, full state table and sources",
   "The scale gate is derived from the model on A2. Weights, normalisation and rank stability are stated so the run is reproducible.")
rule(a3,DIVA,1.16,0.011,3.90,A_RULE)

y=ap_banner(a3,0.55,1.16,4.10,"INDEX SPECIFICATION")
b=box(a3,0.55,y+0.10,4.10,0.44,A_CHIP,radius=0.10)
f=b.text_frame; f.margin_left=In(0.12); f.vertical_anchor=MSO_ANCHOR.MIDDLE
p=par(f,True); p.line_spacing=0.92
run(p,"COI  =  0.40 x norm(campus concentration)  +  0.35 x norm(urban MPCE)  +  0.25 x norm(hostel residents)\n",6.3,A_FG,bold=True)
run(p,"Min-max normalisation applied after the scale gate, so micro-territories cannot dominate on MPCE alone.",6.1,A_MUTE)
gy=y+0.62
spec=[("Scale gate, derived", f"1,400 orders/day / {I.ORD_RES} orders per resident (1.5x Blinkit's disclosed frequency) = {I.CLUSTER:,.0f} residents per cluster; five clusters = {I.GATE:,.0f} to justify a network"),
      ("States clearing the gate", f"{len(I.g)} of {len(I.df)}. The remaining {len(I.df)-len(I.g)} are sub-scale, which answers 'structurally unviable at scale' quantitatively"),
      ("Rank stability", "Tamil Nadu and Delhi hold the top two in every +/-10pp weight scenario. Ranks 3 to 5 rotate among Karnataka, Telangana and Bihar and are reported as one tied band"),
      ("Variable screening", "Digital penetration was tested and dropped. Daily internet use among urban 15-24s has a 4% coefficient of variation across candidate states against 55% for campus concentration, so it cannot separate them. TRAI subscriber density was also rejected: service areas do not map to state boundaries"),
      ("Robustness", f"Concentration ranking 2019-20 against 2023-24 gives Spearman rho {S.rho_conc:.2f} (p = 2.8e-08), so the density axis is structural. Swapping urban MPCE for per-capita NSDP 2024-25 gives rho {S.rho_axis:.2f} with four of the top five unchanged")]
for lab,desc in spec:
    f=tb(a3,0.55,gy,4.10,0.34); p=par(f,True); p.line_spacing=0.92
    run(p,lab+".  ",6.5,A_HI,bold=True); run(p,desc,6.3,A_MUTE)
    gy+=0.425

y=ap_banner(a3,5.20,1.16,4.25,f"FULL RUN  ·  {len(I.g)} STATES CLEARING THE SCALE GATE")
cols5=[("#",0.07,"l"),("State",0.30,"l"),("Conc.",0.16,"r"),("MPCE",0.17,"r"),("Residents",0.19,"r"),("COI",0.11,"r")]
half=(len(I.g)+1)//2
def rws(d):
    return [((f"{int(r.Rank)}",False,A_MUTE), (r.State,r.Rank<=5,A_FG if r.Rank>5 else A_POS),
             f"{int(r.Concentration)}", f"{int(r.MPCE):,}", f"{int(r.Residents):,}",
             f"{r.COI:.3f}") for _,r in d.iterrows()]
ty=ap_table(a3,5.20,y+0.12,2.10,cols5,rws(I.g.iloc[:half]),fs=5.5,rh=0.152)
ap_table(a3,7.42,y+0.12,2.03,cols5,rws(I.g.iloc[half:]),fs=5.5,rh=0.152)

ty=ap_banner(a3,5.20,ty+0.18,4.25,"CODED OCCASION TAXONOMY  ·  LAYER 2 INPUTS")+0.12
tax=[("Mess-off substitution","The institutional calendar publishes mess closure nights in advance, so demand is schedulable."),
     ("Academic-utility emergency","Stationery, print supplies and devices against submission and exam dates."),
     ("Term-start durables","Coolers, blankets, mosquito nets. Rs1,000+ baskets in a two to three week window."),
     ("Late-night, curfew-conditional","Exists only where policy permits. Truncated data where it does not."),
     ("Substitute-scarcity routine","Recurring demand where on-campus retail is weakest, measured inversely.")]
for lab,desc in tax:
    box(a3,5.20,ty,4.25,0.205,A_CHIP,radius=0.16)
    f=tb(a3,5.30,ty+0.030,4.07,0.16); p=par(f,True)
    run(p,lab+"   ",6.2,A_FG,bold=True); run(p,desc,5.9,A_MUTE)
    ty+=0.212

src=("SOURCES.  AISHE 2023-24 (MoE) aishe.gov.in, Tables 31, 42  ·  HCES 2023-24 (MoSPI) mospi.gov.in  ·  "
     "JM Financial, Eternal and Quick Commerce deep-dive models, 2025-26  ·  J.P. Morgan, Eternal, Jul 2026  ·  "
     "Elara Capital, Eternal, Jul 2025  ·  Eternal Ltd earnings calls, Jul 2025 and Jan 2026  ·  Swiggy Ltd investor materials  ·  "
     "Inc42, Flipkart Minutes coverage, Aug 2026  ·  RedSeer dark-store economics  ·  NPCI district UPI statistics  ·  "
     "Published quick-commerce per-order cost structure and dark-store capex benchmarks, 2026.")
mb=box(a3,0.55,4.62,4.10,0.66,None,radius=0.06,line=A_MUTE,lw=0.7)
f=mb.text_frame; f.margin_left=In(0.12); f.vertical_anchor=MSO_ANCHOR.MIDDLE
p=par(f,True); p.line_spacing=0.92
run(p,src,5.7,A_MUTE)

# ---------------- A4 · SOURCES AND CITATIONS ----------------
a4=ap_slide("Appendix A4  ·  Sources and citations",
   "Every figure in the deck traces to one of these. Government census and company disclosure are primary; analyst models are built on company data; industry benchmarks are range-quoted.")
rule(a4,DIVA,1.16,0.011,3.90,A_RULE)

import math as _math
def srcblock(sl,l,t,w,title,items):
    yy=ap_banner(sl,l,t,w,title)+0.10
    for name,detail,url in items:
        nlines=max(1,_math.ceil((len(name)+len(detail))/(w*17.5)))
        h=0.098*nlines
        f=tb(sl,l,yy,w,h+0.14); p=par(f,True); p.line_spacing=0.92
        run(p,name+"  ",6.5,A_FG,bold=True); run(p,detail,6.1,A_MUTE)
        yy+=h
        if url:
            f2=tb(sl,l,yy+0.005,w,0.11); p2=par(f2,True); p2.line_spacing=0.92
            r=p2.add_run(); r.text=url; r.font.size=Pt(5.9); r.font.name=FONT
            r.font.underline=True; r.hyperlink.address=url; r.font.color.rgb=A_LINK
            yy+=0.105
        yy+=0.045
    return yy

y=srcblock(a4,0.55,1.16,4.10,"GOVERNMENT CENSUS AND OFFICIAL STATISTICS",[
 ("AISHE 2023-24, Ministry of Education.","Table 31 hostels, intake and students residing; Table 42 average enrolment per college. Enrolment 4.46 cr, hostel capacity 87.7 L, residents 49.4 L.","https://aishe.gov.in/"),
 ("HCES 2023-24, MoSPI.","Average monthly per capita consumption expenditure, urban, by state.","https://mospi.gov.in/"),
 ("NSS, internet use by age group and state.","Persons 15-24 and 15-29 using the internet, by frequency, gender, residence and device.","https://www.mospi.gov.in/national-sample-survey-nsso"),
 ("NPCI district-wise UPI statistics.","Examined as a district tie-breaker, held outside the index.","https://www.npci.org.in/what-we-do/upi/product-statistics"),
])
y=srcblock(a4,0.55,y+0.12,4.10,"COMPANY DISCLOSURE",[
 ("Eternal Ltd (Blinkit).","Quarterly shareholders' letters and earnings calls, Aug 2024 through Jan 2026. Contribution per order, dark-store network, capex per store, net working capital 18 days, ROCE target above 40%.","https://www.eternal.com/investor-relations/"),
 ("Swiggy Ltd (Instamart).","Quarterly investor materials and earnings calls, Feb 2025 through Jul 2026. Contribution margin as a share of GOV, store count and floor area.","https://www.swiggy.com/investor-relations"),
])

y=srcblock(a4,5.20,1.16,4.25,"ANALYST AND BROKER RESEARCH  ·  VIA BLOOMBERG",[
 ("JM Financial Institutional Securities.","Eternal and Swiggy coverage plus quick-commerce deep dives, 2023-2026. Source of the 19.41% take rate, Rs29.4 contribution per order, the 1,334-1,487 orders per day per store series, 3.6 monthly orders per user, and the n=153 consumer survey.",""),
 ("J.P. Morgan.","Eternal, Swiggy and quick-commerce viability notes, 2022-2026. Source of the illustrative franchised dark-store economics and the 56-month payback.",""),
 ("Elara Capital.","Eternal and Swiggy diet reports and initiating coverage, 2024-2026. Source of the store contribution-margin dispersion, 2% in mature markets against 4% for the top cohort.",""),
])
y=srcblock(a4,5.20,y+0.12,4.25,"INDUSTRY AND MARKET",[
 ("Inc42.","Flipkart Minutes AOV of Rs750-800 and 1,000-1,100 orders per day per store, Aug 2026.","https://inc42.com/"),
 ("RedSeer Strategy Consultants.","Quick-commerce sector economics and the dark-store cost base.","https://redseer.com/articles/the-dark-store-blind-spot-the-part-of-quick-commerce-growth-that-the-topline-doesnt-show/"),
 ("Dark-store unit economics, industry steady state.","Capex per store, asset turn, working-capital days and ROCE benchmarks.","https://officechai.com/learn/the-economics-of-dark-stores-in-india/"),
 ("Quick-commerce per-order cost structure.","Last mile Rs42, store operations Rs39, packaging and support Rs12.","https://businessmodelhub.in/dark-store-business-model/"),
 ("Delivery-partner payout benchmarks, 2026.","Cross-check on the last-mile input.","https://www.pickmywork.com/blinkit-zepto-swiggy-delivery-partner-earnings-2026/"),
])

mb=box(a4,0.55,y+0.14,8.90,0.40,None,radius=0.06,line=A_MUTE,lw=0.7)
f=mb.text_frame; f.margin_left=In(0.12); f.vertical_anchor=MSO_ANCHOR.MIDDLE
p=par(f,True); p.line_spacing=0.92
run(p,"PRIMARY CORPUS.  ",6.6,A_FG,bold=True)
run(p,"Access-regime and behavioural inputs are coded from a keyword-retrieved corpus of public posts by students on Indian campuses. "
      "Thematic coding establishes theme salience; population frequencies cannot be inferred from it. The corpus is not cited by individual post, "
      "and no institution or individual is named anywhere in this deck.",6.2,A_MUTE)

p=f"/sessions/bold-modest-albattani/mnt/Flipkart/{OUT}.pptx"
prs.save(p)

# force the theme hyperlink colour so links stay legible on the appendix background
import zipfile, shutil, re as _re
HLINK = "FFC220" if THEME=="light" else "124AB8"
tmp=p+".tmp"
zin=zipfile.ZipFile(p); zout=zipfile.ZipFile(tmp,"w",zipfile.ZIP_DEFLATED)
for it in zin.infolist():
    data=zin.read(it.filename)
    if it.filename.startswith("ppt/theme/"):
        x=data.decode("utf-8")
        x=_re.sub(r'(<a:hlink>\s*<a:srgbClr val=")[0-9A-Fa-f]{6}', r'\g<1>'+HLINK, x)
        x=_re.sub(r'<a:hlink>\s*<a:sysClr[^/]*/>\s*</a:hlink>',
                  f'<a:hlink><a:srgbClr val="{HLINK}"/></a:hlink>', x)
        data=x.encode("utf-8")
    zout.writestr(it, data)
zin.close(); zout.close(); shutil.move(tmp,p)
print("saved",p)
