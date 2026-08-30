import sys, os
from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import campus_model as M, index_model as I

THEME = sys.argv[1] if len(sys.argv)>1 else "light"
OUT   = sys.argv[2] if len(sys.argv)>2 else "Flipkart_Minutes_WiRED_Round1"
HERE  = os.path.dirname(os.path.abspath(__file__))
ROOT  = os.path.dirname(HERE)
TPL   = os.path.join(ROOT, "Case", "Presentation_template.pptx")
IMG   = os.path.join(HERE, "charts", THEME)
FONT  = "Calibri"

if THEME=="light":
    CONTENT_LAYOUT, APPENDIX_LAYOUT = "CUSTOM_3","CUSTOM_2"
    FG=C(0x0D,0x1F,0x5C); MUTE=C(0x5A,0x67,0x85); RULE=C(0xD3,0xDC,0xF0)
    BANNER_BG=C(0x0D,0x1F,0x5C); BANNER_FG=C(0xFF,0xFF,0xFF)
    CHIP=C(0xF2,0xF6,0xFF); CHIP2=C(0xFF,0xF8,0xE3)
    POS=C(0x1E,0x7A,0x46); NEG=C(0xC0,0x39,0x2B); HI=C(0xB5,0x7F,0x00)
    NAVBG=C(0xE8,0xEE,0xFB); NAVON=C(0x0D,0x1F,0x5C); NAVONFG=C(0xFF,0xFF,0xFF)
    BAND=C(0x0D,0x1F,0x5C); BANDFG=C(0xFF,0xFF,0xFF); BANDHI=C(0xFF,0xC2,0x20)
    A_FG=C(0xFF,0xFF,0xFF); A_MUTE=C(0xC3,0xD2,0xF5); A_LINK=C(0xFF,0xC2,0x20)
else:
    CONTENT_LAYOUT, APPENDIX_LAYOUT = "CUSTOM_2","CUSTOM_3"
    FG=C(0xFF,0xFF,0xFF); MUTE=C(0xBF,0xCE,0xF0); RULE=C(0x3E,0x5C,0xA8)
    BANNER_BG=C(0xFF,0xC2,0x20); BANNER_FG=C(0x0A,0x1A,0x4E)
    CHIP=C(0x12,0x2B,0x74); CHIP2=C(0x1B,0x36,0x86)
    POS=C(0x6B,0xE0,0xA0); NEG=C(0xFF,0x9A,0x8B); HI=C(0xFF,0xC2,0x20)
    NAVBG=C(0x14,0x2C,0x77); NAVON=C(0xFF,0xC2,0x20); NAVONFG=C(0x0A,0x1A,0x4E)
    BAND=C(0xFF,0xC2,0x20); BANDFG=C(0x0A,0x1A,0x4E); BANDHI=C(0x0A,0x1A,0x4E)
    A_FG=C(0x0D,0x1F,0x5C); A_MUTE=C(0x5A,0x67,0x85); A_LINK=C(0x12,0x4A,0xB8)

prs=Presentation(TPL)
L={l.name:l for l in prs.slide_masters[0].slide_layouts}
_ids=prs.slides._sldIdLst
for sid in list(_ids)[1:]:
    prs.part.drop_rel(sid.rId); _ids.remove(sid)

# ---------------- primitives ----------------
def strip(s):
    for sh in list(s.shapes):
        if sh.is_placeholder: sh._element.getparent().remove(sh._element)

def tb(s,l,t,w,h,anchor=MSO_ANCHOR.TOP):
    x=s.shapes.add_textbox(In(l),In(t),In(w),In(h)); f=x.text_frame
    f.word_wrap=True; f.vertical_anchor=anchor
    f.margin_left=f.margin_right=f.margin_top=f.margin_bottom=0
    return f

def par(f,first=False,before=0,space_after=0,lsp=None):
    p=f.paragraphs[0] if (first and not f.paragraphs[0].runs) else f.add_paragraph()
    if before: p.space_before=Pt(before)
    if space_after: p.space_after=Pt(space_after)
    if lsp: p.line_spacing=lsp
    return p

def run(p,txt,sz,col,bold=False,ital=False):
    r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=bold
    r.font.italic=ital; r.font.name=FONT; r.font.color.rgb=col; return r

def text(s,l,t,w,h,txt,sz,col,bold=False,ital=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,lsp=None):
    f=tb(s,l,t,w,h,anchor); p=par(f,True,lsp=lsp); p.alignment=align
    run(p,txt,sz,col,bold,ital); return f

def box(s,l,t,w,h,fill,radius=0.10,line=None,lw=0.75):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,In(l),In(t),In(w),In(h))
    try: sh.adjustments[0]=radius
    except Exception: pass
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=Pt(lw)
    sh.shadow.inherit=False
    sh.text_frame.word_wrap=True
    sh.text_frame.margin_left=In(0.07); sh.text_frame.margin_right=In(0.07)
    sh.text_frame.margin_top=In(0.03); sh.text_frame.margin_bottom=In(0.03)
    return sh

def rule(s,l,t,w,h,col=None):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,In(l),In(t),In(w),In(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=col or RULE
    sh.line.fill.background(); sh.shadow.inherit=False; return sh

from PIL import Image
def pic(s,name,l,t,maxw,maxh,center=True):
    p=f"{IMG}/{name}.png"
    iw,ih=Image.open(p).size; a=ih/iw
    w=maxw; h=w*a
    if h>maxh: h=maxh; w=h/a
    x=l+(maxw-w)/2 if center else l
    return s.shapes.add_picture(p,In(x),In(t),width=In(w),height=In(h))

STEPS=["01  THE FILTER","02  METRIC SYSTEM  (D1)","03  SEGMENTATION ENGINE  (D2)"]
def nav(s,active):
    x=0.55
    for i,lab in enumerate(STEPS):
        w=[1.30,2.05,2.55][i]
        on=(i==active)
        b=box(s,x,0.26,w,0.255,NAVON if on else NAVBG,radius=0.42)
        f=b.text_frame; p=par(f,True); p.alignment=PP_ALIGN.CENTER
        run(p,lab,7.4,NAVONFG if on else MUTE,bold=on)
        f.vertical_anchor=MSO_ANCHOR.MIDDLE
        if i<2: text(s,x+w+0.055,0.30,0.16,0.2,"›",8.5,MUTE)
        x+=w+0.30

def head(s,title,sub):
    text(s,0.55,0.60,8.12,0.40,title,19,FG,bold=True,lsp=0.92)
    text(s,0.55,1.06,9.0,0.30,sub,8.2,MUTE,ital=True,lsp=0.95)

def banner(s,l,t,w,label):
    b=box(s,l,t,w,0.235,BANNER_BG,radius=0.16)
    f=b.text_frame; f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True); run(p,label,7.6,BANNER_FG,bold=True)
    return t+0.235

def divider(s,x,t,h): rule(s,x,t,0.011,h)

def band(s,t,parts,h=0.40):
    b=box(s,0.55,t,8.90,h,BAND,radius=0.10)
    f=b.text_frame; f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True); p.line_spacing=0.98
    for txt,bold in parts: run(p,txt,8.4,BANDHI if bold else BANDFG,bold=bold)
    return b

def footer(s,txt): text(s,0.55,5.37,8.90,0.16,txt,5.5,MUTE,ital=True)

# =========================================================
# SLIDE 1
# =========================================================
s1=prs.slides.add_slide(L[CONTENT_LAYOUT]); strip(s1); nav(s1,0)
head(s1,"Only 11% of India's 4.46 crore students sit inside a dark store's reach",
     "AISHE 2023-24 census. The remaining 89% are day-scholars and distance-mode students. Campus rules then cut the reachable 11% again.")
DIV=4.98
divider(s1,DIV,1.40,3.42)
y=banner(s1,0.55,1.40,4.15,"WHO IS REACHABLE   ·   AISHE 2023-24 CENSUS")
pic(s1,"funnel",0.55,y+0.12,4.15,2.30)
text(s1,0.55,4.30,4.15,0.46,
     "National hostel occupancy is 56.3%. Of built capacity, 44% stands empty, so residential "
     "density and enrolment are different quantities.",7.4,MUTE,lsp=0.95)

y=banner(s1,5.30,1.40,4.15,"WHAT GATES THEM   ·   CODED PRIMARY CORPUS")
gates=[("PERMITTED 24x7","doorstep delivery observed to 3 AM",POS),
       ("GATE-RESTRICTED","walk-to-gate is the norm; 8 PM cut-offs observed",HI),
       ("BANNED OR REVOCABLE","access withdrawn after misuse; one case gender-differentiated",NEG)]
gy=y+0.12
for lab,desc,col in gates:
    b=box(s1,5.30,gy,4.15,0.335,CHIP,radius=0.14)
    rule(s1,5.30,gy,0.052,0.335,col)
    f=b.text_frame; f.margin_left=In(0.13); f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True); p.line_spacing=0.92
    run(p,lab+"   ",7.3,col,bold=True); run(p,desc,7.1,MUTE)
    gy+=0.395
text(s1,5.30,gy+0.02,4.15,0.18,"WHEN THEY CAN ORDER  ·  ILLUSTRATIVE DEMAND SHAPE",6.9,FG,bold=True)
pic(s1,"curfew",5.30,gy+0.20,4.15,1.28)

band(s1,4.86,[("Minutes runs the highest AOV in Indian quick commerce at Rs750-800. ",False),
   ("Students are its lowest-AOV cohort, so a campus store cannot be underwritten on AOV.",True)],h=0.36)
footer(s1,"AISHE 2023-24 (MoE) Tables 31, 42  ·  AOV benchmarks: Inc42, Aug 2026 (Minutes Rs750-800; Instamart Rs746; Blinkit Rs547; Zepto Rs410)  ·  Access regime coded from a keyword-retrieved public corpus; thematic coding establishes salience, not population frequencies.")

# =========================================================
# SLIDE 2
# =========================================================
s2=prs.slides.add_slide(L[CONTENT_LAYOUT]); strip(s2); nav(s2,1)
BE=M.full_breakeven_aov(3.0)
head(s2,"Breakeven needs Rs528 an order, which a student basket does not reach",
     "Calibrated on Blinkit's disclosed economics: at its FY26E basket the model returns Rs29.4 contribution per order against Rs29.4 reported. One input then changes.")
banner(s2,0.55,1.42,8.90,"ONE ORDER, THREE STORES   ·   CONTRIBUTION WATERFALL, RUPEES PER ORDER")
pic(s2,"waterfall",0.55,1.70,8.90,1.66)

DIV2=5.06
divider(s2,DIV2,3.48,1.34)
text(s2,0.55,3.46,4.30,0.18,"THE METRIC SYSTEM  ·  EACH GUARDS ONE BAR ABOVE",7.2,FG,bold=True)
mets=[("Peak-to-Trough Demand Ratio",">4.0x","store ops"),
      ("Minimum Viable Hour","staffing rule","store ops"),
      ("Gate-Drop Consolidation Ratio","≥3.0x","last mile"),
      ("Tiered Wastage Rate","<2x ambient","net revenue"),
      ("Term-Weighted CM + Break Runway","≥1.0x","calendar"),
      ("Cohort-Bounded CAC Payback","≤40% of tenure","calendar"),
      ("Committed Revenue Ratio","≥25%","net revenue"),
      ("Calendar-Linked Labour Share","≥50%","store ops")]
cy=3.70
for i,(name,bandv,guards) in enumerate(mets):
    col = 0 if i<4 else 1
    x = 0.55 + col*2.24
    yy = cy + (i%4)*0.272
    b=box(s2,x,yy,2.16,0.255,CHIP if col==0 else CHIP2,radius=0.16)
    f=b.text_frame; f.margin_left=In(0.09); f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True); p.line_spacing=0.88
    run(p,name+"  ",6.4,FG,bold=True); run(p,bandv,6.4,HI,bold=True)
text(s2,0.55,4.80,4.30,0.16,"Operational (left) and financial (right). All eight run on order, roster and inventory logs Flipkart already holds.",6.2,MUTE)

text(s2,5.38,3.46,4.07,0.18,"WHERE VIABILITY BREAKS  ·  AND WHAT MOVES IT",7.2,FG,bold=True)
pic(s2,"cliff",5.36,3.66,2.20,1.14)
pic(s2,"tornado",7.60,3.66,1.85,1.14)
band(s2,4.96,[("Rs528 is above every student-basket estimate, so standalone viability depends on lifting the basket. ",True),
  ("A 3-month pilot instruments the store and returns a go/no-go read; industry stores need 6-12 months to mature.",False)],h=0.34)
footer(s2,"Take rate 19.41%, AOV Rs694 and contribution Rs29.4/order: JM Financial Eternal model (FY26E), built on company data  ·  Throughput ceiling 1,400/day sits inside Blinkit's observed 1,334-1,487 range across nine quarters  ·  Per-order cost stack: published quick-commerce economics, 2026; Rs12.3 of variable cost is unallocated and carried openly.  Full derivation on appendix A2.")

# =========================================================
# SLIDE 3
# =========================================================
s3=prs.slides.add_slide(L[CONTENT_LAYOUT]); strip(s3); nav(s3,2)
head(s3,f"The minimum viable catchment is {I.CLUSTER:,.0f} hostel residents in one radius",
     f"Breakeven density divided by orders per resident gives the minimum catchment. No single Indian campus reaches it, so the engine scores clusters. Run on {len(I.df)} states, {len(I.g)} clear the scale gate.")
DIV3=4.62
divider(s3,DIV3,1.36,3.78)
y=banner(s3,0.55,1.36,3.80,"THE ENGINE   ·   THREE STAGES")

g0=box(s3,0.55,y+0.08,3.80,0.50,CHIP2,radius=0.13,line=NEG,lw=1.0)
f=g0.text_frame; f.margin_left=In(0.10); f.vertical_anchor=MSO_ANCHOR.MIDDLE
p=par(f,True); p.line_spacing=0.90
run(p,"GATE 0  ACCESS REGIME  ",7.0,NEG,bold=True)
run(p,"screens, and corrects\n",6.3,MUTE,ital=True)
run(p,"A curfewed campus generates a truncated order history, so observed data under-states it. Gate 0 re-weights Layer 2 inputs before scoring.",6.3,MUTE)

def layer(s,t,title,note,rows,tint):
    text(s,0.55,t,3.80,0.16,title,7.0,FG,bold=True)
    text(s,0.55,t+0.155,3.80,0.14,note,6.2,MUTE,ital=True)
    yy=t+0.305
    for lab,w in rows:
        box(s,0.55,yy,3.80,0.205,tint,radius=0.20)
        wmax=max(r[1] for r in rows)
        rule(s,0.55,yy+0.055,0.94+2.10*(w/wmax),0.095,
             C(0x9F,0xB6,0xE8) if THEME=="light" else C(0x37,0x5F,0xC6))
        f=tb(s,0.65,yy+0.030,3.60,0.16); p=par(f,True)
        run(p,lab,6.3,FG,bold=True)
        f2=tb(s,0.62,yy+0.030,3.60,0.16); p2=par(f2,True); p2.alignment=PP_ALIGN.RIGHT
        run(p2,f"{w:.0%}",6.4,FG,bold=True)
        yy+=0.204
    return yy

y2=layer(s3,y+0.66,"LAYER 1  ·  MATHEMATICAL","scored, then tiered x1.0 / x0.7 / x0.4",
   [("Campus concentration  (avg enrolment/college)",0.30),
    ("Addressable base  (hostel residents)",0.30),
    ("Purchasing power  (urban MPCE)",0.25),
    ("Residential intensity  (occupancy %)",0.15)],CHIP)
y3=layer(s3,y2+0.055,"LAYER 2  ·  BEHAVIOURAL","coded occasion index",
   [("Occasion density  (mess-off nights, exams)",0.30),
    ("Habit formation  (weekly order frequency)",0.30),
    ("Category mix  (snacks + term-start durables)",0.25),
    ("Group ordering  (multi-account, one gate)",0.15)],CHIP2)
b=box(s3,0.55,y3+0.045,3.80,0.235,BANNER_BG,radius=0.14)
f=b.text_frame; f.vertical_anchor=MSO_ANCHOR.MIDDLE
p=par(f,True); p.alignment=PP_ALIGN.CENTER
run(p,"PRIORITY  =  BEHAVIOURAL INDEX  x  VIABILITY MULTIPLIER",6.9,BANNER_FG,bold=True)

y=banner(s3,4.94,1.36,4.51,f"THE ENGINE, RUN LIVE   ·   {len(I.g)} STATES CLEARING THE SCALE GATE")
pic(s3,"scatter",4.94,y+0.05,4.51,2.30)
text(s3,4.94,4.02,2.35,0.16,"RANK STABILITY  ·  ±10pp ON EVERY WEIGHT",6.4,FG,bold=True)
pic(s3,"rankrange",4.94,4.18,2.38,0.86)
text(s3,7.46,4.04,1.99,0.72,
     "Tamil Nadu and Delhi hold the top two places in every weight scenario tested. "
     "Ranks 3 to 5 rotate among Karnataka, Telangana and Bihar and should be read as one tied band.",
     6.3,MUTE,lsp=0.94)
text(s3,7.46,4.76,1.99,0.30,
     "State level is the demonstration. Every variable has a campus-level equivalent already named.",
     6.3,HI,ital=True,lsp=0.94)

band(s3,4.96,[("Posture by quadrant:  ",False),("build standalone",True),(" (priority)  ·  ",False),
  ("bounded ecosystem bet, kill trigger",True),(" (dense, low spend)  ·  ",False),
  ("serve from the city store",True),(" (affluent, fragmented)  ·  ",False),
  ("do not serve",True),(" (sub-scale)",False)],h=0.30)
footer(s3,"AISHE 2023-24 (MoE) Table 42 average enrolment per college, Table 31 hostel residents  ·  HCES 2023-24 (MoSPI) urban MPCE  ·  Scale gate derived from the contribution model on slide 2 at an assumed 0.25 orders per resident per day  ·  Behavioural weights from coded primary corpus.")

# =========================================================
# APPENDIX
# =========================================================
s4=prs.slides.add_slide(L[APPENDIX_LAYOUT]); strip(s4)
text(s4,0.60,0.44,8.80,0.34,"Appendix  ·  Sources, method and model",19,A_FG,bold=True)
text(s4,0.60,0.88,8.80,0.20,
     "Index specification, weight-sensitivity runs and the full contribution model are in the accompanying workbook.",
     8.4,A_MUTE,ital=True)
links=[("AISHE 2023-24 — All India Survey on Higher Education (MoE)",
  "Enrolment (4.46 cr), hostel capacity and residents (T-31), average enrolment per college (T-42)","https://aishe.gov.in/"),
 ("HCES 2023-24 — Household Consumption Expenditure Survey (MoSPI)",
  "Urban MPCE by state, used as the purchasing-power axis","https://mospi.gov.in/"),
 ("Eternal Ltd — quarterly shareholders' letters",
  "Blinkit contribution per order (Rs30), net AOV (Rs525), dark-store network","https://www.eternal.com/investor-relations/"),
 ("Swiggy Ltd — quarterly investor materials",
  "Instamart contribution margin as % of GOV, store count and floor area","https://www.swiggy.com/investor-relations"),
 ("Inc42 — quick commerce coverage, Aug 2026",
  "Flipkart Minutes AOV (Rs750-800) and per-store order volume","https://inc42.com/"),
 ("RedSeer — dark-store sector economics",
  "Network-level cost base and the store-maturity ramp","https://redseer.com/articles/the-dark-store-blind-spot-the-part-of-quick-commerce-growth-that-the-topline-doesnt-show/"),
 ("Dark-store unit economics, industry steady-state",
  "Capex per store, asset turn, working-capital days, ROCE benchmarks","https://officechai.com/learn/the-economics-of-dark-stores-in-india/"),
 ("Quick-commerce per-order cost structure",
  "Last mile (Rs42), store operations (Rs39), packaging and support (Rs12)","https://businessmodelhub.in/dark-store-business-model/"),
 ("Delivery-partner payout benchmarks, 2026",
  "Rider payout per short-distance order, cross-checks the last-mile input","https://www.pickmywork.com/blinkit-zepto-swiggy-delivery-partner-earnings-2026/"),
 ("NPCI — district-wise UPI statistics",
  "District-level digital-payment concentration, used only as a tie-breaker","https://www.npci.org.in/what-we-do/upi/product-statistics"),
]
cw=4.20; xs=[0.60,0.60+cw+0.30]; cols=[links[:5],links[5:]]
for cx,items in zip(xs,cols):
    f=tb(s4,cx,1.34,cw,3.5); first=True
    for lab,use,url in items:
        p=par(f,first,before=0 if first else 14); p.line_spacing=0.96
        run(p,lab,8.4,A_FG,bold=True)
        p1=par(f,before=2); p1.line_spacing=0.96
        run(p1,use,7.3,A_MUTE)
        p2=par(f,before=2); p2.line_spacing=0.96
        r=p2.add_run(); r.text=url; r.font.size=Pt(6.4); r.font.name=FONT
        r.font.underline=True; r.hyperlink.address=url
        r.font.color.rgb=A_LINK
        first=False
mb=box(s4,0.60,4.96,8.80,0.44,None,radius=0.06,line=A_MUTE,lw=0.75)
f=mb.text_frame; f.margin_left=In(0.12); f.vertical_anchor=MSO_ANCHOR.MIDDLE
p=par(f,True); p.line_spacing=0.92
run(p,"METHOD.  ",7.0,A_FG,bold=True)
run(p,"Access-regime and behavioural inputs are coded from a keyword-retrieved public corpus; thematic coding "
      "establishes theme salience, not population frequencies. Cost inputs carry mixed vintages and are "
      "range-quoted in the workbook. The blended take rate is derived from disclosed figures, not reported. "
      "Two inputs have no source and are solved for rather than assumed: campus AOV and orders per resident per day.",
      6.6,A_MUTE)

p=os.path.join(ROOT, f"{OUT}.pptx")
prs.save(p)

# force the theme hyperlink colour so links stay legible on the appendix background
import zipfile, shutil, re as _re
HLINK = "FFC220" if THEME=="light" else "124AB8"
tmp=p+".tmp"
zin=zipfile.ZipFile(p); zout=zipfile.ZipFile(tmp,"w",zipfile.ZIP_DEFLATED)
for it in zin.infolist():
    data=zin.read(it.filename)
    if it.filename.startswith("ppt/theme/"):
        x=data.decode("utf-8")
        x=_re.sub(r'(<a:hlink>\s*<a:srgbClr val=")[0-9A-Fa-f]{6}', r'\g<1>'+HLINK, x)
        x=_re.sub(r'<a:hlink>\s*<a:sysClr[^/]*/>\s*</a:hlink>',
                  f'<a:hlink><a:srgbClr val="{HLINK}"/></a:hlink>', x)
        data=x.encode("utf-8")
    zout.writestr(it, data)
zin.close(); zout.close(); shutil.move(tmp,p)
print("saved",p)
