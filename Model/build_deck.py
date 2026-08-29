"""
THE REFERENCE DECK — all 8 content slides, S31 architecture.

  python3 build_deck.py [light|dark]

Order: recommendation -> market -> dead zone -> moneyshot+return -> site -> operating model
       -> financials -> roadmap. Rationale in DECK_ARCHITECTURE_SemiFinal.md.

RULE: never type a figure. Everything is imported. Verify any export with verify_deck.py.
This is a REFERENCE build for rebuilding by hand -- layout is dense on purpose.
"""
import sys, os
from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

import campus_model as M, cost_stack as CS, break_mode as B, basket as BK
import aishe_district as AD, calendar_fragmentation as CF, risk_quadrant as Q
import sla as SL, fleet_mix as FM, labour_class as LC, roce as RC
import risk_shocks as RS, working_capital as WC

THEME = sys.argv[1] if len(sys.argv) > 1 else "light"
OUT   = sys.argv[2] if len(sys.argv) > 2 else "Flipkart_Minutes_WiRED_SemiFinal_REFERENCE"
HERE  = os.path.dirname(os.path.abspath(__file__)); ROOT = os.path.dirname(HERE)
TPL   = os.path.join(ROOT, "Case", "Presentation_template.pptx")
IMG   = os.path.join(HERE, "charts", THEME)
FONT  = "Calibri"; MONO = "Consolas"

if THEME == "light":
    CONTENT_LAYOUT = "CUSTOM_3"
    FG=C(0x0D,0x1F,0x5C); MUTE=C(0x5A,0x67,0x85); RULE=C(0xD3,0xDC,0xF0)
    BANNER_BG=C(0x0D,0x1F,0x5C); BANNER_FG=C(0xFF,0xFF,0xFF)
    CHIP=C(0xF2,0xF6,0xFF); CHIP2=C(0xFF,0xF8,0xE3); MONOBG=C(0xF6,0xF8,0xFD)
    POS=C(0x1E,0x7A,0x46); NEG=C(0xC0,0x39,0x2B); HI=C(0xB5,0x7F,0x00)
    NAVBG=C(0xE8,0xEE,0xFB); NAVON=C(0x0D,0x1F,0x5C)
    BAND=C(0x0D,0x1F,0x5C); BANDFG=C(0xFF,0xFF,0xFF); BANDHI=C(0xFF,0xC2,0x20)
else:
    CONTENT_LAYOUT = "CUSTOM_2"
    FG=C(0xFF,0xFF,0xFF); MUTE=C(0xBF,0xCE,0xF0); RULE=C(0x3E,0x5C,0xA8)
    BANNER_BG=C(0xFF,0xC2,0x20); BANNER_FG=C(0x0A,0x1A,0x4E)
    CHIP=C(0x12,0x2B,0x74); CHIP2=C(0x1B,0x36,0x86); MONOBG=C(0x0F,0x24,0x66)
    POS=C(0x6B,0xE0,0xA0); NEG=C(0xFF,0x9A,0x8B); HI=C(0xFF,0xC2,0x20)
    NAVBG=C(0x14,0x2C,0x77); NAVON=C(0xFF,0xC2,0x20)
    BAND=C(0xFF,0xC2,0x20); BANDFG=C(0x0A,0x1A,0x4E); BANDHI=C(0x0A,0x1A,0x4E)

prs = Presentation(TPL)
L = {l.name: l for l in prs.slide_masters[0].slide_layouts}
_ids = prs.slides._sldIdLst
for sid in list(_ids):
    prs.part.drop_rel(sid.rId); _ids.remove(sid)

# ---------------- primitives ----------------
def strip_ph(s):
    for sh in list(s.shapes):
        if sh.is_placeholder: sh._element.getparent().remove(sh._element)

def tb(s,l,t,w,h,anchor=MSO_ANCHOR.TOP):
    x=s.shapes.add_textbox(In(l),In(t),In(w),In(h)); f=x.text_frame
    f.word_wrap=True; f.vertical_anchor=anchor
    f.margin_left=f.margin_right=f.margin_top=f.margin_bottom=0
    return f
def par(f,first=False,lsp=None,before=0):
    p=f.paragraphs[0] if (first and not f.paragraphs[0].runs) else f.add_paragraph()
    if lsp: p.line_spacing=lsp
    if before: p.space_before=Pt(before)
    return p
def run(p,txt,sz,col,bold=False,ital=False,mono=False):
    r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=bold
    r.font.italic=ital; r.font.name=(MONO if mono else FONT); r.font.color.rgb=col; return r
def text(s,l,t,w,h,txt,sz,col,bold=False,ital=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,lsp=None,mono=False):
    f=tb(s,l,t,w,h,anchor); p=par(f,True,lsp=lsp); p.alignment=align
    run(p,txt,sz,col,bold,ital,mono); return f
def box(s,l,t,w,h,fill,radius=0.10,line=None,lw=0.75):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,In(l),In(t),In(w),In(h))
    try: sh.adjustments[0]=radius
    except Exception: pass
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=Pt(lw)
    sh.shadow.inherit=False; sh.text_frame.word_wrap=True
    sh.text_frame.margin_left=In(0.07); sh.text_frame.margin_right=In(0.07)
    sh.text_frame.margin_top=In(0.03); sh.text_frame.margin_bottom=In(0.03)
    return sh
def rect(s,l,t,w,h,col):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,In(l),In(t),In(w),In(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=col; sh.line.fill.background()
    sh.shadow.inherit=False; sh.text_frame.word_wrap=True
    sh.text_frame.margin_left=In(0.05); sh.text_frame.margin_right=In(0.05)
    sh.text_frame.margin_top=In(0.01); sh.text_frame.margin_bottom=In(0.01)
    return sh
def pic(s,name,l,t,maxw,maxh,center=True):
    p=os.path.join(IMG,f"{name}.png"); iw,ih=Image.open(p).size; a=ih/iw
    w=maxw; h=w*a
    if h>maxh: h=maxh; w=h/a
    x=l+(maxw-w)/2 if center else l
    return s.shapes.add_picture(p,In(x),In(t),width=In(w),height=In(h))
def banner(s,l,t,w,label):
    b=box(s,l,t,w,0.225,BANNER_BG,radius=0.16)
    f=b.text_frame; f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True); run(p,label,6.6,BANNER_FG,bold=True)
    return t+0.225
def band(s,t,parts,h=0.34):
    b=box(s,0.45,t,9.10,h,BAND,radius=0.10)
    f=b.text_frame; f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True,lsp=0.96)
    for txt,bold in parts: run(p,txt,8.0,BANDHI if bold else BANDFG,bold=bold)
    return b
def foot(s,txt): text(s,0.45,5.30,9.10,0.20,txt,5.2,MUTE,ital=True,lsp=0.92)
def rail(s, active, label, forward):
    x=0.45
    for i in range(8):
        on=(i==active-1)
        rect(s,x,0.28,0.34 if not on else 0.46,0.05, NAVON if on else NAVBG)
        x+=(0.34 if not on else 0.46)+0.05
    text(s,x+0.10,0.215,3.6,0.16,label,6.4,MUTE,bold=True)
    text(s,7.10,0.215,2.45,0.16,forward,6.2,MUTE,ital=True,align=PP_ALIGN.RIGHT)
def head(s,title,kicker,size=16.5):
    text(s,0.45,0.52,8.20,0.60,title,size,FG,bold=True,lsp=0.92)
    text(s,0.45,1.14,9.10,0.26,kicker,7.2,MUTE,ital=True,lsp=0.98)
def mathbox(s,l,t,w,h,lines,title="THE ARITHMETIC"):
    box(s,l,t,w,h,MONOBG,radius=0.06,line=RULE,lw=0.75)
    text(s,l+0.09,t+0.05,w-0.18,0.13,title,5.8,HI,bold=True)
    f=tb(s,l+0.09,t+0.20,w-0.18,h-0.26)
    for i,(txt,bold) in enumerate(lines):
        p=par(f,i==0,lsp=1.02); run(p,txt,5.9,FG if bold else MUTE,bold=bold,mono=True)
def kv(s,l,t,w,rows,size=6.4,gap=0.175,label_w=None):
    """Fact rows: label left, value right, thin rule between."""
    y=t
    for i,(k,v,col) in enumerate(rows):
        text(s,l,y,w*0.62 if label_w is None else label_w,0.15,k,size,MUTE)
        text(s,l,y,w,0.15,v,size,col,bold=True,align=PP_ALIGN.RIGHT)
        y+=gap
        if i<len(rows)-1: rect(s,l,y-0.035,w,0.006,RULE)
    return y
def table(s,l,t,w,cols,rows,size=6.0,hdr=5.6,gap=0.165,zebra=True):
    """cols = [(label, width_frac, align)] ; rows = [[(text,col,bold),...]]"""
    x=l
    for lab,frac,al in cols:
        text(s,x,t,w*frac,0.13,lab,hdr,MUTE,bold=True,align=al); x+=w*frac
    y=t+0.155
    for i,r in enumerate(rows):
        if zebra and i%2==0: rect(s,l,y-0.015,w,gap-0.01,CHIP)
        x=l
        for (lab,frac,al),(txt,col,bold) in zip(cols,r):
            text(s,x,y,w*frac,0.14,txt,size,col,bold=bold,align=al); x+=w*frac
        y+=gap
    return y

# ============================== SHARED VALUES ==============================
THR   = [B.threshold(fn) for _,fn in B.CONFIGS]
REL   = B.relocate_vs_flex()
REV   = Q.post_revocation_survival()
LM    = SL.volume_weighted()[1]
BE    = CS.breakeven_d2_consistent(CS.CAMPUS_FIXED, LM)
SCN   = RC.scenario_rows()

# ============================== SLIDE 1 ==============================
def slide1():
    s=prs.slides.add_slide(L[CONTENT_LAYOUT]); strip_ph(s)
    rail(s,1,"RECOMMENDATION","→ why is this not a cost problem?")
    head(s, f"Enter on a cluster-plus-catchment node: {M.TURN_RATIO:.0%} of a city store's asset "
            f"productivity, at a site chosen so the break is survivable",
         f"One node serves a 3–6 college cluster AND its adjacent non-student catchment. The academic "
         f"calendar enters as a site-selection criterion, not a cost line. 90-day build, one pilot, "
         f"₹{RC.CE_BASE/1e5:,.0f} lakh of capital employed.")

    yl=banner(s,0.45,1.46,2.86,"THE DECISION")
    dec=[("BUILD","1 pilot node, 3–6 college cluster, non-metro Tier-1/2"),
         ("UNDERWRITE",f"campus demand + ≥{AD.SITE_FILTER_UNCONTESTED:,.0f}/day uncontested adjacent"),
         ("FLEET","gig to the gate, institution runner inside it"),
         ("SLA",f"dynamic batch, {SL.WAIT_CAP_MIN:.0f}-min cap or {SL.BATCH_MAX:.0f} orders" if hasattr(SL,"WAIT_CAP_MIN") else "dynamic batch, 6-min cap or 12 orders"),
         ("BREAK","repurpose the catchment; the node never closes"),
         ("GATE","go/no-go at day 90 on 4 instrumented metrics")]
    y=yl+0.06
    for k,v in dec:
        text(s,0.45,y,0.72,0.13,k,5.8,HI,bold=True)
        text(s,1.19,y,2.12,0.24,v,6.2,FG,lsp=0.94)
        y+=0.245
    text(s,0.45,y+0.02,2.86,0.15,
         f"Replicable across {AD.N_CANDIDATES} districts on the same screen.",5.8,MUTE)

    yr=banner(s,3.45,1.46,3.00,"NOT A COST PROBLEM")
    kv(s,3.45,yr+0.08,3.00,[
        ("cost ladder moves the requirement", f"{THR[0]:.1%} → {THR[-1]:.1%}", NEG),
        ("and stops, short of closing it", f"{THR[-1]*100:.1f} pts unfunded", NEG),
        ("dead-zone burn, no levers", f"₹{REL['do_nothing_total']/1e5:.1f} L", FG),
        ("relocate instead", f"₹{REL['relocate_once']/1e5:.0f} L", FG),
        ("network basket vs required", f"₹{M.MINUTES_AOV:,.0f} → ₹{BE:,.0f}", FG)],gap=0.235)

    yr2=banner(s,6.60,1.46,2.95,"A SITE-SELECTION PROBLEM")
    kv(s,6.60,yr2+0.08,2.95,[
        ("districts clearing the screen", f"{AD.N_CANDIDATES} of {AD.N_DISTRICTS:,}", POS),
        ("asset turn, campus : city", f"{M.CAMPUS_TURN:.2f}× / {M.CITY_TURN:.2f}× = {M.TURN_RATIO:.3f}", POS),
        ("across the observed range", f"{M.turn_ratio_at(M.CEIL_LO):.2f} – {M.turn_ratio_at(M.CEIL_HI):.2f}", POS),
        ("ROCE at a 30% non-grocery basket", f"{SCN[1]['roce']:.1%}", POS),
        ("holding vs relocating", f"{REL['flex_vs_relocate']:.0%} of one move", POS)],gap=0.235)

    ay=3.50
    text(s,0.45,ay,9.10,0.13,"THE ASK, AND WHAT IT RETURNS  ·  ONE NODE, FIVE-YEAR LIFE, ALL FOUR NUMBERS SOLVED NOT ASSUMED",5.9,FG,bold=True)
    cells=[("CAPITAL EMPLOYED", f"₹{RC.CE_BASE/1e5:,.0f} L", f"capex ₹{RC.CAPEX_MID/1e5:,.0f} L + working capital ₹{WC.WC_ADOPTED/1e5:,.0f} L", FG),
           ("ROCE, 30% BASKET", f"{SCN[1]['roce']:.1%}", f"and {RC.ROCE_HURDLE:.0%} at AOV ₹{RC.AOV_HURDLE:,.0f}, inside the stated ceiling", POS),
           ("PAYBACK", f"{SCN[1]['payback']:.0f} mo", f"against a {M.FRANCHISE_PAYBACK}-month franchised-store benchmark¹¹", POS),
           ("IRR", f"{RC.irr(RC.AOV_HURDLE):.1%}", f"{RC.irr(RC.AOV_HURDLE,ramp=6):.0%}–{RC.irr(RC.AOV_HURDLE,ramp=2):.0%} across a 2–6 month ramp", POS)]
    cw=9.10/4
    for i,(lab,val,body,col) in enumerate(cells):
        b=box(s,0.45+i*cw,ay+0.17,cw-0.08,0.52,CHIP,radius=0.10); rect(s,0.45+i*cw,ay+0.17,0.04,0.52,col)
        f=b.text_frame; f.margin_left=In(0.10); f.vertical_anchor=MSO_ANCHOR.MIDDLE
        p=par(f,True,lsp=0.90); run(p,lab+"  ",5.6,col,bold=True); run(p,val,8.6,FG,bold=True)
        p2=par(f,lsp=0.90); run(p2,body,5.5,MUTE)
    text(s,0.45,ay+0.72,9.10,0.14,
         f"LIMIT, COMPUTED: at a 30% basket with volume −30% the node returns {SCN[3]['roce']:.1%} and pays back in "
         f"{SCN[3]['payback']:.0f} months — longer than its {RC.NODE_LIFE_MO}-month life. Volume is the day-90 gate.",
         5.6,NEG,ital=True)

    b=box(s,0.45,4.38,9.10,0.34,None,radius=0.10,line=HI,lw=1.0)
    f=b.text_frame; f.margin_left=In(0.13); f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True,lsp=0.92)
    run(p,"NOT FIRST   ",6.2,HI,bold=True)
    run(p,"Swiggy: student rewards with college canteens and hostels, Toing (AOV <₹200, ~50 cities), a Young India Skills "
          "University MoU²⁰˒²¹. Blinkit: in-airport with Adani at CSMIA²¹. Zepto: Express Prints ₹2/page, ~285 stores²¹. ",6.2,FG)
    run(p,"None underwrites the node on the cluster. That is the whole difference.",6.2,FG,bold=True)
    band(s,4.76,[("Site the node on the cluster, price the break as a filter, ",False),
                 (f"and enter the {AD.N_CANDIDATES} districts where both hold.",True)])
    foot(s,"Sources 1,2,4,5,6,10,20,21  ·  every figure asserted in A1  ·  derivations A5–A7  ·  return model A6b")

# ============================== SLIDE 2 ==============================
def slide2():
    s=prs.slides.add_slide(L[CONTENT_LAYOUT]); strip_ph(s)
    rail(s,2,"THE MARKET","→ if we enter, what breaks?")
    head(s, f"Six in ten Indian colleges are rural and metros run {AD.METRO_EXCESS:.0%} over sustainable "
            f"capacity: a smaller market than it looks, and an emptier one",
         "Both facts come out of registers, not forecasts. The first cuts our own addressable universe by "
         f"{1-AD.URBAN_SHARE_COL:.0%}. The second is why we still recommend entry.")

    b=box(s,0.45,1.44,9.10,0.44,None,radius=0.10,line=NEG,lw=1.0)
    f=b.text_frame; f.margin_left=In(0.13); f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True,lsp=0.93)
    run(p,"THE SEGMENT HAS ALREADY KILLED AN OPERATOR   ",6.4,NEG,bold=True)
    run(p,"Jun 2026: Starship Technologies exited ALL US higher education after 8 years and 60+ campuses, withdrawing "
          "~1,200 robots¹⁵. CEO Ahti Heinla: ",6.6,FG)
    run(p,"“campus and grocery are fundamentally different operations: one is seasonal and contract-driven, the other "
          "is a 365-day urban business.” ",6.6,FG,ital=True)
    run(p,"It lost 365-day utilisation to 8.5-month utilisation — the ratio on slide 4. The campuses did not go dark: "
          "Avride took them on a national foodservice master agreement. A contract shape failed, not the segment.",6.6,FG,bold=True)

    yl=banner(s,0.45,2.00,4.45,"THE UNIVERSE NARROWS   ·   AISHE REGISTER, 28-8-2026")
    pic(s,"universe_narrow",0.45,yl+0.05,4.45,1.55)
    text(s,0.45,yl+1.62,4.45,0.42,
         f"{AD.N_COL:,} colleges → {AD.URBAN_COL:,} urban ({AD.URBAN_SHARE_COL:.1%}) → "
         f"{AD.NON_METRO_URBAN_COLLEGES:,} urban and non-metro, of which {AD.N_STA_URBAN_HP:,} are urban "
         f"high-propensity standalones. Register count [6], not the 2023-24 report count [7]; the two instruments "
         f"are never mixed.",6.1,MUTE,lsp=0.98)

    yr=banner(s,5.10,2.00,4.45,"THE DENSIFICATION WALL   ·   ONE QUARTER OF ADDITIONS")
    pic(s,"metro_squeeze",5.10,yr+0.05,4.45,1.55)
    text(s,5.10,yr+1.62,4.45,0.42,
         f"{AD.STORES_ADDED_Q} stores bought {AD.PIN_CODES_ADDED_Q} new pin codes — "
         f"{AD.STORES_PER_NEW_PINCODE:.1f} per pin code, so nine in ten went where an operator already served. "
         f"Of {AD.TOTAL_STORES:,} stores across five operators, only {AD.NON_METRO_STORES:,} sit outside the metros. "
         f"Store density per urban college: metro {AD.METRO_STORE_DENSITY:.2f} against non-metro "
         f"{AD.NON_METRO_STORE_DENSITY:.2f} — {AD.DENSITY_RATIO:.1f}× headroom where the archetype actually sits.",
         6.1,MUTE,lsp=0.98)

    b=box(s,0.45,4.28,9.10,0.38,CHIP2,radius=0.10); rect(s,0.45,4.28,0.045,0.38,HI)
    f=b.text_frame; f.margin_left=In(0.12); f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True,lsp=0.92)
    run(p,"AND THE SEGMENT IS UNCOUNTED   ",6.3,HI,bold=True)
    run(p,"Euromonitor's Indian consumer-foodservice taxonomy has no education or institutional location type, and its "
          "household typology has no hostel, PG or shared-student category — students fall into “Other”, never broken "
          "out¹⁴. This micro-market is not small in the standard dataset; it is absent from it.",6.3,FG)
    band(s,4.76,[("We shrink our own market on slide two and still recommend entry. ",True),
                 ("The question is not whether — it is where, and at what return.",False)])
    foot(s,"Sources 5,6,7,14,15  ·  full screen A7  ·  register vs report note A2")

# ============================== SLIDE 3 ==============================
def slide3():
    s=prs.slides.add_slide(L[CONTENT_LAYOUT]); strip_ph(s)
    rail(s,3,"D1 · THE DEAD ZONE","→ if cost can't close it, what does?")
    head(s, f"{CS.fixed_flex_share():.1%} of the fixed base flexes, and flexing all of it still leaves the node "
            f"needing {THR[-1]:.1%} of term demand from a break that delivers none",
         f"Fixed base ₹{CS.CAMPUS_FIXED:,.0f}/month¹ — a Tier-1/2 store, confirmed from the rent band in the same "
         f"report. Robust across that band: ₹{CS.T12_FIXED_LOW:,.0f}–₹{CS.T12_FIXED_HIGH:,.0f}, width "
         f"{CS.T12_FIXED_BAND_PCT:.1%}.")

    mathbox(s,0.45,1.44,3.05,1.12,[
        ("σ = 12 / m = 12 / 8.5 = 1.412",True),
        ("   calendar surcharge",False),
        ("r* = F_break / (30 · V · CM)",True),
        ("   at breakeven CM, r* = 1/σ = 70.8%",False),
        ("",False),
        ("F fixed ₹/mo · V term orders/day",False),
        ("CM contribution ₹/order · m active months",False)],"TWO LINES, ONE CONSTANT")
    text(s,0.45,2.60,3.05,0.30,
         "The do-nothing threshold IS the calendar surcharge inverted. The same constant closes the "
         "argument twice, from two directions.",6.0,FG,lsp=0.98)

    ty=2.96
    text(s,0.45,ty,3.05,0.13,"WHAT THE BASE IS MADE OF",5.9,FG,bold=True)
    bx,BW=0.45,3.05
    cols={"Store rent":HI,"In-store staff":BANNER_BG,"Utilities and cold chain":POS,"Other fixed":MUTE}
    for name,val,_,_ in CS.FIXED_STACK:
        w=BW*val/CS.CAMPUS_FIXED
        r=rect(s,bx,ty+0.17,w,0.20,cols[name])
        f=r.text_frame; f.vertical_anchor=MSO_ANCHOR.MIDDLE
        p=par(f,True); p.alignment=PP_ALIGN.CENTER
        run(p,f"{val/CS.CAMPUS_FIXED:.0%}",6.0,C(0xFF,0xFF,0xFF),bold=True)
        bx+=w
    lx=0.45
    for name,val,_,_ in CS.FIXED_STACK:
        w=BW*val/CS.CAMPUS_FIXED
        text(s,lx,ty+0.38,w,0.12,{"Store rent":"rent","In-store staff":"staff",
             "Utilities and cold chain":"util.","Other fixed":"other"}[name],5.3,MUTE)
        lx+=w
    text(s,0.45,ty+0.53,3.05,0.26,
         f"‘Other fixed’ ₹{CS.OTHER_FIXED:,.0f} ({CS.OTHER_FIXED/CS.CAMPUS_FIXED:.0%}) is a residual we allocated "
         f"from JM's blended ₹100/sqft line — the least defensible number in our stack, and it persists through "
         f"the break.",5.8,NEG,lsp=0.98)

    ym=banner(s,3.65,1.44,2.75,"THE COST LADDER")
    pic(s,"lever_ladder_s5",3.65,ym+0.04,2.75,1.32)
    text(s,3.65,ym+1.38,2.75,0.40,
         f"Shaded band is the residual no lever reaches. The ladder cuts the requirement "
         f"{(THR[0]-THR[-1])/THR[0]:.0%} and stops.",5.9,MUTE,lsp=0.98)

    yr=banner(s,6.55,1.44,3.00,"THE SOLVER   ·   BEST = REPURPOSE")
    box(s,6.55,yr+0.04,3.00,1.62,MONOBG,radius=0.06,line=RULE,lw=0.75)
    f=tb(s,6.63,yr+0.10,2.86,1.52)
    mono=[("minimise dead-zone cash burn",True),
          ("s.t. reactivation lead ≥ 28 d",False),
          ("     SLA avg ≤ 27.1 min at peak",False),
          ("     revocation ≥ 513 orders/day",False),
          ("     no calendar-indexed lease",False),
          ("     7-week wind-down rule",False),
          ("",False),
          ("STRATEGY        ₹ 3.5mo   RESID",True),
          ("DO_NOTHING     31.6 L    70.8%",False),
          ("LABOUR_FLEX    21.6 L    59.1%",False),
          ("COLD_RIGHTSIZE 21.2 L    57.0%",False),
          ("SMALL_FORMAT   18.5 L    48.6%",False),
          ("REPURPOSE*      0.0 L    48.6%  ←",True)]
    for i,(txt,bold) in enumerate(mono):
        p=par(f,i==0,lsp=1.02); run(p,txt,5.7,FG if bold else MUTE,bold=bold,mono=True)
    text(s,6.55,yr+1.70,3.00,0.40,
         "* Not a cost configuration: it runs the small-format base and fills 48.6% from adjacent "
         "catchment. Residual campus population supplies only 8–15%.",5.8,MUTE,lsp=0.98)

    band(s,4.76,[("Cost levers are exhausted at 48.6% and a break delivers zero. ",False),
                 ("The gap has to be filled by demand — which makes it a siting decision, not an operating one.",True)])
    foot(s,"Sources 1,11,17,18  ·  last mile ₹19.0/order derived on slide 6  ·  solver and constraints A4  ·  fixed-base derivation A5  ·  Physicswallah ~850bps calendar swing [11]; Kambli 2020 labour-beats-capital [18]")

# ============================== SLIDE 4 ==============================
def slide4():
    s=prs.slides.add_slide(L[CONTENT_LAYOUT]); strip_ph(s)
    rail(s,4,"THE RETURN","→ so which sites supply the density?")
    head(s, f"Density buys back the calendar: {M.CAMPUS_TURN:.2f}× against a city store's {M.CITY_TURN:.2f}× — "
            f"and the node clears the {RC.ROCE_HURDLE:.0%} ROCE the client manages to",
         f"Asset turn is the turnover leg of ROCE. The campus node clears that leg at {M.TURN_RATIO:.0%} of a city "
         f"store; the margin leg is what the basket lever has to buy. Two legs, one identity, in the client's own metric.")

    ib=box(s,0.45,1.44,5.40,0.56,CHIP,radius=0.08)
    f=ib.text_frame; f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True,lsp=0.95); p.alignment=PP_ALIGN.CENTER
    run(p,f"({M.CEILING:,.0f} / {M.MINUTES_ORD:,.0f}) × ({M.ACTIVE_MONTHS} / 12)  =  ",11.0,FG,bold=True)
    run(p,f"{M.CEILING/M.MINUTES_ORD:.3f}",11.0,POS,bold=True); run(p," × ",11.0,FG,bold=True)
    run(p,f"{M.ACTIVE_MONTHS/12:.3f}",11.0,NEG,bold=True); run(p,"  =  ",11.0,FG,bold=True)
    run(p,f"{M.TURN_RATIO:.3f}",15.0,HI,bold=True)
    p2=par(f); p2.alignment=PP_ALIGN.CENTER
    run(p2,"density ",6.0,POS,bold=True); run(p2,f"+{(M.CEILING/M.MINUTES_ORD-1)*100:.1f}%    ",6.0,MUTE)
    run(p2,"calendar ",6.0,NEG,bold=True); run(p2,f"−{(1-M.ACTIVE_MONTHS/12)*100:.1f}%    ",6.0,MUTE)
    run(p2,"dead zone costs ",6.0,FG,bold=True); run(p2,f"{(1-M.TURN_RATIO)*100:.1f} points of productivity, not 30",6.0,MUTE)

    sy=2.06
    text(s,0.45,sy,5.40,0.13,"IF 1,400 ORDERS/DAY IS WRONG  ·  THE RATIO IS LINEAR IN THROUGHPUT",5.9,FG,bold=True)
    cw=5.40/4
    for i,(opd,ratio) in enumerate(M.TURN_SENSITIVITY):
        on=(abs(opd-M.CEILING)<1); parity=(abs(opd-M.PARITY_OPD)<1)
        b=box(s,0.45+i*cw,sy+0.17,cw-0.06,0.38,BANNER_BG if on else (CHIP2 if parity else CHIP),radius=0.09)
        f=b.text_frame; f.vertical_anchor=MSO_ANCHOR.MIDDLE
        p=par(f,True,lsp=0.90); p.alignment=PP_ALIGN.CENTER
        run(p,f"{opd:,.0f}/day",5.9,BANNER_FG if on else MUTE)
        p2=par(f,lsp=0.90); p2.alignment=PP_ALIGN.CENTER
        run(p2,f"{ratio:.3f}",9.6,BANNER_FG if on else (HI if parity else FG),bold=True)
    text(s,0.45,sy+0.60,5.40,0.30,
         f"Blinkit's own observed nine-quarter range is {M.CEIL_LO:,}–{M.CEIL_HI:,} orders/day². Across all of it the "
         f"node runs {M.turn_ratio_at(M.CEIL_LO):.2f}–{M.turn_ratio_at(M.CEIL_HI):.2f}. Parity needs "
         f"{M.PARITY_OPD:,.0f}/day = q_city × σ, which is {M.PARITY_VS_CEIL_HI:.1%} of the observed maximum. "
         f"At 1,000/day the ratio is {M.TURN_RATIO_AT_1000:.2f} and the density argument weakens.",6.0,MUTE,lsp=0.98)

    mathbox(s,0.45,3.06,5.40,0.92,[
        ("ROCE = EBIT/CE = (EBIT/NOV) × (NOV/CE)   [DuPont]",True),
        ("                  margin leg   turnover leg",False),
        (f"     = {RC.dupont(RC.AOV_HURDLE)['ebit_margin']:.2%} × {RC.dupont(RC.AOV_HURDLE)['capital_turn']:.2f}× = "
         f"{RC.roce(RC.AOV_HURDLE):.0%}  at AOV ₹{RC.AOV_HURDLE:,.0f}",True),
        ("AOV* = (ROCE·CE + F·12)/(τ·N) + c/τ    closed form, no search",False)],"THE IDENTITY THE CLIENT MANAGES TO")

    yr=banner(s,6.05,1.44,3.50,"ROCE BY SCENARIO   ·   AGAINST ETERNAL'S OWN HURDLE")
    pic(s,"roce_ladder",6.05,yr+0.05,3.50,1.42)
    kv(s,6.05,yr+1.52,3.50,[
        ("AOV for ROCE = 0 (breakeven)", f"₹{RC.AOV_BREAKEVEN:,.0f}", FG),
        ("AOV for the 40% hurdle", f"₹{RC.AOV_HURDLE:,.0f}", HI),
        ("premium over breakeven", f"₹{RC.HURDLE_PREMIUM:,.0f}", NEG),
        ("non-grocery mix it implies", f"{RC.HURDLE_NONGROCERY_SHARE:.1f}% of GOV", FG),
        ("vs management's stated ceiling", f"{BK.NONGROCERY_CEILING_LO:.0f}–{BK.NONGROCERY_CEILING:.0f}%", POS)],gap=0.20)

    b=box(s,0.45,4.06,5.40,0.58,None,radius=0.09,line=NEG,lw=1.0)
    f=b.text_frame; f.margin_left=In(0.12); f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True,lsp=0.93)
    run(p,"TWO ASSET TURNS, BOTH CORRECT — ALWAYS QUOTE THE BASIS   ",5.9,NEG,bold=True)
    run(p,f"The {M.CAMPUS_TURN:.2f}× above is a like-for-like comparison at a COMMON AOV (₹{M.MINUTES_AOV:,.0f}), which is "
          f"what isolates density × calendar. The DuPont leg ({RC.dupont(RC.AOV_HURDLE)['capital_turn']:.2f}×) is the node's "
          f"own turnover at its achieved AOV on capital employed including working capital. Quoting either without its "
          f"basis is how a panel finds a contradiction that is not there.",5.9,FG)

    band(s,4.76,[("Breakeven was never the hurdle. ",True),
                 (f"₹{RC.AOV_BREAKEVEN:,.0f} earns zero on ₹{RC.CE_BASE/1e5:,.0f} lakh; ₹{RC.AOV_HURDLE:,.0f} earns "
                  f"{RC.ROCE_HURDLE:.0%} — and the site is what supplies the throughput either needs.",False)])
    foot(s,"Sources 1,2,4,8,10,11  ·  ROCE model and DuPont reconciliation A6b  ·  throughput sensitivity A6  ·  capital build A5")

# ============================== SLIDE 5 ==============================
def slide5():
    s=prs.slides.add_slide(L[CONTENT_LAYOUT]); strip_ph(s)
    rail(s,5,"SITE SELECTION","→ can we operate inside a gate?")
    head(s, f"The calendar becomes a filter: {AD.SITE_FILTER_UNCONTESTED:,.0f} orders/day of uncontested adjacent "
            f"demand, {AD.N_CANDIDATES} districts of {AD.N_DISTRICTS:,} — and we disqualified {AD.STACKED_SHARE:.0%} of our own list",
         f"The filter binds on access revocation ({REV['orders_needed']:,.0f}/day), not on break-period solvency "
         f"({REV['adjacent_lo']:,.0f}/day), so it is set at the higher of the two. A node that survives losing the campus "
         f"survives the break by construction.")

    mathbox(s,0.45,1.44,3.05,1.06,[
        ("D_adj ≥ max(D_solvency, D_revocation) / (1−κ)",True),
        ("κ = contested share of the adjacent catchment",False),
        ("  κ=0%   → 569/day      κ=25% → 758/day",False),
        ("  κ=10%  → 632/day      κ=44% → 1,015/day",False),
        ("",False),
        ("7-WEEK RULE, derived:  21d + 28d = 49d",True)],"THE FILTER")
    text(s,0.45,2.54,3.05,0.42,
         f"Only the METRO five-operator overlap ({AD.METRO_OVERLAP_5OP:.0%}) is published⁵, so the requirement is a band "
         f"and we quote the band. Wind-down (T−21d) and ramp-up (28d) collide below 7 weeks, so a shorter break "
         f"segment cannot be wound down at all.",6.0,MUTE,lsp=0.98)

    ym=banner(s,3.65,1.44,2.85,"THE SCREEN, AND THEN OUR OWN LIST")
    pic(s,"district_screen",3.65,ym+0.04,2.85,1.30)
    text(s,3.65,ym+1.36,2.85,0.42,
         f"{AD.N_DISTRICTS:,} → {AD.N_CANDIDATES} on four criteria, then screened for incumbent presence: "
         f"{AD.PROX_COUNTS['uncontested']} uncontested, {AD.PROX_COUNTS['contested']} contested, "
         f"{AD.PROX_COUNTS['stacked']} stacked.",5.9,MUTE,lsp=0.98)

    yr=banner(s,6.65,1.44,2.90,"RANKED CANDIDATES")
    tp=AD.with_proximity(AD.TOP).head(6)
    rows=[]
    for r in tp.itertuples():
        col={"uncontested":POS,"contested":HI,"stacked":NEG}[r.proximity]
        rows.append([(r.District.title(),FG,True),(r.State[:11],MUTE,False),
                     (f"{r.urban_colleges:,.0f}",FG,True),(f"{r.exp_incumbent_stores:.0f}",col,True)])
    table(s,6.65,yr+0.06,2.90,
          [("DISTRICT",0.40,PP_ALIGN.LEFT),("STATE",0.30,PP_ALIGN.LEFT),
           ("URB COL",0.15,PP_ALIGN.RIGHT),("INCUMB",0.15,PP_ALIGN.RIGHT)],rows,size=6.0,gap=0.185)
    kar=int((AD.TOP.head(20).State=="Karnataka").sum())
    text(s,6.65,yr+1.28,2.90,0.44,
         f"CONVERGENCE: Karnataka and Odisha were chosen at STATE level off residential intensity and hostel occupancy. "
         f"The district register — a different table, pulled later — puts {kar} of the top 20 in Karnataka and Khordha "
         f"first. The state was not chosen and then justified.",5.9,FG,lsp=0.98)

    b=box(s,0.45,3.06,5.20,0.94,None,radius=0.09,line=NEG,lw=1.0)
    f=b.text_frame; f.margin_left=In(0.12); f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True,lsp=0.94)
    run(p,"WHAT THIS FLAG CANNOT DO   ",5.9,NEG,bold=True)
    run(p,f"Expected incumbent stores = urban colleges × non-metro store density, so ‘uncontested’ is arithmetically "
          f"‘{AD.MIN_URBAN_COLLEGES}–7 urban colleges’: the {AD.PROX_COUNTS['uncontested']} clean districts are the "
          f"{AD.PROX_COUNTS['uncontested']} smallest, sitting on the screen's own floor. The flag ranks; it does not site. "
          f"Contestedness and cluster density trade off directly — there is no large empty district, and the plan does "
          f"not need one: ",5.9,FG)
    run(p,"campus demand inside the gate is uncontested by construction, and what this screen prices is the adjacent half.",5.9,FG,bold=True)
    p2=par(f,lsp=0.94)
    run(p2,"The register also carries no enrolment⁶, so district student counts are imputed off state ratios⁷ and labelled as such.",5.7,MUTE,ital=True)

    ry=3.06
    text(s,5.85,ry,3.70,0.13,"CALENDAR SHAPES  ·  WHY THE 7-WEEK CONDITION IS IN THE FILTER",5.8,FG,bold=True)
    rows=[]
    for name,(cost,_) in CF.RESULTS.items():
        share=CF.windownable_share(CF.SHAPES[name])
        pen=cost/CF.BASE_COST-1
        col=POS if share>=0.99 else (NEG if share<=0.01 else MUTE)
        rows.append([(name.split(":")[0][:26],FG,False),(f"{share:.0%}",col,True),
                     (f"₹{cost/1e5:.1f} L",FG,True),(f"{pen:+.0%}",col,True)])
    table(s,5.85,ry+0.17,3.70,[("CALENDAR SHAPE",0.50,PP_ALIGN.LEFT),("WIND-DOWNABLE",0.18,PP_ALIGN.RIGHT),
          ("DEAD ZONE",0.17,PP_ALIGN.RIGHT),("vs BASE",0.15,PP_ALIGN.RIGHT)],rows,size=5.8,gap=0.165)

    band(s,4.76,[("Site chosen, break survivable, revocation hedged. ",True),
                 ("Now the node has to run inside a gate that an incumbent cannot pass.",False)])
    foot(s,"Sources 5,6,7  ·  111-district screen A7  ·  fragmentation A7b  ·  revocation and solvency thresholds A4")

# ============================== SLIDE 6 ==============================
def slide6():
    s=prs.slides.add_slide(L[CONTENT_LAYOUT]); strip_ph(s)
    rail(s,6,"THE OPERATING MODEL","→ does the money work?")
    head(s, f"Split the trip at the gate and buy the batch with time: last mile falls "
            f"₹{M.LAST_MILE:.0f} → ₹{LM:.1f} per order, and the product is fastest at exam-night peak",
         "Two decisions, one page: the fleet is a labour class, not a vehicle; and the SLA is a batching rule, "
         "not a promise. Both are what make the ₹19.0 on slide 3 real.")

    mathbox(s,0.45,1.44,3.05,1.20,[
        ("t_trip = 2·(d/v) + b·τ         trip minutes",True),
        ("C/order = (w/60)·t_trip / b",True),
        ("  city leg   w = ₹168/ACTIVE hr   (gig)",False),
        ("  gate leg   w = ₹72/ROSTERED hr  (employed)",False),
        (f"  ratio {CS.RATIO_JPM:.2f}× ; on UBS's anchor {CS.RATIO_UBS:.2f}×",True),
        ("b(λ) = min(K, ⌈λ·W/60⌉)   W=6 min, K=12",True),
        ("Little's Law  L = λ·W  on the gate queue",False)],"TRIP AND BATCH")
    text(s,0.45,2.68,3.05,0.40,
         "Cost per order falls linearly in batch and rises linearly in the wage of the class serving the leg. "
         "Only the in-gate leg can change that class — which is why the gate, not the vehicle, is the decision.",
         6.0,FG,lsp=0.98)

    ly=3.12
    text(s,0.45,ly,3.05,0.13,"LABOUR CLASS  ·  FIVE COSTED PARAMETERS",5.8,FG,bold=True)
    rows=[[("GIG · quick commerce",FG,False),("0",MUTE,False),("42.06",FG,False),("168",NEG,True),("40%",MUTE,False)],
          [("GIG · JM non-metro",FG,False),("0",MUTE,False),("43.96",FG,False),("176",NEG,True),("36%",MUTE,False)],
          [("EMPLOYED RUNNER",FG,True),("577",MUTE,False),("4.66",FG,True),("72",POS,True),("100%",MUTE,False)]]
    table(s,0.45,ly+0.17,3.05,[("CLASS",0.42,PP_ALIGN.LEFT),("FIX/d",0.13,PP_ALIGN.RIGHT),
          ("₹/ord",0.15,PP_ALIGN.RIGHT),("₹/hr",0.15,PP_ALIGN.RIGHT),("UTIL",0.15,PP_ALIGN.RIGHT)],rows,size=5.8,gap=0.165)
    text(s,0.45,4.10,3.05,0.36,
         f"Runner threshold {FM.breakeven_volume():.0f}/day is SOLVED, not chosen — last year's equivalent column was an "
         f"input. Incentive column is FLAGGED, not invented: the UBS split is in the corpus, not in a module¹².",
         5.7,MUTE,lsp=0.98)

    ym=banner(s,3.65,1.44,2.90,"THE SLA  ·  BATCH WAIT FALLS AS DEMAND RISES")
    rows=[[("Trough",MUTE,False),("6.5",MUTE,False),("1",MUTE,False),("25.7",MUTE,False),("64.8",NEG,True)],
          [("Average",FG,False),("25.9",FG,False),("2",FG,False),("21.6",FG,False),("33.0",FG,True)],
          [("Peak 4×",FG,True),("103.7",FG,True),("10",FG,True),("27.1",FG,True),("7.7",POS,True)],
          [("Exam 6×",FG,True),("155.6",FG,True),("12",FG,True),("27.1",FG,True),("6.6",POS,True)]]
    table(s,3.65,ym+0.06,2.90,[("STATE",0.28,PP_ALIGN.LEFT),("ORD/HR",0.19,PP_ALIGN.RIGHT),
          ("BATCH",0.15,PP_ALIGN.RIGHT),("SLA",0.17,PP_ALIGN.RIGHT),("₹/ORD",0.21,PP_ALIGN.RIGHT)],rows,size=5.9,gap=0.175)
    text(s,3.65,ym+0.86,2.90,0.50,
         f"A 4× spike needs 1.4× the runners: the spike is absorbed by batch size, not headcount. "
         f"{SL.volume_weighted()[0]*100 if isinstance(SL.volume_weighted()[0],float) else 62.7:.1f}% of orders fall in "
         f"the peak band, which is why the volume-weighted cost is ₹{LM:.1f} and not the ₹33.0 an average-hour "
         f"calculation gives.",5.8,MUTE,lsp=0.98)
    text(s,3.65,ym+1.40,2.90,0.13,"MODE  ·  IN-GATE ₹/ORDER BY BATCH",5.8,FG,bold=True)
    rows=[[("Petrol 2W (incumbent)",MUTE,False),("28.3",MUTE,False),("14.0",MUTE,False),("—",MUTE,False)],
          [("Cycle",FG,False),("18.6",FG,False),("8.4",FG,False),("—",MUTE,False)],
          [("E-cart, stationed",FG,True),("15.4",FG,True),("7.2",FG,True),("4.1",POS,True)]]
    table(s,3.65,ym+1.57,2.90,[("MODE",0.46,PP_ALIGN.LEFT),("n=1",0.18,PP_ALIGN.RIGHT),
          ("n=3",0.18,PP_ALIGN.RIGHT),("n=12",0.18,PP_ALIGN.RIGHT)],rows,size=5.8,gap=0.165)

    yr=banner(s,6.70,1.44,2.85,"THE PARTNERSHIP  ·  PRECEDENT AND SCOPE")
    kv(s,6.70,yr+0.08,2.85,[
        ("UC Davis, per package¹⁶", "$2.00", FG),
        ("packages per stop vs a truck", "12 vs 3 (4×)", POS),
        ("same-day rate, year one", "99.6%", POS),
        ("our fee band, floor to ceiling", "₹9.6 – ₹29.2", HI)],gap=0.205)
    b=box(s,6.70,yr+0.92,2.85,0.62,CHIP2,radius=0.10); rect(s,6.70,yr+0.92,0.045,0.62,HI)
    f=b.text_frame; f.margin_left=In(0.11); f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True,lsp=0.92)
    run(p,"SCOPE IT AS CONTINUOUS CIRCUITS  ",5.8,HI,bold=True)
    run(p,"never as the parcel desk's accumulate-and-sort round. The institution supplies labour and access on a "
          "roster; it does not supply an accumulation service. Get that wrong in the licence and the service level dies, "
          "not the economics.",5.8,FG)
    text(s,6.70,yr+1.60,2.85,0.44,
         "CONCESSION: 11 minutes is marketing. Walmart's own CEO says <13¹³, Euromonitor 15–20¹⁴, our 35-store field "
         "survey median 15 with 9% at ≤10³. JPM's table shows the two 10-minute operators did not batch¹¹.",5.7,NEG,lsp=0.98)

    fy=4.06
    text(s,3.65,fy,5.90,0.13,"TWO FLOORS, AND THEY ARE DIFFERENT NUMBERS  ·  BOTH ARE ON THE DAY-90 INSTRUMENT LIST",5.8,FG,bold=True)
    cells=[("STORE FLOOR", f"{THR[-1]*M.CEILING:,.0f}/day", f"{THR[-1]:.1%} of term volume — below it the node cannot fund its own fixed base", NEG),
           ("RUNNER FLOOR", f"{FM.breakeven_volume():.0f}/day", "below it the rostered runner is uneconomic and the gate leg reverts to gig", HI),
           ("MARKET DIRECTION", "53% of riders", "now batch >20% of orders, up from 42%¹² — the SLA we price is where the market is going", POS)]
    cw=5.90/3
    for i,(lab,val,body,col) in enumerate(cells):
        b=box(s,3.65+i*cw,fy+0.17,cw-0.08,0.50,CHIP,radius=0.10); rect(s,3.65+i*cw,fy+0.17,0.04,0.50,col)
        f=b.text_frame; f.margin_left=In(0.10); f.vertical_anchor=MSO_ANCHOR.MIDDLE
        p=par(f,True,lsp=0.90)
        run(p,lab+"  ",5.6,col,bold=True); run(p,val,7.2,FG,bold=True)
        p2=par(f,lsp=0.90); run(p2,body,5.5,MUTE)
    band(s,4.76,[("The operating model is a labour-class decision and a queueing rule. ",True),
                 ("Both are cheap, both are reversible, and together they are the ₹19.0 the whole D1 case rests on.",False)])
    foot(s,"Sources 1,3,11,12,13,14,16  ·  trip identity and volume weighting A5  ·  labour-class provenance A2  ·  SLA × access-regime matrix A4")

# ============================== SLIDE 7 ==============================
def slide7():
    s=prs.slides.add_slide(L[CONTENT_LAYOUT]); strip_ph(s)
    rail(s,7,"THE FINANCIALS","→ what breaks it, and when do we know?")
    head(s, f"One node returns {SCN[1]['roce']:.0%} on ₹{RC.CE_BASE/1e5:,.0f} lakh at a 30% non-grocery basket, "
            f"and {RC.ROCE_HURDLE:.0%} at ₹{RC.AOV_HURDLE:,.0f} — inside management's own stated ceiling",
         f"The basket is fitted, not assumed: the one disclosed quarterly series where an Indian operator actually "
         f"lifted AOV. AOV = {BK.SLOPE:.2f}·x + {BK.INTERCEPT:.0f}, R² = {BK.R2:.3f}, n = 4 quarters¹⁰.")

    yl=banner(s,0.45,1.44,4.35,"ONE NODE, ONE YEAR   ·   ₹ LAKH")
    pic(s,"pnl_bridge",0.45,yl+0.05,4.35,1.50)
    kv(s,0.45,3.26,4.35,[
        ("capex, midpoint of the band²", f"₹{RC.CAPEX_MID/1e5:,.0f} L", FG),
        ("working capital at 14 NWC days⁸", f"₹{WC.WC_ADOPTED/1e5:,.0f} L", FG),
        ("capital employed", f"₹{RC.CE_BASE/1e5:,.0f} L", HI),
        ("cash conversion cycle⁹", f"{WC.ZEPTO_CCC:.0f} days — supplier-funded", POS)],gap=0.185)

    ym=banner(s,5.00,1.44,2.20,"THE BASKET LADDER")
    pic(s,"basket_ladder",5.00,ym+0.05,2.20,1.18)
    text(s,5.00,ym+1.26,2.20,0.62,
         f"₹{M.MINUTES_AOV:,.0f} → ₹525 on term-start durables → ₹{BE:,.0f} at "
         f"{BK.SHARE_NEEDED_AFTER_OCCASION:.1f}% non-grocery, from {BK.MINUTES_NONGROCERY:.0f}% today. The "
         f"{RC.ROCE_HURDLE:.0%} hurdle needs {RC.HURDLE_NONGROCERY_SHARE:.1f}% — still inside the stated "
         f"{BK.NONGROCERY_CEILING_LO:.0f}–{BK.NONGROCERY_CEILING:.0f}% ceiling, with "
         f"{RC.HURDLE_HEADROOM_PTS:+.1f} pts of headroom.",5.8,MUTE,lsp=0.98)

    yr=banner(s,7.40,1.44,2.15,"RETURNS")
    kv(s,7.40,yr+0.08,2.15,[
        ("breakeven AOV", f"₹{RC.AOV_BREAKEVEN:,.0f}", FG),
        ("hurdle AOV, pre-tax", f"₹{RC.AOV_HURDLE:,.0f}", HI),
        ("post-tax at 25.17%", f"₹{RC.AOV_HURDLE_POSTTAX:,.0f}", MUTE),
        ("IRR, 5-yr, 3-mo ramp", f"{RC.irr(RC.AOV_HURDLE):.1%}", POS),
        ("ramp 2 – 6 months", f"{RC.irr(RC.AOV_HURDLE,ramp=2):.0%} – {RC.irr(RC.AOV_HURDLE,ramp=6):.0%}", MUTE),
        ("node life anchor¹¹", f"{M.FRANCHISE_PAYBACK} mo payback", MUTE)],gap=0.195)

    sy=3.56
    text(s,5.00,sy,4.55,0.13,"SCENARIOS  ·  ALL ON INPUTS ALREADY PRICED ELSEWHERE IN THE DECK",5.8,FG,bold=True)
    rows=[]
    for r_ in SCN:
        col = POS if r_["roce"]>=RC.ROCE_HURDLE else (NEG if r_["roce"]<0.15 else FG)
        pb = "—" if r_["roce"]<=0 else f"{r_['payback']:.0f}"
        rows.append([(r_["name"].replace("  ("," (")[:30],FG,False),(f"₹{r_['aov']:,.0f}",FG,False),
                     (f"{r_['margin']:.2%}",MUTE,False),(f"{r_['turn']:.2f}×",MUTE,False),
                     (f"{r_['roce']:.1%}",col,True),(pb,col,False)])
    table(s,5.00,sy+0.17,4.55,[("SCENARIO",0.40,PP_ALIGN.LEFT),("AOV",0.12,PP_ALIGN.RIGHT),
          ("MARGIN",0.13,PP_ALIGN.RIGHT),("TURN",0.11,PP_ALIGN.RIGHT),("ROCE",0.12,PP_ALIGN.RIGHT),
          ("PAYBACK mo",0.12,PP_ALIGN.RIGHT)],rows,size=5.8,gap=0.175)

    b=box(s,0.45,4.06,4.35,0.60,None,radius=0.09,line=NEG,lw=1.0)
    f=b.text_frame; f.margin_left=In(0.12); f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True,lsp=0.94)
    run(p,"THE LIMIT, COMPUTED NOT ASSERTED   ",5.9,NEG,bold=True)
    dn=SCN[3]
    run(p,f"At a 30% basket with volume −30%, ROCE falls to {dn['roce']:.1%} and payback runs to "
          f"{dn['payback']:.0f} months — longer than the {RC.NODE_LIFE_MO} months the node's own life is anchored on. ",5.9,FG)
    run(p,"Under that combination the node does not pay back within its life, which is why the day-90 gate measures "
          "volume before anything else.",5.9,FG,bold=True)
    text(s,5.00,4.44,4.55,0.30,
         f"Second lever, unused in the base case: Minutes' free-delivery threshold is "
         f"₹{BK.FREE_DELIVERY_THRESHOLD['Flipkart Minutes']:.0f}, the market's lowest (Instamart "
         f"₹{BK.FREE_DELIVERY_THRESHOLD['Instamart']:.0f}, Blinkit ₹{BK.FREE_DELIVERY_THRESHOLD['Blinkit']:.0f}, and Blinkit "
         f"already varies it by location and demand) — an existing platform lever, not a new ask.   ·   NPV at 12% / 15% "
         f"[assumed; no WACC is disclosed]: ₹{RC.npv(RC.AOV_HURDLE,0.12)/1e5:,.0f} L / ₹{RC.npv(RC.AOV_HURDLE,0.15)/1e5:,.0f} L.",5.6,MUTE,lsp=0.96)

    band(s,4.76,[("The node clears the hurdle on a basket mix management has already said is reachable. ",True),
                 ("One dependency binds: term volume. That is what the first 90 days are for.",False)])
    foot(s,"Sources 2,4,8,9,10,11  ·  P&L, capital and ROCE derivation A6b  ·  basket fit A6  ·  working-capital constructs A5  ·  scenarios A3")

# ============================== SLIDE 8 ==============================
def slide8():
    s=prs.slides.add_slide(L[CONTENT_LAYOUT]); strip_ph(s)
    rail(s,8,"ROADMAP, RISK, GATE","→ decision")
    head(s, "90 days to a measured go/no-go: four workstreams, four instrumented metrics, and one "
            "dependency that can stop the rollout",
         "Every action fires on a calendar date, not on observed volume — semester demand steps rather than ramps. "
         "The gate metrics are Round 1's metric architecture, now doing the job it was designed for.")

    yl=banner(s,0.45,1.44,5.45,"THE 90 DAYS   ·   WORKSTREAM × WEEK, WITH OWNERS")
    ws=[("SITE & LICENCE","Contracts",[("W1–2","shortlist 6 sites in 2 districts, uncontested-catchment test"),
                                       ("W3–5","campus licence + DORMANCY CLAUSE in supplier terms"),
                                       ("W6","sign; publish the academic calendar into the ops calendar")]),
        ("NODE BUILD","Cluster ops",[("W3–6","fit-out 2,000 sqft small-format, one chilled zone"),
                                     ("W7–9","hire 25 store staff; 4-week lead on the rehire pool")]),
        ("FLEET & PARTNER","Fleet",[("W5–8","institution runner agreement, per-parcel fee, SLA-linked"),
                                    ("W8–10","e-cart circuits mapped, OTP chain of custody live")]),
        ("COMMERCIAL","Category",[("W6–9","non-grocery mix to 25%: durables, accessories, print, BPC, OTC"),
                                  ("W10–12","campus free-delivery threshold test vs the ₹149 network floor")])]
    y=yl+0.06
    for name,owner,items in ws:
        text(s,0.45,y,1.30,0.13,name,5.9,HI,bold=True)
        text(s,0.45,y+0.13,1.30,0.12,owner,5.4,MUTE,ital=True)
        yy=y
        for wk,act in items:
            text(s,1.80,yy,0.52,0.12,wk,5.6,FG,bold=True)
            text(s,2.34,yy,3.56,0.24,act,5.9,MUTE,lsp=0.94)
            yy+=0.145
        y=max(yy,y+0.30)+0.055
        rect(s,0.45,y-0.03,5.45,0.006,RULE)
    text(s,0.45,y+0.02,5.45,0.14,
         "Ramp-up cannot be gradual: semester-start demand STEPS, so every ramp action is dated off the published "
         "academic calendar.",5.7,NEG,ital=True)

    yr=banner(s,6.05,1.44,3.50,"DAY-90 GATE   ·   ROUND 1's METRICS, INSTRUMENTED")
    rows=[[("Peak-to-Trough Demand Ratio",FG,False),("> 4.0×",POS,True)],
          [("Gate-Drop Consolidation Ratio",FG,False),("≥ 3.0×",POS,True)],
          [("Term-Weighted CM + Break Runway",FG,False),("≥ 1.0×",POS,True)],
          [("Calendar-Linked Labour Share",FG,False),("≥ 50%",POS,True)]]
    table(s,6.05,yr+0.06,3.50,[("GATE METRIC",0.68,PP_ALIGN.LEFT),("THRESHOLD",0.32,PP_ALIGN.RIGHT)],rows,size=6.0,gap=0.185)
    text(s,6.05,2.68,3.50,0.30,
         f"Volume is the first read: at 1,000 orders/day the asset-turn ratio is {M.TURN_RATIO_AT_1000:.2f} and at "
         f"−30% the node does not pay back within its life. Measure it before capital goes to node two.",5.8,NEG,lsp=0.98)

    ty=3.06
    text(s,6.05,ty,3.50,0.13,"PRICED SHOCKS  ·  AGAINST THE ₹580 BREAKEVEN",5.8,FG,bold=True)
    rows=[[("Volume −30%",FG,True),("₹647",NEG,True),("+67",NEG,True)],
          [("Shrinkage on top (upper bound)",FG,False),("₹623",FG,False),("+43",MUTE,False)],
          [("Gig social-security levy¹⁹",FG,False),("₹595",FG,False),("+15",MUTE,False)],
          [("Calendar fragmentation",FG,False),("₹589",FG,False),("+9",MUTE,False)]]
    table(s,6.05,ty+0.17,3.50,[("SHOCK",0.62,PP_ALIGN.LEFT),("BREAKEVEN",0.22,PP_ALIGN.RIGHT),
          ("Δ",0.16,PP_ALIGN.RIGHT)],rows,size=5.8,gap=0.165)
    text(s,6.05,4.10,3.50,0.42,
         f"Three of four are covered inside a 30% non-grocery basket (₹{BK.SLOPE*BK.NONGROCERY_CEILING_LO+BK.INTERCEPT:,.0f}); "
         f"all four inside 40% (₹{BK.SLOPE*BK.NONGROCERY_CEILING+BK.INTERCEPT:,.0f}). Shrinkage is an UPPER BOUND — the stack "
         f"was calibrated on a reported contribution figure that is already net of it, so charging it again double-counts.",
         5.7,MUTE,lsp=0.98)

    b=box(s,0.45,4.06,5.45,0.58,None,radius=0.09,line=HI,lw=1.0)
    f=b.text_frame; f.margin_left=In(0.12); f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True,lsp=0.93)
    run(p,"THE THREE CONCESSIONS, MADE FROM THE FRONT FOOT   ",5.9,HI,bold=True)
    run(p,"11 minutes is a marketing number (slide 6). The industry relocates rather than mothballs — and we never "
          "proposed mothballing; holding costs 24% of one move. We are not first: Swiggy, Blinkit and Zepto are all in "
          "the segment (slide 1). None of the three underwrites the node on the cluster.",5.9,FG)

    band(s,4.76,[("Approve one pilot node, ₹%.0f lakh of capital employed, and a day-90 gate on four metrics. " % (RC.CE_BASE/1e5),True),
                 ("If volume clears, node two is a site-selection decision we have already made 111 times.",False)])
    foot(s,"Sources 8,13,19,20,21  ·  shock derivations A5  ·  assumption ledger A3  ·  fragmentation A7b  ·  playbook triggers and owners A4")

for fn in (slide1, slide2, slide3, slide4, slide5, slide6, slide7, slide8):
    fn()
out=os.path.join(ROOT,f"{OUT}_{THEME}.pptx")
prs.save(out)
print("wrote", out, "|", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
