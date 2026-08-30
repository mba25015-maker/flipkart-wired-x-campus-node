"""
VERIFY THE DECK, NOT THE DOCUMENT.

`verify_docs.py` parses figures out of HANDOFF.md. This does the same job one step later:
it opens the BUILT .pptx, pulls every run of text out of it, and asserts that the model's
figures are the ones physically on the slides -- and that no superseded Round 1 figure is.

  python3 verify_deck.py [path.pptx]

A slide number that drifts from the model fails here even if `audit.py` passes, because
audit.py checks the model against itself; this checks the artifact the panel will see.
"""
import sys, os, re
from pptx import Presentation
import campus_model as M, cost_stack as CS, break_mode as B
import aishe_district as AD, risk_quadrant as Q, calendar_fragmentation as CF
import basket as BK, sla as SL, roce as RC, fleet_mix as FM, working_capital as WC
from deck_checks import DECK_CHECK_COUNT
import risk_shocks as RS, working_capital as WC

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
import params as P
import check_counts as CC
# The default used to be the SUPERSEDED four-slide reference build. Run bare - as the deck
# itself instructs a judge to run it - this verifier reported 51/93 and DECK VERIFICATION
# FAILED, and nothing caught it because it was only ever invoked by hand with an explicit
# path. The default is now the single constant that names the final deck.
PPTX = sys.argv[1] if len(sys.argv) > 1 else os.path.join(ROOT, P.FINAL_DECK)

def _chart_text(ch):
    """A native chart's numbers live in the chart part, not in a text frame. Once a table
    becomes a chart the figures vanish from the deck's readable text and every check on
    them passes silently. So read the categories, the plotted values and any literal data
    labels back out, and treat them as text on the slide."""
    buf = []
    try:
        buf += [str(c) for c in ch.plots[0].categories]
    except Exception:
        pass
    for ser in ch.series:
        try:
            for v in ser.values:
                if v is None:
                    continue
                buf.append(f"{v:g}")
                buf.append(f"{v:,.0f}")
                buf.append(f"{v:.1f}")
        except Exception:
            pass
    for dLbl in ch._chartSpace.iter():
        if dLbl.tag.endswith("}t") and dLbl.text:
            buf.append(dLbl.text)
    return " ".join(buf)

def deck_text(path):
    prs = Presentation(path); out = []
    for s in prs.slides:
        buf = []
        for sh in s.shapes:
            if sh.has_text_frame:
                buf.append(" ".join(r.text for p in sh.text_frame.paragraphs for r in p.runs))
            if getattr(sh, "has_chart", False):
                buf.append(_chart_text(sh.chart))
        out.append(" ".join(buf))
    return out

pages = deck_text(PPTX)
ALL = " ".join(pages)
checks = []
def _norm(t):
    """Compare on content, not on glyphs. The deck writes ₹, ×, en dashes and non-breaking
    spaces; the model prints Rs, x and hyphens. A check must not fail on typography."""
    for a, b in (("\u20b9","Rs"), ("\u00d7","x"), ("\u2013","-"), ("\u2014","-"),
                 ("\u2212","-"), ("\u00a0"," "), ("\u2019","'"), ("\u201c",'"'), ("\u201d",'"'), (" lakh"," L")):
        t = t.replace(a, b)
    return t

def must(label, needle, page=None):
    """Figures are matched exactly; prose is matched case-insensitively, because a slide
    may shout a phrase the model spells in lower case. The page is taken from the label's
    own prefix when the caller does not pass one, so re-ordering the deck is a one-line change."""
    if page is None:
        page = PAGE.get(label.split()[0])
    hay = _norm(ALL if page is None else pages[page]); nd = _norm(needle)
    checks.append((nd in hay or nd.lower() in hay.lower(), label, needle))
def must_not(label, needle):
    checks.append((_norm(needle) not in _norm(ALL), label, "ABSENT: " + needle))
def must_not_on(label, needle, page=None):
    """Absence on ONE page. Needed where a figure is correct on one slide and the wrong
    basis on another: Rs729 is a legitimate scenario AOV on slide 7 and a basis error in
    slide 8's coverage sentence, so a deck-wide must_not would be false."""
    if page is None:
        page = PAGE.get(label.split()[0])
    hay = _norm(ALL if page is None else pages[page])
    checks.append((_norm(needle) not in hay, label, f"ABSENT from p{(page or 0)+1}: " + needle))

thr   = [B.threshold(fn) for _, fn in B.CONFIGS]
rev   = Q.post_revocation_survival()
rel   = B.relocate_vs_flex()
be    = CS.breakeven_d2_consistent(CS.CAMPUS_FIXED, SL.volume_weighted()[1])

# Slide order. The 8-slide REFERENCE deck (S31 architecture) and the earlier 4-slide build have
# different page indices, so the page is INFERRED from the check's own label prefix rather than
# hard-coded twice. A deck with a different slide count fails loudly instead of silently passing.
# A cover slide, when present, shifts every content page by one. Detect it by CONTENT, not by
# slide count: counting broke the moment an appendix page was added, silently zeroing the page
# map and turning 106 page-scoped checks into deck-wide ones. Slide 1 always carries its rail
# label; the cover never does.
COVER = 0 if "RECOMMENDATION" in _norm(pages[0]).upper() else 1
if COVER:
    pages = pages[0:]          # keep the cover in ALL for the guillemet check below

if len(pages) == 4:                       # the earlier build: slides 1, 2, 5, 6 only
    PAGE = {"S1":0,"S2":1,"S5":2,"S6":3}
elif len(pages) - COVER >= 8:             # S31: rec, market, dead zone, return, site, ops, fin, roadmap
    # The first EIGHT are the content slides; anything after them is appendix. Keying off the
    # total slide count silently emptied PAGE the moment the appendix was added, which turned
    # every page-scoped check into a deck-wide one -- a check that cannot fail is not a check.
    PAGE = {k: i + COVER for i, k in enumerate(["S1","S2","S3","S4","S5","S6","S7","S8"])}
else:
    PAGE = {}
S1, S2, S5, S6 = PAGE.get("S1"), PAGE.get("S2"), PAGE.get("S5"), PAGE.get("S6")

# The appendix, when it is present. Content slides are 1-8; A0 onwards follow in order.
APX_ORDER = ["A0","A1","A2","A3","A4","A5","A5b","A5c","A6","A6b","A7","A7b","A8","A9","A10"]
HAS_APX = len(pages) - COVER - 8 == len(APX_ORDER)
if HAS_APX:
    PAGE.update({k: 8 + COVER + i for i, k in enumerate(APX_ORDER)})
# AND IF IT DOES NOT MATCH, SAY SO. This list is typed and the appendix is built, so the two
# can disagree - and when they did, HAS_APX simply went False and 36 page-scoped checks
# quietly became deck-wide ones. The run still reported "all pass", on 36 fewer checks. A
# structure mirrored by hand needs an assertion that the mirror is still true; silence is the
# defect, not the mismatch. Adding a5c to the appendix is what surfaced this.
checks.append((HAS_APX, "SELF  appendix page map covers every appendix slide",
               f"{len(APX_ORDER)} mapped, {len(pages) - COVER - 8} in the deck"))

# ---- slide 1: the executive summary --------------------------------------------
must("S1  headline asset-turn ratio",            f"{M.TURN_RATIO:.0%}")
must("S1  ladder start",                         f"{thr[0]:.1%}")
must("S1  ladder floor",                         f"{thr[-1]:.1%}")
must("S1  do-nothing burn through the break",    f"Rs{rel['do_nothing_total']/1e5:.1f} lakh")
must("S1  cost of one relocation",               f"Rs{rel['relocate_once']/1e5:.0f} lakh")
must("S1  holding as a share of relocating",     f"{rel['flex_vs_relocate']:.0%}")
must("S1  D1/D2-consistent breakeven",           f"Rs{be:,.0f}")
must("S1  candidate districts",                  f"{AD.N_CANDIDATES} of {AD.N_DISTRICTS:,}")
must("S1  campus asset turn",                    f"{M.CAMPUS_TURN:.2f}x")
must("S7  non-grocery share required",           f"{BK.SHARE_NEEDED_AFTER_OCCASION:.1f}%")
must("S7  the range disclosed by Swiggy",          f"{BK.NONGROCERY_CEILING_LO:.0f}–{BK.NONGROCERY_CEILING:.0f}%")
must("S1  the not-first concession is present",  "not first")

# ---- slide 2: the market -------------------------------------------------------
must("S2  metro oversupply",                     f"{AD.METRO_EXCESS:.0%}")
must("S2  colleges in the register",             f"{AD.N_COL:,}")
must("S2  urban colleges",                       f"{AD.URBAN_COL:,}")
must("S2  urban non-metro colleges",             f"{AD.NON_METRO_URBAN_COLLEGES:,}")
must("S2  urban high-propensity standalones",    f"{AD.N_STA_URBAN_HP:,}")
must("S2  the Starship hook is on the slide",    "Starship")
must("S2  the quote is verbatim",                "365-day urban business")
must("S2  Euromonitor white space",              "Euromonitor")

# ---- the dead-zone slide: every headline figure (slide 5 in the 4-slide build, 3 in S31)
must("S3  flexible share of the fixed base",     f"{CS.fixed_flex_share():.1%}")
must("S3  do-nothing residual threshold",        f"{thr[0]:.1%}")
must("S3  post-ladder residual threshold",       f"{thr[-1]:.1%}")
must("S3  fixed base",                           f"Rs{CS.CAMPUS_FIXED:,.0f}")
must("S3  tier rent band, low",                  f"Rs{CS.T12_FIXED_LOW:,.0f}")
must("S3  tier rent band, high",                 f"Rs{CS.T12_FIXED_HIGH:,.0f}")
must("S3  'other fixed' residual, stated openly", f"Rs{CS.OTHER_FIXED:,.0f}")
must("S4  the moneyshot ratio",                  f"{M.TURN_RATIO:.3f}")
must("S4  density term",                         f"{M.CEILING/M.MINUTES_ORD:.3f}")
must("S4  calendar term",                        f"{M.ACTIVE_MONTHS/12:.3f}")
must("S4  campus asset turn",                    f"{M.CAMPUS_TURN:.2f}x")
must("S4  city asset turn",                      f"{M.CITY_TURN:.2f}x")
must("S4  sensitivity, ratio at 1,000/day",      f"{M.TURN_RATIO_AT_1000:.3f}")
must("S4  sensitivity, ratio at 1,200/day",      f"{M.TURN_RATIO_AT_1200:.3f}")
must("S4  parity throughput",                    f"{M.PARITY_OPD:,.0f}/day")
must("S4  Blinkit observed range, low",          f"{M.CEIL_LO:,}")
must("S4  Blinkit observed range, high",         f"{M.CEIL_HI:,}")
must("S4  Minutes AOV on the adopted basis",     f"Rs{M.MINUTES_AOV:,.0f}")

# ---- the site-selection slide (6 in the 4-slide build, 5 in S31) -----------------
must("S5  site filter, uncontested",             f"{AD.SITE_FILTER_UNCONTESTED:,.0f}")
must("S5  candidate districts",                  f"{AD.N_CANDIDATES}")
must("S5  districts in the register",            f"{AD.N_DISTRICTS:,}")
must("S5  stacked share of our own list",        f"{AD.STACKED_SHARE:.0%}")
must("S5  uncontested districts",                f"{AD.N_UNCONTESTED}")
must("S5  revocation binding requirement",       f"{rev['orders_needed']:,.0f}")
must("S5  break-solvency requirement",           f"{rev['adjacent_lo']:,.0f}")
must("S5  gross demand at metro overlap",        f"{AD.gross_demand_required(AD.METRO_OVERLAP_5OP):,.0f}/day")
must("S5  metro five-operator overlap",          f"{AD.METRO_OVERLAP_5OP:.0%}")
must("S5  the 7-week rule is derived on-slide",   "7-WEEK RULE")
must("S5  wind-downable share, typical calendar", f"{CF.TYPICAL_WINDOWNABLE:.0%}")
must("S5  fragmentation penalty, typical",       f"{CF.TYPICAL_PENALTY:+.0%}")
must("S2  non-metro store base",                 f"{AD.NON_METRO_STORES:,}")
must("S2  five-operator store total",            f"{AD.TOTAL_STORES:,}")

# ---- slide 4: the return (S31) --------------------------------------------------
must("S4  ROCE hurdle",                          f"{RC.ROCE_HURDLE:.0%}")
must("S4  AOV at the hurdle",                    f"{RC.AOV_HURDLE:,.0f}")
must("S4  AOV at breakeven, ROCE basis",         f"{RC.AOV_BREAKEVEN:,.0f}")
must("S4  hurdle premium",                       f"{RC.HURDLE_PREMIUM:,.0f}")
must("S4  non-grocery share the hurdle implies", f"{RC.HURDLE_NONGROCERY_SHARE:.1f}%")
must("S4  DuPont turnover leg",                  f"{RC.dupont(RC.AOV_HURDLE)['capital_turn']:.2f}")
must("S4  DuPont margin leg",                    f"{RC.dupont(RC.AOV_HURDLE)['ebit_margin']:.2%}")
must("S4  capital employed",                     f"{RC.CE_BASE/1e5:,.0f}")

# ---- slide 6: the operating model ----------------------------------------------
must("S6  runner breakeven volume, solved",      f"{FM.breakeven_volume():.0f}/day")
must("S6  gig:runner ratio, JPM anchor",         f"{CS.RATIO_JPM:.2f}")
must("S6  gig:runner ratio, UBS anchor",         f"{CS.RATIO_UBS:.2f}")

# ---- slide 7: the financials ----------------------------------------------------
must("S7  capex midpoint",                       f"{RC.CAPEX_MID/1e5:,.0f}")
must("S7  working capital",                      f"{WC.WC_ADOPTED/1e5:,.0f}")
must("S7  cash conversion cycle",                f"{WC.ZEPTO_CCC:.0f}")
must("S7  IRR at the hurdle AOV",                f"{RC.irr(RC.AOV_HURDLE):.1%}")
must("S7  ROCE at a 30% basket",                 f"{RC.roce(RC.AOV_UP):.1%}")
must("S7  post-tax hurdle AOV",                  f"{RC.AOV_HURDLE_POSTTAX:,.0f}")

_SC = RC.scenario_rows()

# ---- slide 8: roadmap and gate --------------------------------------------------
must("S8  the day-90 gate is named",             "day-90")
must("S8  volume shock breakeven",               f"{RS.AOV_VOLUME:.0f}")
must("S8  the three concessions are present",    "concessions")
# The coverage verdict "three of four" is computed at RS.AOV_CEILING_LO -- the fit's SLOPE
# applied as a delta from Minutes' post-occasion basket. roce.py's scenario AOVs (Rs729 /
# Rs842, correct and present on slide 7) are the regression read as a LEVEL. Quoting the
# level beside the verdict makes the sentence false: at Rs729 all four shocks clear.
must("S8  basket coverage at the 30% floor",     f"{RS.AOV_CEILING_LO:.0f}")
must("S8  basket coverage at the 40% ceiling",   f"{RS.AOV_CEILING_HI:.0f}")
must_not_on("S8  scenario AOV is not the coverage basis", f"{RC.AOV_UP:,.0f}")

# ---- the two native charts: their values must survive into the chart part ----------
must("S4  chart carries the downside ROCE",      f"{_SC[3]['roce']*100:.1f}")
must("S4  chart carries the base-case ROCE",     f"{_SC[1]['roce']*100:.1f}")
must("S4  chart carries the 40% basket ROCE",    f"{_SC[2]['roce']*100:.1f}")
must("S4  chart carries the hurdle line",        "hurdle 40%")
must("S8  tornado carries shrinkage",            f"{RS.AOV_SHRINKAGE:.0f}")
must("S8  tornado carries the levy",             f"{RS.AOV_LEVY:.0f}")
must("S8  tornado carries fragmentation",        f"{RS.AOV_FRAGMENTATION:.0f}")
must("S8  tornado carries the basket line",      f"{RS.AOV_CEILING_LO:,.0f}")

# ---- the appendix, if it is built --------------------------------------------------
if HAS_APX:
    must("A0   the verification counts",          f"{CC.AUDIT_COUNT}")
    must("A0   declared check count",             f"{DECK_CHECK_COUNT} deck checks")
    must("A0   the repository is named",          "github.com/mba25015-maker/flipkart-wired-x-campus-node")
    # A-1 reports the package's verification of itself, which is circular for THIS run:
    # run_all.py writes _verification.json only after verify_deck.py has finished, so the
    # results file this process can read is the PREVIOUS run's. Asserting the passed form
    # here would therefore be asserting yesterday's outcome.
    #
    # So the two questions are split. Here: is the claim a LEGAL FORM at all - either
    # "N checks defined" or "N / N", with N the number of checks that actually exist? That
    # catches a stale or invented count either way. Whether the passed form is EARNED is
    # release.py step 4's job, which re-verifies the exact stamped file and then requires it.
    _legal = [f"{CC.AUDIT_COUNT} / {CC.AUDIT_COUNT}", f"{CC.AUDIT_COUNT} checks defined"]
    _hay = _norm(ALL)
    checks.append((any(_norm(f) in _hay for f in _legal),
                   "A1   the audit self-report is a legal form",
                   " or ".join(_legal)))
    must("A2   the register is complete",         "VOC 2026 consumer survey")
    must("A3   the e-cart reductio",              "8.4 lakh")
    must("A5b  the per-gate roster at plan",      f"{FM.plan_roster()[0]//FM.plan_roster()[3]}")
    must("A5b  the in-gate cost at plan",         f"{FM.plan_roster()[2]:.2f}")
    must("A5b  pooling is shown as an upside",    f"{FM.pooled_upside()[2]:.2f}")
    must("A5b  the pooling saving is quantified", f"{FM.pooled_upside()[3]:.2f}")
    must("A5b  all-gig on the same leg",          f"{FM.gig_leg_cost():.2f}")
    must("A5b  the block shelf is priced",        f"{FM.shelf_handoff_value()[0]:.2f}")
    must("A5b  the closed-form rule is stated",   "R* =")

    # ---- A-5c  the brief's three micro-markets -------------------------------------
    # The four geometry figures AND the claim they support. A page can carry every correct
    # number and still make the wrong argument, so the ordering sentence is checked as prose
    # and the ordering itself is asserted in audit.py.
    _SEG = CS.segment_cpo()
    must("A5c  standard residential benchmark",   f"{_SEG['Standard 2-3 km residential']:.2f}")
    must("A5c  campus gate geometry, unconsolidated", f"{_SEG['Type A campus, gate-drop']:.2f}")
    must("A5c  urban PG cluster, doorstep",       f"{_SEG['Type B urban PG cluster']:.2f}")
    must("A5c  PG cluster, common-drop",          f"{_SEG['Type B PG cluster, common-drop']:.2f}")
    must("A5c  the plan is quoted with its basis", f"{SL.cost_legs()['total']:.2f}")
    must("A5c  the city leg is shown separately", f"{SL.cost_legs()['city']:.2f}")
    must("A5c  the in-gate leg is shown separately", f"{SL.cost_legs()['in_gate']:.2f}")
    must("A5c  the advantage is named as consolidation", "not the gate")
    must("A5c  PG demand is excluded from the base case", "EXCLUDED from")
    must("A5c  the PG position is Phase 2",       "Phase-2 adjacency")
    must("A5c  the promotion test is stated",     "common-drop feasibility")
    must("A5c  the brief's three segments are named", "paying-guest")
    must("S6   the runner row is the block-shelf unit cost", f"{FM.shelf_handoff_value()[2]:.2f}")
    must_not_on("S6   door-drop cost is not on the plan slide", f"{FM.shelf_handoff_value()[1]:.2f}")
    must("A6   the LEVEL basis is named",         f"{RC.AOV_UP:,.0f}")
    must("A6   the ANCHORED basis is named",      f"{RS.AOV_CEILING_LO:,.0f}")
    # These two used to be the literals "8.50" and "2.14". The deck hardcoded the same
    # literals, so checker and slide agreed with each other while both disagreed with the model
    # (turnover leg 8.44). A check whose expected value is typed rather than computed cannot
    # detect drift - it can only confirm that two copies of the same typo match.
    _du = RC.dupont(RC.AOV_HURDLE)
    must("A6b  the DuPont turnover leg",          f"{_du['capital_turn']:.2f}")
    must("A6b  the DuPont margin leg",            f"{_du['ebit_margin']*100:.2f}%")
    must("A6b  the day-count note is present",    f"{RC.DAYCOUNT_GAP:.2f}")
    must("A6b  capital employed",                 f"{RC.CE_BASE/1e5:,.1f} L")
    must("A6b  AOV at ROCE = 0",                  f"{RC.AOV_BREAKEVEN:,.0f}")
    must("A6b  AOV at the hurdle",                f"{RC.AOV_HURDLE:,.0f}")
    must("A6b  post-tax hurdle AOV",              f"{RC.AOV_HURDLE_POSTTAX:,.0f}")
    must("A6b  the like-for-like asset turn",     f"{RC.TURN_SLIDE4:.2f}")
    # A-5: working-capital constructs and the priced shocks, none of which were read before
    must("A5   working capital, adopted",         f"{WC.WC_ADOPTED/1e5:,.1f} L")
    must("A5   working capital, 12-day target",   f"{WC.WC_TARGET/1e5:,.1f} L")
    must("A5   the rejected COGS construct",      f"{WC.WC_OLD/1e5:,.1f} L")
    must("A5   shock, volume -30%",               f"{RS.AOV_VOLUME:,.0f}")
    must("A5   shock, shrinkage upper bound",     f"{RS.AOV_SHRINKAGE:,.0f}")
    must("A5   shock, gig levy",                  f"{RS.AOV_LEVY:,.0f}")
    must("A5   shock, calendar fragmentation",    f"{RS.AOV_FRAGMENTATION:,.0f}")
    must("A7   districts clearing the screen",    f"{AD.N_CANDIDATES}")
    must("A7b  the seven-week derivation",        "21 + 28 = 49")
    must("A9   the load-bearing paper",           "Kambli")
    must("A10  the gaps are counted",             "Seven gaps")

# ---- PHRASE BANS: the wording itself, in the built file --------------------------------
# The CLAIM checks in audit.py assert a property of the MODEL (argmin SLA is not the peak
# band). They do not read the deck, so the deck could say "fastest at peak" again while both
# CLAIM checks kept passing. These read the deck.
for _p in ("fastest at peak", "fastest at exam-night peak", "fastest and cheapest",       # scan:allow
           "the client manages", "client's own hurdle", "management's own",                # scan:allow
           "management's stated", "management has already", "management ceiling",         # scan:allow
           "the stated ceiling", "stated ceiling", "their stated ceiling"):                    # scan:allow
    must_not_on(f"BAN  '{_p}'", _p)
# and the corrected sentence must actually be there, not merely the banned one absent
# NOTE the numbering: the operating-model page is the deck's S6 and the FILE's slide 7.
must("S6  the corrected peak claim is present", "ost per order is lowest at exam-night peak")
# stale verification counts must not survive anywhere, screenshots included
for _s in ("311/311", "311 / 311", "80/80", "80 / 80", "83/83", "83 / 83"):
    must_not_on(f"BAN  stale verification count '{_s}'", _s)

# ---- nothing superseded may appear anywhere ---------------------------------------
must_not("Round 1 breakeven Rs528",   "Rs528")
must_not("Round 1 Minutes AOV Rs775", "Rs775")
must_not("Round 1 city asset turn",   "12.6x")
must_not("superseded working capital", "Rs95.7")
must_not("superseded breakeven Rs554", "Rs554")
must_not("stale NWC days",            "NWC 18d")
# ---- the pre-correction number set must not survive ANYWHERE, appendix included ----------
# Slides 15 and 18 carried a complete parallel set of these in hardcoded text through every
# passing layer, because no page-scoped check read those panels and the two that did compared
# a literal to a literal.
for _lbl, _v in (("pre-correction CE", "Rs325.0 L"), ("pre-correction WC", "Rs90.0 L"),
                 ("pre-correction WC target", "Rs77.1 L"), ("pre-correction rejected WC", "Rs117.8 L"),
                 ("pre-correction breakeven", "Rs578"), ("pre-correction hurdle", "Rs763"),
                 ("pre-correction post-tax", "Rs825"), ("pre-correction spine", "Rs580"),
                 ("pre-correction turnover leg", "8.50x"), ("pre-correction margin leg", "4.71%"),
                 ("pre-correction shock, volume", "Rs647"), ("pre-correction shock, shrinkage", "Rs623"),
                 ("pre-correction shock, levy", "Rs595"), ("pre-correction shock, fragmentation", "Rs589"),
                 ("pre-correction last mile", "Rs19.0")):
    must_not(f"superseded  {_lbl}", _v)
must_not("placeholder text",          "TBD")
must_not("cover team fields unfilled", "\u00ab")   # \u00abTEAM NAME\u00bb must be replaced before shipping
# ...and absence of the placeholder is not the same as PRESENCE of the name. An empty TEAM_NAME
# would clear the guillemets and ship a blank cover past the check above.
must("cover carries the team name", os.environ.get("TEAM_NAME") or P.TEAM_NAME)
must_not("placeholder text, lorem",   "Lorem")

# The count on A0 is only trustworthy if the constant it is printed from matches the number of
# checks this file actually runs. Asserting one without the other is how 83 survived on a slide
# while the checker ran 106.
# ---- HYPERLINKS MUST EXIST AND POINT SOMEWHERE REAL ------------------------------------
# The deck carried QR codes and ZERO hyperlink relationships across all 23 slides, so a judge
# reading on a laptop had to photograph their own screen to follow a link. Presence is asserted
# here because a QR is not a link, and nothing else in the stack can tell the difference.
_links = []
for _sl in Presentation(PPTX).slides:
    for _rid, _rel in _sl.part.rels.items():
        if _rel.reltype.endswith("/hyperlink"): _links.append(_rel.target_ref)
checks.append((len(_links) >= 6, "LINKS  the deck carries real hyperlinks",
               f"{len(_links)} hyperlink relationships"))
checks.append((all(u.startswith("https://") for u in _links),
               "LINKS  every hyperlink is absolute https",
               "all https" if all(u.startswith("https://") for u in _links)
               else [u for u in _links if not u.startswith("https://")][:2]))

# ---- NO CHECK MAY BE REGISTERED TWICE -------------------------------------------------
# A stray `for` loop once swallowed the ten appendix checks that followed it, because they were
# already indented and the loop body was a single line. They ran SIX times each - 50 phantom
# rows - and the deck printed "183 deck checks" against a checker that genuinely held 133.
# Every row still passed, so nothing looked wrong.
#
# The SELF count check below could not catch it, because the constant it compares against was
# being realigned to whatever the run produced. A number regenerated from reality cannot police
# reality. Registering the same (label, value) twice is never intentional here, so it is an error.
import collections as _co
_dup = {k: n for k, n in _co.Counter(tuple(str(x) for x in c[1:]) for c in checks).items() if n > 1}
checks.append((not _dup, "SELF  no check is registered twice",
               (f"{len(_dup)} duplicated, first: {list(_dup)[0][0]}" if _dup else "0 duplicates")))

# THE DELETED DEMAND CONSTANT MUST NOT REAPPEAR. campus_model carried a second
# orders/resident/day of 0.25, with a 5,600-resident cluster and a 28,000-resident state gate
# derived from it. Nothing imported them and no assertion read them, so they read as fact for
# weeks. params.py is now the only home; these three are the deck-side half of the guard.
for _lbl, _v in (("cluster", "5,600"), ("cluster, unpunctuated", "5600"),
                 ("state gate", "28,000")):
    must_not(f"superseded  pre-params {_lbl} residents", _v)

checks.append((len(checks) + 1 == DECK_CHECK_COUNT,
               "SELF  deck_checks.DECK_CHECK_COUNT matches this run",
               f"{DECK_CHECK_COUNT} (actual {len(checks) + 1})"))

if __name__ == "__main__":
    bad = [c for c in checks if not c[0]]
    print(f"\nVERIFYING {os.path.basename(PPTX)}  ({len(pages)} slides)\n" + "="*78)
    for ok, lab, needle in checks:
        print(("  OK  " if ok else "  FAIL") + f"  {lab:<44} {needle}")
    print(f"\n{len(checks)-len(bad)}/{len(checks)} deck checks pass")
    if bad: raise SystemExit("DECK VERIFICATION FAILED")
