"""
VERIFY THE COMPANION ARTEFACTS.

    python3 verify_artifacts.py

WHY THIS LAYER EXISTS. Four layers passed - model, HANDOFF, spec, built deck - while the
workbook said "of 760 districts", the notebooks carried saved outputs claiming 344/344 and
133/133, and the asset manifest pointed at screenshot files that had been renamed months
earlier. Everything the package SHIPS ALONGSIDE the deck was outside the scope of every check.

The pattern, met once more: **a check can only fail for something it reads.** Adding generators
for the workbook, the README and the manifest did not put them in scope; it only meant they were
regenerated. This layer reads them.

Scope: Campus_Store_Model.xlsx, the six notebooks, deck_asset_manifest.csv, and Model/data.
"""
import csv, json, os, re, subprocess, sys

HERE = os.path.dirname(os.path.abspath(__file__)); ROOT = os.path.dirname(HERE)
sys.path.insert(0, HERE)
import aishe_district as AD, check_counts as CC, params as P, roce as RC, working_capital as WC

checks = []
def chk(ok, label, detail=""):
    checks.append((bool(ok), label, detail))

# Any of these appearing in a companion artefact is a pre-correction figure.
SUPERSEDED = ["760 districts", "of 760", "760-district", "760 ", "311/311", "344/344", "133/133", "83/83", "80/80",
              "Rs325.0", "₹325.0", "Rs90.0 L", "₹90.0 L", "Rs578", "₹578", "Rs763", "₹763",
              "8.50x", "8.50×", "Rs647", "₹647", "Rs19.0", "₹19.0", "32.7%", "57.1%"]

# ---------------------------------------------------------------- the workbook
XL = os.path.join(ROOT, "Campus_Store_Model.xlsx")
if os.path.exists(XL):
    from openpyxl import load_workbook
    wb = load_workbook(XL)
    cells = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            for v in row:
                if v is not None: cells.append(str(v))
    workbook_blob = "\n".join(cells)
    chk(all(s not in workbook_blob for s in SUPERSEDED), "WORKBOOK  no superseded figure",
        [s for s in SUPERSEDED if s in workbook_blob][:3] or "clean")
    chk(str(AD.N_DISTRICTS) in workbook_blob, "WORKBOOK  carries the pair-safe district count",
        f"{AD.N_DISTRICTS}")
    chk(f"{RC.CE_BASE/1e5:.1f}" in workbook_blob, "WORKBOOK  capital employed matches the model",
        f"{RC.CE_BASE/1e5:.1f}")
    chk(f"{WC.WC_ADOPTED/1e5:.1f}" in workbook_blob, "WORKBOOK  working capital matches the model",
        f"{WC.WC_ADOPTED/1e5:.1f}")
    chk(str(CC.AUDIT_COUNT) in workbook_blob, "WORKBOOK  release-check counts match the run",
        f"{CC.AUDIT_COUNT}")
    chk(P.TEAM_NAME in workbook_blob, "WORKBOOK  names the team", P.TEAM_NAME)
else:
    chk(False, "WORKBOOK  present", "missing")

# ---------------------------------------------------------------- the notebooks
NBD = next((d for d in (os.path.join(ROOT, "notebooks"),
                        os.path.join(ROOT, "_release_repo", "notebooks")) if os.path.isdir(d)), None)
if NBD:
    nbs = sorted(f for f in os.listdir(NBD) if f.endswith(".ipynb"))
    # Structure alone is not enough. L5 once passed every package verifier while its public
    # cell failed because sla.py's display loop still unpacked the old eight-column row after
    # volume_weighted() had grown to ten columns. Execute every non-recursive notebook
    # entrypoint here; L1 is this verifier's parent (run_all.py), so assert that it is wired to
    # the supported command rather than recursively launching it.
    _entrypoints = {
        "L2_Dead_Zone_Solver.ipynb": ["solver.py"],
        "L3_Return_Model.ipynb": ["roce.py"],
        "L4_District_Screen.ipynb": ["aishe_district.py"],
        "L5_Fulfilment_Model.ipynb": ["sla.py", "fleet_mix.py"],
        "L6_Basket_Regression.ipynb": ["basket.py"],
    }
    _entry_fail = []
    _display_missing = []
    for _nb in nbs:
        _raw = open(os.path.join(NBD, _nb), encoding="utf-8").read()
        if "capture_output=True" not in _raw or "result.check_returncode()" not in _raw:
            _display_missing.append(_nb)
    for _nb, _scripts in _entrypoints.items():
        for _script in _scripts:
            _run = subprocess.run([sys.executable, os.path.join(HERE, _script)],
                                  cwd=ROOT, stdout=subprocess.DEVNULL,
                                  stderr=subprocess.PIPE, text=True)
            if _run.returncode:
                _entry_fail.append(f"{_nb}:{_script}:exit {_run.returncode}")
    _l1 = os.path.join(NBD, "L1_Audit_Verification.ipynb")
    _l1_wired = os.path.exists(_l1) and "Model/run_all.py" in open(_l1, encoding="utf-8").read()
    chk(len(nbs) == 6 and _l1_wired and not _entry_fail and not _display_missing,
        "NOTEBOOKS  entrypoints execute and surface their output",
        ((_entry_fail + [f"{n}:output hidden" for n in _display_missing])[:2]
         if (_entry_fail or _display_missing) else f"{len(nbs)} found; all commands visible, exit 0"))
    saved_out, homepaths, stale = [], [], []
    for f in nbs:
        nb = json.load(open(os.path.join(NBD, f), encoding="utf-8"))
        for c in nb.get("cells", []):
            if c.get("outputs"): saved_out.append(f)
            body = "".join(c.get("source", []))
            for o in c.get("outputs", []):
                body += "".join(o.get("text", [])) if "text" in o else json.dumps(o.get("data", {}))
            if re.search(r"/Users/[A-Za-z0-9_.-]+", body): homepaths.append(f)
            for s in SUPERSEDED:
                if s in body: stale.append(f"{f}:{s}")
    # A saved output is a screenshot in JSON: it looks like evidence and never regenerates.
    chk(not saved_out, "NOTEBOOKS  no saved outputs", sorted(set(saved_out))[:3] or "clean")
    chk(not homepaths, "NOTEBOOKS  no local home path leaked", sorted(set(homepaths))[:3] or "clean")
    chk(not stale, "NOTEBOOKS  no superseded figure", sorted(set(stale))[:3] or "clean")
else:
    chk(False, "NOTEBOOKS  directory present", "missing")

# ---------------------------------------------------------------- the asset manifest
MF = os.path.join(HERE, "assets", "deck_asset_manifest.csv")
if os.path.exists(MF):
    rows = list(csv.reader(open(MF, encoding="utf-8")))
    manifest_blob = "\n".join(",".join(r) for r in rows)
    chk(all(s not in manifest_blob for s in SUPERSEDED), "MANIFEST  no superseded figure",
        [s for s in SUPERSEDED if s in manifest_blob][:3] or "clean")
    chk(str(CC.AUDIT_COUNT) in manifest_blob, "MANIFEST  audit count matches the run", f"{CC.AUDIT_COUNT}")
    chk(str(AD.N_DISTRICTS) in manifest_blob, "MANIFEST  district count matches the model", f"{AD.N_DISTRICTS}")
    # every asset it names must actually exist - it pointed at renamed screenshots for weeks
    missing = []
    for r in rows[1:]:
        for cell in r[-2:]:
            for a in [x.strip() for x in cell.split(";") if x.strip()]:
                if not os.path.exists(os.path.join(HERE, a)): missing.append(a)
    chk(not missing, "MANIFEST  every asset it names exists", missing[:3] or "clean")
else:
    chk(False, "MANIFEST  present", "missing")

# ---------------------------------------------------------------- the layer count
# Adding the fifth verifier turned "four layers" into a stale value in eight places at once,
# because the number was typed everywhere it appeared instead of counted. This bans the outdated
# wording wherever the package describes itself.
_N = len(CC.DEFINED)
_STALE_WORDING = ["ALL FOUR LAYERS", "All four layers", "all four layers", "four layers",
                  "All four rows", "all four verification layers", "four verification commands"]
_where = []
for _p, _read in (("README.md", lambda p: open(p, encoding="utf-8").read()),
                  ("Campus_Store_Model.xlsx", None)):
    pass
_txt = ""
_rd = os.path.join(ROOT, "README.md")
if os.path.exists(_rd): _txt += open(_rd, encoding="utf-8").read()
if os.path.exists(XL): _txt += workbook_blob
if NBD:
    for f in sorted(os.listdir(NBD)):
        if f.endswith(".ipynb"): _txt += open(os.path.join(NBD, f), encoding="utf-8").read()
# THE DECK IS THE POINT. The first version of this check read the README, the workbook and the
# notebooks and stopped there - so slide 10's "checked by four scripts" and "Each of the four"
# were never in scope, and the check reported clean while the deck was wrong. A ban that does not
# read the artefact it is protecting is decoration.
_DECK = os.path.join(ROOT, P.FINAL_DECK)
_deck_txt = ""
if os.path.exists(_DECK):
    from pptx import Presentation
    for _s in Presentation(_DECK).slides:
        for _sh in _s.shapes:
            if _sh.has_text_frame: _deck_txt += _sh.text_frame.text + "\n"
            if _sh.has_table:
                for _r in _sh.table.rows: _deck_txt += " | ".join(c.text for c in _r.cells) + "\n"
_STALE_WORDING += ["four scripts", "Each of the four", "each of the four", "the four is a command"]
_hit  = [w for w in _STALE_WORDING if w in _txt] if _N != 4 else []
_dhit = [w for w in _STALE_WORDING if w in _deck_txt] if _N != 4 else []
chk(not _hit,  "LAYERS  no outdated layer-count wording in docs/workbook/notebooks",
    _hit[:3] or f"{_N} layers, consistent")
chk(not _dhit, "LAYERS  no outdated layer-count wording IN THE DECK",
    _dhit[:3] or f"{_N} layers, consistent")
# and the deck must name every verifier it claims to run
_VERIFIERS = ("audit.py","verify_docs.py","verify_spec.py","verify_deck.py","verify_artifacts.py")
_missing = [k for k in _VERIFIERS if k not in _deck_txt]
chk(not _missing, "LAYERS  the deck names every verifier", _missing or f"all {_N} named")
# The README says how many layers there are AND lists them; both must agree, and both must agree
# with what actually runs. The two lists were maintained separately, which is how "five layers"
# ended up above a table of four.
_rm = open(_rd, encoding="utf-8").read() if os.path.exists(_rd) else ""
_rmiss = [k for k in _VERIFIERS if k not in _rm]
chk(not _rmiss, "LAYERS  the README lists every verifier", _rmiss or f"all {_N} listed")
# and PUSH_INSTRUCTIONS, if present, must quote the counts that actually run
# ALWAYS a check, never a conditional one. PUSH_INSTRUCTIONS.md is internal and absent from the
# published tree by design, so a conditional made the artefact COUNT differ between the working
# folder (22) and the staged tree (21) - which meant check_counts could not certify the staged
# tree at all. A verifier whose check count depends on which files happen to exist cannot be
# compared against a constant.
_pi = os.path.join(ROOT, "PUSH_INSTRUCTIONS.md")
_want = " / ".join(str(CC.DEFINED[k]) for k in ("audit", "docs", "spec", "deck", "artefacts"))
chk((_want in open(_pi, encoding="utf-8").read()) if os.path.exists(_pi) else True,
    "LAYERS  push instructions quote the live counts",
    _want if os.path.exists(_pi) else "not in this tree (internal, by design)")

# ---------------------------------------------------------------- the STAGED tree
# THE TREE THAT ACTUALLY GETS PUSHED, CHECKED AGAINST THE ONE THAT WAS VERIFIED.
# _release_repo/ was assembled and refreshed BY HAND. After the last round the working folder
# verified at one set of counts while the staged tree still held the PREVIOUS build's deck,
# workbook, README and four .py files. Nothing could detect it: every verifier ran inside one
# tree or the other and never compared them. That is the same defect the package was built to
# eliminate, sitting in the one tree a judge reads.
#
# It is now derived by stage_release.py from a manifest, and this asserts the derivation is
# current. In the staged tree itself there is no _release_repo/, so the check passes on absence
# and the artefact COUNT stays the same in both trees.
import hashlib as _hl
def _md5(p_):
    with open(p_, "rb") as fh: return _hl.md5(fh.read()).hexdigest()
_stage = os.path.join(ROOT, "_release_repo")
_stale = []
if os.path.isdir(_stage):
    try:
        import stage_release as _SR
        _pairs = [(os.path.join(ROOT, f), os.path.join(_stage, f)) for f in _SR.ROOT_FILES]
        for _t in _SR.TREES:
            for _dp, _dn, _fn in os.walk(os.path.join(ROOT, _t)):
                _dn[:] = [d for d in _dn if d not in _SR.EXCLUDE]
                for _f in _fn:
                    if _f in _SR.EXCLUDE or _f.startswith(".") or _f.endswith((".pyc", ".pyo")):
                        continue
                    _a = os.path.join(_dp, _f)
                    _pairs.append((_a, os.path.join(_stage, os.path.relpath(_a, ROOT))))
        for _a, _b in _pairs:
            if not os.path.exists(_b) or _md5(_a) != _md5(_b):
                _stale.append(os.path.relpath(_a, ROOT))
    except Exception as _e:
        _stale = [f"could not compare: {_e}"]
chk(not _stale, "STAGE  the staged tree matches the verified tree",
    (f"{len(_stale)} stale, first: {_stale[0]}" if _stale
     else ("current" if os.path.isdir(_stage) else "no staged tree here (this IS the staged tree)")))

# ---------------------------------------------------------------- the public data
DATA = os.path.join(HERE, "data")
chk(os.path.exists(os.path.join(DATA, "README.md")), "DATA  provenance file present")
chk(int(AD.DIST["District"].isna().sum()) == 0, "DATA  no blank district in the universe")
chk(AD.N_DISTRICTS > AD.N_DISTRICT_NAMES, "DATA  denominator counts pairs, not names",
    f"{AD.N_DISTRICTS} pairs > {AD.N_DISTRICT_NAMES} names")

if __name__ == "__main__":
    bad = [c for c in checks if not c[0]]
    print(f"\nVERIFYING THE COMPANION ARTEFACTS  ({os.path.basename(XL)}, notebooks, manifest, data)")
    print("=" * 86)
    for ok, lab, detail in checks:
        print(("  OK  " if ok else "  FAIL") + f"  {lab:<52} {detail}")
    print(f"\n{len(checks)-len(bad)}/{len(checks)} artefact checks pass")
    if bad: raise SystemExit("ARTEFACT VERIFICATION FAILED")
