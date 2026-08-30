#!/usr/bin/env python3
"""
THE RELEASE SEQUENCE:  build -> verify -> stamp -> verify the exact final file.

    TEAM_NAME="..." TEAM_ID="..." python3 Model/release.py

WHY IT IS FOUR STEPS AND NOT TWO. Slide A-1 reports the package's own verification result,
which is a self-reference: the deck cannot contain the outcome of verifying itself unless the
outcome is computed first and stamped in afterwards. The earlier build did this by carrying the
counts as constants, so the slide printed "116 / 116" on a build the verifier was actually
failing 115/116. A number a document states about itself must come from a run.

  1 BUILD      with no results file present, A-1 prints "N checks defined" - a claim about how
               many checks exist, which is true whether or not they pass.
  2 VERIFY     run_all.py runs every layer and writes _verification.json.
  3a CAPTURE   photograph THAT run into A-1's two terminal images. They were previously
               captured by hand and went stale exactly like a literal - showing 311/311 and
               80/80 beside tiles reading 336 and 116, on the page arguing nothing drifts.
               An image cannot be asserted, so it has to be generated.
  3 STAMP      rebuild. A-1 now prints "N / N", but ONLY because step 2 recorded a clean run.
  4 RE-VERIFY  run every layer again over the EXACT final package. Steps 3a onward regenerate
               the screenshots, README, workbook, notebooks and manifest, so the artefacts step 2
               checked are not the artefacts that ship. This step is what makes the claim true
               rather than merely consistent.

Any step failing stops the sequence and leaves the results file recording the failure, so the
next build reverts to "checks defined" instead of keeping a stale pass.
"""
import os, subprocess, sys

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
RESULTS = os.path.join(HERE, "_verification.json")

def run(cmd, label):
    print(f"\n>>> {label}")
    r = subprocess.run(cmd, cwd=ROOT)
    if r.returncode != 0:
        if r.returncode == 1 and "BUILD" in label:
            print("    (If this is a clean clone, the organisers' template is absent by design —"
                  "\n     run_all.py verifies without it.)")
        print(f"\n!!! {label} FAILED - stopping. The deck does not get to claim a pass.")
        raise SystemExit(r.returncode)

def lock_files():
    """Office writes ~$name.pptx while a file is open. It carries the editor's NAME, it ends up
    committed to the repo, and - worse - its presence means PowerPoint holds the deck open and
    may write over anything this script builds. Checked before the build, not after."""
    import glob
    return [os.path.basename(p) for p in glob.glob(os.path.join(ROOT, "~$*"))]

def main():
    py = sys.executable
    stale = lock_files()
    if stale:
        print("\n!!! Office lock file(s) present: " + ", ".join(stale))
        print("    The deck is open in PowerPoint. Close it before releasing: a lock file gets")
        print("    committed to the repo, it carries the editor's name, and PowerPoint can")
        print("    overwrite the file this script is about to build.")
        raise SystemExit(2)
    # Step 1 must build with NO standing claim. Blanked rather than deleted: an empty record
    # is unambiguous, and some mounts refuse deletion.
    with open(RESULTS, "w", encoding="utf-8") as fh: fh.write("{}\n")
    run([py, "Model/build_full.py", "light"], "1 BUILD  (A-1 prints 'checks defined')")
    # STAGE BEFORE VERIFY, AND AGAIN BEFORE RE-VERIFY. verify_artifacts asserts that
    # _release_repo/ matches the tree that just verified, and step 1 has this moment made the
    # working tree newer than the staged one - so staging has to bracket both verification
    # passes, not sit between them. Staging is an rsync from a manifest; running it twice costs
    # nothing and means neither pass can be run against a stale published tree.
    run([py, "Model/stage_release.py"],       "1a STAGE _release_repo/ from the manifest")
    run([py, "Model/run_all.py"],             "2 VERIFY (records _verification.json)")
    run([py, "Model/make_screenshots.py"],    "3a CAPTURE the clean run into A-1's screenshots")
    run([py, "Model/make_readme.py"],         "3a' REGENERATE README.md from that run")
    run([py, "Model/make_workbook.py"],       "3a\" REGENERATE Campus_Store_Model.xlsx from the model")
    run([py, "Model/clean_notebooks.py"],     "3a\"' CLEAN the notebooks (outputs, paths, typed counts)")
    run([py, "Model/sync_push_instructions.py"], "3a\"'' SYNC the internal push note to the live counts")
    run([py, "Model/build_full.py", "light"], "3b STAMP  (A-1 prints passed/total)")
    run([py, "Model/stage_release.py"],       "3c STAGE the stamped package into _release_repo/")
    # RE-VERIFY THE WHOLE PACKAGE, not just the deck. Steps 3a-3a"' regenerate the
    # screenshots, README, workbook, notebooks and manifest AFTER step 2 checked them, so a
    # verify_deck-only step 4 left every regenerated artefact unverified in its final form - the
    # same architecture gap the fifth layer was added to close, reopened one step later.
    run([py, "Model/run_all.py"],             "4 RE-VERIFY the exact final package")

    # And the claim must be the EARNED one. verify_deck only checks the form is legal, because
    # it cannot see its own outcome. This is the step that requires the passed form, on the
    # exact file that carries it, after that file has just verified clean.
    from pptx import Presentation
    import check_counts as _CC0, importlib as _il
    _il.reload(_CC0)
    _txt = " ".join(sh.text_frame.text for sl in Presentation(
        os.path.join(ROOT, _CC0.__dict__.get("FINAL", "") or __import__("params").FINAL_DECK)
    ).slides for sh in sl.shapes if sh.has_text_frame)
    _want = f"{_CC0.AUDIT_COUNT} / {_CC0.AUDIT_COUNT}"
    if _want not in _txt:
        print(f"\n!!! the stamped deck does not carry the earned claim {_want!r} - NOT RELEASED.")
        raise SystemExit(1)
    print(f"    stamped claim {_want!r} is present on the verified file.")
    import importlib, sys as _s
    _s.path.insert(0, HERE); import check_counts as CC; importlib.reload(CC)
    print("\n" + "=" * 78)
    print("  RELEASED." if _CC0.all_verified() else "  NOT RELEASED - results file is not clean.")
    for k in _CC0.DEFINED:          # every layer, not a list that has to be maintained
        print(f"    {k:<10} {_CC0.tile(k)}")
    print("=" * 78)
    return 0 if _CC0.all_verified() else 1

if __name__ == "__main__":
    raise SystemExit(main())
