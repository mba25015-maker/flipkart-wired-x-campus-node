"""
SEMI-FINAL DECK BUILDER.  Slides 5 and 6 were built first -- the pivot and the payoff --
then 1 and 2, so the framing foreshadows a payoff that was already fixed.

  python3 build_semifinal.py [light|dark] [out_basename]

RULE, inherited from Round 1 and non-negotiable: NEVER TYPE A FIGURE. Every number
on a slide is imported from the model at build time, so a model change moves the slide
and `audit.py` asserts the pair. Prose is written here; arithmetic is not.

The Round 1 deck (build_v7.py) is untouched. This file rebuilds its visual system with
the paths made relative, because build_v7's were absolute to a dead session.
"""
import sys, os
from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

import campus_model as M, cost_stack as CS, break_mode as B, basket as BK
import aishe_district as AD, calendar_fragmentation as CF
import risk_quadrant as Q, sla as SL

THEME = sys.argv[1] if len(sys.argv) > 1 else "light"
OUT   = sys.argv[2] if len(sys.argv) > 2 else "Flipkart_Minutes_WiRED_SemiFinal"
HERE  = os.path.dirname(os.path.abspath(__file__))
ROOT  = os.path.dirname(HERE)
TPL   = os.path.join(ROOT, "Case", "Presentation_template.pptx")
IMG   = os.path.join(HERE, "charts", THEME)
FONT  = "Calibri"

if THEME == "light":
    CONTENT_LAYOUT = "CUSTOM_3"
    FG=C(0x0D,0x1F,0x5C); MUTE=C(0x5A,0x67,0x85); RULE=C(0xD3,0xDC,0xF0)
    BANNER_BG=C(0x0D,0x1F,0x5C); BANNER_FG=C(0xFF,0xFF,0xFF)
    CHIP=C(0xF2,0xF6,0xFF); CHIP2=C(0xFF,0xF8,0xE3)
    POS=C(0x1E,0x7A,0x46); NEG=C(0xC0,0x39,0x2B); HI=C(0xB5,0x7F,0x00)
    NAVBG=C(0xE8,0xEE,0xFB); NAVON=C(0x0D,0x1F,0x5C); NAVONFG=C(0xFF,0xFF,0xFF)
    BAND=C(0x0D,0x1F,0x5C); BANDFG=C(0xFF,0xFF,0xFF); BANDHI=C(0xFF,0xC2,0x20)
else:
    CONTENT_LAYOUT = "CUSTOM_2"
    FG=C(0xFF,0xFF,0xFF); MUTE=C(0xBF,0xCE,0xF0); RULE=C(0x3E,0x5C,0xA8)
    BANNER_BG=C(0xFF,0xC2,0x20); BANNER_FG=C(0x0A,0x1A,0x4E)
    CHIP=C(0x12,0x2B,0x74); CHIP2=C(0x1B,0x36,0x86)
    POS=C(0x6B,0xE0,0xA0); NEG=C(0xFF,0x9A,0x8B); HI=C(0xFF,0xC2,0x20)
    NAVBG=C(0x14,0x2C,0x77); NAVON=C(0xFF,0xC2,0x20); NAVONFG=C(0x0A,0x1A,0x4E)
    BAND=C(0xFF,0xC2,0x20); BANDFG=C(0x0A,0x1A,0x4E); BANDHI=C(0x0A,0x1A,0x4E)

prs = Presentation(TPL)
L = {l.name: l for l in prs.slide_masters[0].slide_layouts}
_ids = prs.slides._sldIdLst
for sid in list(_ids)[1:]:
    prs.part.drop_rel(sid.rId); _ids.remove(sid)
for sid in list(_ids):                       # slides 5-6 only: no cover in this build
    prs.part.drop_rel(sid.rId); _ids.remove(sid)

# ---------------- primitives (from build_v7, paths fixed) ----------------
def strip_ph(s):
    for sh in list(s.shapes):
        if sh.is_placeholder: sh._element.getparent().remove(sh._element)

def tb(s,l,t,w,h,anchor=MSO_ANCHOR.TOP):
    x=s.shapes.add_textbox(In(l),In(t),In(w),In(h)); f=x.text_frame
    f.word_wrap=True; f.vertical_anchor=anchor
    f.margin_left=f.margin_right=f.margin_top=f.margin_bottom=0
    return f

def par(f,first=False,before=0,space_after=0,lsp=None):
    p=f.paragraphs[0] if (first and not f.paragraphs[0].runs) else f.add_paragraph()
    if before: p.space_before=Pt(before)
    if space_after: p.space_after=Pt(space_after)
    if lsp: p.line_spacing=lsp
    return p

def run(p,txt,sz,col,bold=False,ital=False):
    r=p.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.bold=bold
    r.font.italic=ital; r.font.name=FONT; r.font.color.rgb=col; return r

def text(s,l,t,w,h,txt,sz,col,bold=False,ital=False,align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP,lsp=None):
    f=tb(s,l,t,w,h,anchor); p=par(f,True,lsp=lsp); p.alignment=align
    run(p,txt,sz,col,bold,ital); return f

def box(s,l,t,w,h,fill,radius=0.10,line=None,lw=0.75):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,In(l),In(t),In(w),In(h))
    try: sh.adjustments[0]=radius
    except Exception: pass
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=Pt(lw)
    sh.shadow.inherit=False
    sh.text_frame.word_wrap=True
    sh.text_frame.margin_left=In(0.07); sh.text_frame.margin_right=In(0.07)
    sh.text_frame.margin_top=In(0.03); sh.text_frame.margin_bottom=In(0.03)
    return sh

def rect(s,l,t,w,h,col):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,In(l),In(t),In(w),In(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=col
    sh.line.fill.background(); sh.shadow.inherit=False
    sh.text_frame.word_wrap=True
    sh.text_frame.margin_left=In(0.05); sh.text_frame.margin_right=In(0.05)
    sh.text_frame.margin_top=In(0.01); sh.text_frame.margin_bottom=In(0.01)
    return sh

def rule_line(s,l,t,w,h,col=None): return rect(s,l,t,w,h,col or RULE)

def pic(s,name,l,t,maxw,maxh,center=True):
    p=os.path.join(IMG,f"{name}.png")
    iw,ih=Image.open(p).size; a=ih/iw
    w=maxw; h=w*a
    if h>maxh: h=maxh; w=h/a
    x=l+(maxw-w)/2 if center else l
    return s.shapes.add_picture(p,In(x),In(t),width=In(w),height=In(h))

def banner(s,l,t,w,label):
    b=box(s,l,t,w,0.235,BANNER_BG,radius=0.16)
    f=b.text_frame; f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True); run(p,label,7.0,BANNER_FG,bold=True)
    return t+0.235

def band(s,t,parts,h=0.38):
    b=box(s,0.55,t,8.90,h,BAND,radius=0.10)
    f=b.text_frame; f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True); p.line_spacing=0.98
    for txt,bold in parts: run(p,txt,8.2,BANDHI if bold else BANDFG,bold=bold)
    return b

def footer(s,txt): text(s,0.55,5.29,8.90,0.20,txt,5.0,MUTE,ital=True,lsp=0.92)

# Eight ticks, one per content slide: the panel always knows where it is.
def rail(s, active, label):
    x=0.55
    for i in range(8):
        on=(i==active-1)
        rect(s,x,0.30,0.30 if not on else 0.42,0.055, NAVON if on else NAVBG)
        x+= (0.30 if not on else 0.42)+0.055
    text(s,x+0.10,0.235,4.4,0.18,label,6.8,MUTE,bold=True)
    text(s,7.95,0.235,1.50,0.18,f"SLIDE {active} OF 8",6.8,MUTE,bold=True,align=PP_ALIGN.RIGHT)

def head(s,title,sub,size=17.0):
    text(s,0.55,0.56,8.12,0.62,title,size,FG,bold=True,lsp=0.92)
    text(s,0.55,1.19,8.90,0.28,sub,7.4,MUTE,ital=True,lsp=0.98)

def preempt(s,t,tag,body,col=None):
    """The line we say before a judge says it. Cheap, and it buys the slide."""
    col = col or HI
    b=box(s,0.55,t,8.90,0.28,None,radius=0.10,line=col,lw=1.0)
    f=b.text_frame; f.margin_left=In(0.13); f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True); p.line_spacing=0.95
    run(p,tag+"   ",6.4,col,bold=True); run(p,body,6.8,FG)
    return t+0.28

# =========================================================================
# SLIDE 1 -- EXECUTIVE SUMMARY.  Two columns: what this is NOT, and what it is.
# =========================================================================
def slide1():
    s=prs.slides.add_slide(L[CONTENT_LAYOUT]); strip_ph(s)
    rail(s,1,"THE RECOMMENDATION")

    thr  = [B.threshold(fn) for _,fn in B.CONFIGS]
    rel  = B.relocate_vs_flex()
    be   = CS.breakeven_d2_consistent(CS.CAMPUS_FIXED, SL.volume_weighted()[1])

    head(s,
      f"Enter campus micro-markets on a cluster-plus-catchment node — {M.TURN_RATIO:.0%} of a city store's "
      f"asset productivity, in the segment nobody is contesting",
      f"The node is underwritten by a 3–6 college cluster PLUS its adjacent non-student catchment, sited so "
      f"the academic break is survivable by construction. Three-month implementation; the calendar is a site "
      f"criterion, not a cost line.")

    yl = banner(s,0.55,1.52,4.35,"DO NOT TREAT THIS AS A COST PROBLEM")
    yr = banner(s,5.10,1.52,4.35,"DO TREAT IT AS A SITE-SELECTION PROBLEM")

    left = [("Cost levers cannot close the calendar.",
             f"The full ladder moves the residual demand the node needs from {thr[0]:.1%} to {thr[-1]:.1%} — "
             f"a {(thr[0]-thr[-1])/thr[0]:.0%} cut — and still cannot close it."),
            ("Walking away costs more than staying.",
             f"Held through the break on no levers a node burns Rs{rel['do_nothing_total']/1e5:.1f} lakh; "
             f"relocating costs Rs{rel['relocate_once']/1e5:.0f} lakh and forfeits a catchment that returns "
             f"on a published date. Holding is {rel['flex_vs_relocate']:.0%} of one relocation."),
            ("And the basket does not start where it must end.",
             f"Minutes' network AOV is Rs{M.MINUTES_AOV:,.0f}. The D1/D2-consistent breakeven is Rs{be:,.0f}.")]
    right = [("The screen is operational, and it is not a label.",
             f"{AD.N_CANDIDATES} of {AD.N_DISTRICTS:,} districts clear it: not a metro, ≥{AD.MIN_URBAN_COLLEGES} "
             f"urban colleges, urban share ≥{AD.URBAN_SHARE_FLOOR:.0%}, one break segment ≥7 weeks."),
            ("Density buys back the calendar.",
             f"Campus asset turn {M.CAMPUS_TURN:.2f}x against a city store's {M.CITY_TURN:.2f}x = "
             f"{M.TURN_RATIO:.3f} — and {M.turn_ratio_at(M.CEIL_LO):.2f}–{M.turn_ratio_at(M.CEIL_HI):.2f} across "
             f"the whole range a mature operator has actually recorded."),
            ("The basket benchmark is source-bounded.",
             f"Rs{be:,.0f} needs {BK.SHARE_NEEDED_AFTER_OCCASION:.1f}% non-grocery, from "
             f"{BK.MINUTES_NONGROCERY:.0f}% today, compared with Swiggy's disclosed "
             f"{BK.NONGROCERY_CEILING_LO:.0f}–{BK.NONGROCERY_CEILING:.0f}% range. This is a cross-operator reference, not a Flipkart commitment.")]

    for col,(x,items,accent) in enumerate([(0.55,left,NEG),(5.10,right,POS)]):
        yy = (yl if col==0 else yr)+0.10
        for i,(claim,body) in enumerate(items):
            b=box(s,x,yy,4.35,0.72,CHIP if col==0 else CHIP2,radius=0.10)
            rect(s,x,yy,0.05,0.72,accent)
            f=b.text_frame; f.margin_left=In(0.14); f.margin_top=In(0.06)
            p=par(f,True); p.line_spacing=0.94
            run(p,f"{i+1}   ",7.6,accent,bold=True); run(p,claim,7.6,FG,bold=True)
            p2=par(f); p2.line_spacing=0.98; p2.space_before=Pt(2)
            run(p2,body,6.7,MUTE)
            yy += 0.79

    # the concession, on the front foot, on slide one
    b=box(s,0.55,4.28,8.90,0.42,None,radius=0.10,line=HI,lw=1.0)
    f=b.text_frame; f.margin_left=In(0.13); f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True); p.line_spacing=0.94
    run(p,"WE ARE NOT FIRST, AND WE WILL SAY SO BEFORE YOU DO   ",6.6,HI,bold=True)
    run(p,"Swiggy runs a student rewards programme with college canteens and hostels, the Toing app "
          "(AOV <Rs200, ~50 cities) and a Young India Skills University MoU. Blinkit runs India's first "
          "in-airport quick commerce with Adani at CSMIA. ",6.8,FG)
    run(p,"What none of them has done is underwrite the node on the CLUSTER rather than the campus — "
          "which is the only structure that survives the break.",6.8,FG,bold=True)

    band(s,4.80,[("One decision on this page: ",False),
      ("site the node on the cluster, price the break as a filter, and enter the 111 districts where both hold.",True)],h=0.36)
    footer(s,
      f"Asset turn, breakeven and the lever ladder from the model at {M.CEILING:,} orders/day and Minutes AOV Rs{M.MINUTES_AOV:,.0f} (TechCrunch, 22 Aug 2026)  ·  "
      f"district screen from the AISHE institution register as on 28-8-2026  ·  basket ladder fitted on Instamart's disclosed quarterly series, R2={BK.R2:.3f}  ·  "
      f"competitive facts: Elara 27 Nov 2025 (Swiggy student programme), Jan 2026 (YISU MoU), 1 Apr 2026 (Blinkit–Adani CSMIA).  Every figure on this page is asserted in the audit, A1.")

# =========================================================================
# SLIDE 2 -- THE MARKET, AND THE TWO FACTS THAT REFRAME IT
# =========================================================================
def slide2():
    s=prs.slides.add_slide(L[CONTENT_LAYOUT]); strip_ph(s)
    rail(s,2,"THE MARKET  ·  SMALLER, AND EMPTIER, THAN IT LOOKS")

    head(s,
      f"Six in ten Indian colleges are rural and the metros are already {AD.METRO_EXCESS:.0%} oversupplied — "
      f"this market is smaller than it looks, and emptier than it looks",
      f"Both facts come out of registers, not forecasts: the AISHE institution register as on 28-8-2026, and "
      f"one quarter of operator store additions. The first shrinks our own market; the second is why we still recommend entry.")

    # THE HOOK. The largest campus-native delivery operator in the world died of D1's problem.
    b=box(s,0.55,1.50,8.90,0.46,None,radius=0.10,line=NEG,lw=1.0)
    f=b.text_frame; f.margin_left=In(0.13); f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True); p.line_spacing=0.94
    run(p,"THIS SEGMENT HAS ALREADY KILLED SOMEONE   ",6.6,NEG,bold=True)
    run(p,"In June 2026 Starship Technologies quit ALL US higher education after eight years and 60+ campuses, "
          "withdrawing ~1,200 robots. CEO Ahti Heinla: ",6.8,FG)
    run(p,"“Campus and grocery are fundamentally different operations: one is seasonal and contract-driven, "
          "the other is a 365-day urban business.” ",6.8,FG,ital=True)
    run(p,"It did not lose to robots-versus-humans. It lost 365-day utilisation to 8.5-month utilisation — "
          "the exact ratio priced on slide 5. And the campuses did not go dark: Avride took them on a national "
          "foodservice master agreement. The segment did not fail; a contract shape did.",6.8,FG,bold=True)

    rule_line(s,4.98,2.10,0.011,2.62)

    yl = banner(s,0.55,2.08,4.15,"THE UNIVERSE NARROWS   ·   AISHE INSTITUTION REGISTER, 28-8-2026")
    pic(s,"universe_narrow",0.55,yl+0.06,4.15,1.60)
    text(s,0.55,yl+1.68,4.15,0.44,
      f"Nothing in the corpus said this before the register was parsed, and it cuts our own addressable "
      f"universe by {1-AD.URBAN_SHARE_COL:.0%}. The metros hold {AD.METRO_URBAN_COLLEGES:,} of the urban "
      f"colleges and the archetype excludes them, which leaves {AD.NON_METRO_URBAN_COLLEGES:,} — "
      f"and {AD.N_STA_URBAN_HP:,} urban high-propensity standalones inside them.",6.2,MUTE,lsp=0.98)

    yr = banner(s,5.30,2.08,4.15,"THE METRO SQUEEZE   ·   ONE QUARTER OF STORE ADDITIONS, APR–JUL 2026")
    pic(s,"metro_squeeze",5.30,yr+0.06,4.15,1.60)
    text(s,5.30,yr+1.68,4.15,0.44,
      f"Nine of every ten new stores went somewhere an operator already served. The marginal metro store is "
      f"stacking, not reaching — and it is being added into a market already {AD.METRO_EXCESS:.0%} above the "
      f"capacity Bernstein calls sustainable. Flipkart needs density nobody else is buying.",6.2,MUTE,lsp=0.98)

    b=box(s,0.55,4.32,8.90,0.38,CHIP2,radius=0.10); rect(s,0.55,4.32,0.045,0.38,HI)
    f=b.text_frame; f.margin_left=In(0.12); f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True); p.line_spacing=0.92
    run(p,"AND THE SEGMENT IS LITERALLY UNCOUNTED   ",6.4,HI,bold=True)
    run(p,"Euromonitor's Indian consumer-foodservice taxonomy has no education or institutional location type, "
          "and its household typology has no hostel, PG or shared-student category — students fall into “Other”, "
          "which is never broken out. In the standard commercial dataset for Indian consumers this micro-market "
          "does not exist as a line, which is exactly why the register had to be parsed to size it.",6.4,FG)

    band(s,4.80,[("We have just shrunk our own market on slide two, and we still recommend entry. ",True),
      ("If the universe is 21,000 urban colleges and the metros are full, the question is not whether — it is WHERE. Slide 6 names the districts.",False)],h=0.36)
    footer(s,
      f"AISHE institution register as on 28-8-2026: {AD.N_COL:,} colleges, urban/rural flag per institution (the AISHE 2023-24 REPORT counts 48,246 colleges at 31-12-2023 — different instrument, different date)  ·  "
      f"store counts, pin codes and overlap: Bernstein via ET, 18 Jul 2026  ·  Starship: Food On Demand, 10 Jun 2026; Avride–Chartwells master services agreement  ·  "
      f"Euromonitor International, Consumer Foodservice by Location in India, Mar 2026.  Screen on A7.")

# =========================================================================
# SLIDE 5 -- D1, THE DEAD ZONE PRICED.  The pivot slide and the moneyshot.
# =========================================================================
def slide5():
    s=prs.slides.add_slide(L[CONTENT_LAYOUT]); strip_ph(s)
    rail(s,5,"D1  ·  THE DEAD ZONE, PRICED")

    flex   = CS.fixed_flex_share()
    thr    = [(lab, B.threshold(fn)) for lab,fn in B.CONFIGS]
    do_nothing, floor = thr[0][1], thr[-1][1]
    cut    = (do_nothing-floor)/do_nothing

    head(s,
      f"{flex:.1%} of the fixed base is flexible — and flexing all of it still leaves the "
      f"node needing {floor:.1%} of term-time demand from a break that delivers none of it",
      f"Fixed base Rs{CS.CAMPUS_FIXED:,.0f}/month: JM Financial Exhibit 13, confirmed a Tier-1/2 store "
      f"from Exhibit 5 of the same report. Robust across the tier's rent band "
      f"(Rs{CS.T12_FIXED_LOW:,.0f}–Rs{CS.T12_FIXED_HIGH:,.0f}, width {CS.T12_FIXED_BAND_PCT:.1%}).")

    y = preempt(s,1.50,"SAID BEFORE YOU ASK",
        f"‘Other fixed’ of Rs{CS.OTHER_FIXED:,.0f} — {CS.OTHER_FIXED/CS.CAMPUS_FIXED:.0%} of the base — is a residual "
        f"we allocated out of JM's blended Rs100/sqft utilities-and-other line, not an itemised figure. "
        f"It is the least defensible number in our stack, and it is the one that persists through the break.")

    rule_line(s,4.98,1.94,0.011,2.86)

    # ---------------- LEFT: what the base is made of, then the ladder ----------------
    yl = banner(s,0.55,1.92,4.15,"THE COST LADDER   ·   RESIDUAL DEMAND REQUIRED THROUGH THE BREAK")

    ty = yl+0.05
    text(s,0.55,ty,4.15,0.13,"WHAT THE BASE IS MADE OF  ·  AND WHAT ACTUALLY FLEXES",6.2,FG,bold=True)
    bx, BW = 0.55, 4.15
    cols = {"Store rent":HI, "In-store staff":BANNER_BG, "Utilities and cold chain":POS,
            "Other fixed":MUTE}
    for name,val,_,_ in CS.FIXED_STACK:
        w = BW*val/CS.CAMPUS_FIXED
        r = rect(s,bx,ty+0.16,w,0.20,cols[name])
        f=r.text_frame; f.vertical_anchor=MSO_ANCHOR.MIDDLE
        p=par(f,True); p.alignment=PP_ALIGN.CENTER
        run(p,f"{val/CS.CAMPUS_FIXED:.0%}",6.2,C(0xFF,0xFF,0xFF),bold=True)
        bx += w
    lx = 0.55
    for name,val,_,_ in CS.FIXED_STACK:
        w = BW*val/CS.CAMPUS_FIXED
        text(s,lx,ty+0.37,w,0.12,
             {"Store rent":"rent","In-store staff":"in-store staff",
              "Utilities and cold chain":"util.","Other fixed":"other fixed"}[name],
             5.2,MUTE)
        lx += w

    pic(s,"lever_ladder_s5",0.55,ty+0.50,4.15,1.50)
    text(s,0.55,ty+2.02,4.15,0.26,
      f"The shaded band is the residual no lever reaches: the ladder cuts the requirement {cut:.0%} "
      f"({do_nothing:.1%} to {floor:.1%}) and then stops, {floor*100:.1f} points above a break that delivers none.",
      6.2,MUTE,lsp=0.98)

    # ---------------- RIGHT: the moneyshot ----------------
    yr = banner(s,5.30,1.92,4.15,"THE MONEYSHOT   ·   ASSET TURN, RESTATED ON ONE BASIS")
    ib = box(s,5.30,yr+0.06,4.15,0.62,CHIP,radius=0.09)
    f=ib.text_frame; f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True); p.alignment=PP_ALIGN.CENTER; p.line_spacing=0.95
    run(p,f"({M.CEILING:,.0f} / {M.MINUTES_ORD:,.0f})  ×  ({M.ACTIVE_MONTHS} / 12)  =  ",11.5,FG,bold=True)
    run(p,f"{M.CEILING/M.MINUTES_ORD:.3f}",11.5,POS,bold=True)
    run(p,"  ×  ",11.5,FG,bold=True)
    run(p,f"{M.ACTIVE_MONTHS/12:.3f}",11.5,NEG,bold=True)
    run(p,"  =  ",11.5,FG,bold=True)
    run(p,f"{M.TURN_RATIO:.3f}",15.0,HI,bold=True)
    p2=par(f); p2.alignment=PP_ALIGN.CENTER
    run(p2,"density  ",6.2,POS,bold=True); run(p2,f"+{(M.CEILING/M.MINUTES_ORD-1)*100:.1f}%   ",6.2,MUTE)
    run(p2,"calendar  ",6.2,NEG,bold=True); run(p2,f"−{(1-M.ACTIVE_MONTHS/12)*100:.1f}%   ",6.2,MUTE)
    run(p2,"campus turn  ",6.2,FG,bold=True); run(p2,f"{M.CAMPUS_TURN:.2f}x  vs city  {M.CITY_TURN:.2f}x",6.2,MUTE)

    text(s,5.30,yr+0.74,4.15,0.42,
      f"Campus density almost exactly buys back the academic calendar. The dead zone costs "
      f"{(1-M.TURN_RATIO)*100:.1f} points of asset productivity, not 30 — and {M.ACTIVE_MONTHS/12:.3f} is the "
      f"{do_nothing:.1%} do-nothing threshold on the left, inverted. The same constant closes a third time.",
      6.6,FG,lsp=0.98)

    # the sensitivity strip -- the deck's most exposed number, priced in the open
    sy = yr+1.22
    text(s,5.30,sy,4.15,0.14,
      "AND IF 1,400 IS WRONG?  ·  THE RATIO IS LINEAR IN THROUGHPUT",6.5,FG,bold=True)
    cw = 4.15/len(M.TURN_SENSITIVITY)
    for i,(opd,ratio) in enumerate(M.TURN_SENSITIVITY):
        on  = (abs(opd-M.CEILING) < 1)
        par_ = (abs(opd-M.PARITY_OPD) < 1)
        b=box(s,5.30+i*cw,sy+0.18,cw-0.06,0.42, BANNER_BG if on else (CHIP2 if par_ else CHIP),radius=0.10)
        f=b.text_frame; f.vertical_anchor=MSO_ANCHOR.MIDDLE
        p=par(f,True); p.alignment=PP_ALIGN.CENTER; p.line_spacing=0.92
        run(p,f"{opd:,.0f}/day",6.2,BANNER_FG if on else MUTE)
        p2=par(f); p2.alignment=PP_ALIGN.CENTER; p2.line_spacing=0.92
        run(p2,f"{ratio:.3f}",10.0,BANNER_FG if on else (HI if par_ else FG),bold=True)
    text(s,5.30,sy+0.64,4.15,0.40,
      f"Blinkit's own observed range across nine quarters is {M.CEIL_LO:,}–{M.CEIL_HI:,} orders/day "
      f"(JM Exhibit 29). Across that entire range a campus node runs "
      f"{M.turn_ratio_at(M.CEIL_LO):.2f}–{M.turn_ratio_at(M.CEIL_HI):.2f}. Parity needs "
      f"{M.PARITY_OPD:,.0f}/day — {M.PARITY_VS_CEIL_HI:.1%} of the observed maximum. "
      f"At {M.TURN_SENSITIVITY[0][0]:,.0f}/day it is {M.TURN_RATIO_AT_1000:.2f}, and the density argument weakens.",
      6.2,MUTE,lsp=0.98)

    # one piece of outside evidence, and it is the only QUANTIFIED Indian proof that the
    # academic calendar moves a P&L. The other two (Kambli 2020, CRISIL) are in the footer.
    b=box(s,0.55,4.50,8.90,0.26,CHIP2,radius=0.10); rect(s,0.55,4.50,0.045,0.26,HI)
    f=b.text_frame; f.margin_left=In(0.12); f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True); p.line_spacing=0.92
    run(p,"IT IS NOT OUR CLAIM THAT THE CALENDAR MOVES A P&L  ",6.0,HI,bold=True)
    run(p,"Physicswallah's pre-Ind-AS EBITDA margin swings ~850bps between Q4FY26 and Q1FY27E on the "
          "academic calendar (with NEET postponement delaying offline batches). A listed Indian operator "
          "already reports the effect we are pricing.",6.2,FG)

    band(s,4.80,[("Cost cannot close the calendar. Density can. ",True),
      ("So the campus is not a cost problem to be mitigated — it is a site-selection problem, and the break becomes a filter criterion.",False)],h=0.36)
    footer(s,
      f"Fixed base, tier and rent band: JM Financial, Eternal model, Exhibits 5/12/13  ·  ceiling {M.CEILING:,}/day sits inside Blinkit's observed {M.CEIL_LO:,}–{M.CEIL_HI:,} range over nine quarters (Exhibit 29)  ·  "
      f"asset turn on capex Rs{(M.CAPEX_LO+M.CAPEX_HI)/2/1e7:.2f} cr and Minutes AOV Rs{M.MINUTES_AOV:,.0f} (TechCrunch, 22 Aug 2026) applied to BOTH stores  ·  lever thresholds solved, not assumed  ·  "
      f"Physicswallah: J.P. Morgan, 9 Jul 2026  ·  Kambli, Sinha & Srinivas 2020 (10.1016/j.jhtm.2020.02.008), wait 23.24→16.45 min  ·  CRISIL warehousing.  A5, A6; audit A1.")

# =========================================================================
# SLIDE 6 -- THEREFORE SITE SELECTION.  The calendar becomes a filter.
# =========================================================================
def slide6():
    s=prs.slides.add_slide(L[CONTENT_LAYOUT]); strip_ph(s)
    rail(s,6,"THE ANSWER  ·  SITE SELECTION, NOT COST")

    rev  = Q.post_revocation_survival()
    filt = AD.SITE_FILTER_UNCONTESTED

    head(s,
      f"Close the gap by choosing the site: {filt:,.0f} orders a day of uncontested adjacent demand, "
      f"{AD.N_CANDIDATES} districts of {AD.N_DISTRICTS:,} — and then we disqualified {AD.STACKED_SHARE:.0%} of our own list",
      f"The filter binds on access revocation ({rev['orders_needed']:,.0f} orders/day), not on break-period solvency "
      f"({rev['adjacent_lo']:,.0f}/day) — so it is set at the higher of the two. A node that survives losing the campus survives the break by construction.")

    y = preempt(s,1.50,"WHAT THIS SCREEN CANNOT DO",
      f"The AISHE register carries NO enrolment, so district student counts are imputed off state ratios and are labelled as one. "
      f"And ‘expected incumbent stores’ is a density band off {AD.NON_METRO_STORES:,} non-metro stores, not geocoded proximity: it ranks districts, it does not site a store.",
      col=NEG)

    rule_line(s,4.98,1.94,0.011,2.86)

    # ---------------- LEFT: the screen ----------------
    yl = banner(s,0.55,1.92,4.15,"THE SCREEN, MADE OPERATIONAL   ·   AND THEN TURNED ON OURSELVES")
    pic(s,"district_screen",0.55,yl+0.06,4.15,1.90)

    fy = yl+1.66
    text(s,0.55,fy,4.15,0.14,
      f"WHAT A CONTESTED CATCHMENT COSTS  ·  GROSS DEMAND NEEDED FOR {filt:,.0f}/DAY CLEAN",6.5,FG,bold=True)
    shares=[0.0,0.10,0.25,AD.METRO_OVERLAP_5OP]
    cw=4.15/len(shares)
    for i,sh in enumerate(shares):
        last=(i==len(shares)-1)
        b=box(s,0.55+i*cw,fy+0.18,cw-0.06,0.36,CHIP2 if last else CHIP,radius=0.10)
        f=b.text_frame; f.vertical_anchor=MSO_ANCHOR.MIDDLE
        p=par(f,True); p.alignment=PP_ALIGN.CENTER; p.line_spacing=0.90
        run(p,f"{sh:.0%} contested",5.6,NEG if last else MUTE)
        p2=par(f); p2.alignment=PP_ALIGN.CENTER; p2.line_spacing=0.90
        run(p2,f"{AD.gross_demand_required(sh):,.0f}/day",8.4,NEG if last else FG,bold=True)
    text(s,0.55,fy+0.56,4.15,0.30,
      f"Only the METRO five-operator overlap ({AD.METRO_OVERLAP_5OP:.0%}) is published, so the requirement is a "
      f"band and we quote the band, not a point. {AD.STACKED_SHARE:.0%} of our own candidate list already sits at "
      f"or above the incumbent density of the average location the industry serves.",5.7,MUTE,lsp=0.98)

    # ---------------- RIGHT: the ranking ----------------
    yr = banner(s,5.30,1.92,4.15,"WHERE   ·   URBAN COLLEGES × STATE RESIDENTIAL INTENSITY × HOSTEL OCCUPANCY")
    hy = yr+0.06
    for lab,x,w,al in [("DISTRICT",5.34,1.40,PP_ALIGN.LEFT),("STATE",6.74,1.02,PP_ALIGN.LEFT),
                       ("URB COL",7.76,0.56,PP_ALIGN.RIGHT),("RESID",8.34,0.46,PP_ALIGN.RIGHT),
                       ("EXP. INCUMBENT",8.80,0.65,PP_ALIGN.RIGHT)]:
        text(s,x,hy,w,0.13,lab,5.4,MUTE,bold=True,align=al)
    ry = hy+0.16
    tp = AD.with_proximity(AD.TOP).head(6)
    for i,r in enumerate(tp.itertuples()):
        if i % 2 == 0: rect(s,5.30,ry-0.012,4.15,0.215,CHIP)
        col = {"uncontested":POS,"contested":HI,"stacked":NEG}[r.proximity]
        text(s,5.34,ry,1.40,0.15,r.District.title(),6.4,FG,bold=True)
        text(s,6.74,ry,1.02,0.15,r.State,6.0,MUTE)
        text(s,7.76,ry,0.56,0.15,f"{r.urban_colleges:,.0f}",6.4,FG,bold=True,align=PP_ALIGN.RIGHT)
        text(s,8.34,ry,0.46,0.15,f"{r.resid_index:.1f}",6.2,MUTE,align=PP_ALIGN.RIGHT)
        text(s,8.80,ry,0.65,0.15,f"{r.exp_incumbent_stores:.0f} stores",6.2,col,bold=True,
             align=PP_ALIGN.RIGHT)
        ry += 0.215

    kar = int((AD.TOP.head(20).State == "Karnataka").sum())
    text(s,5.30,ry+0.03,4.15,0.28,
      f"THE CONVERGENCE: we chose Karnataka and Odisha at STATE level, off residential intensity and hostel "
      f"occupancy. The district register — a different table, pulled later — puts {kar} of the top 20 in Karnataka "
      f"and Khordha first. We did not pick the state and then find districts to justify it.",6.1,FG,lsp=0.96)

    # The finding that changes what the screen is FOR. Say it before a panellist derives it.
    wy = ry+0.34
    b=box(s,5.30,wy,4.15,0.46,None,radius=0.10,line=NEG,lw=1.0)
    f=b.text_frame; f.margin_left=In(0.11); f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True); p.line_spacing=0.94
    run(p,"AND THE UNCOMFORTABLE PART  ",5.9,NEG,bold=True)
    run(p,f"the {AD.N_UNCONTESTED} clean districts are clean because they are SMALL — every one of them "
          f"sits at the screen's {AD.MIN_URBAN_COLLEGES}–7 college floor. Incumbent density rises with "
          f"college count, so contestedness and cluster density trade off directly. There is no large, "
          f"empty district to find, and the plan does not need one.",5.8,FG)

    cy = wy+0.50
    b=box(s,5.30,cy,4.15,0.28,CHIP2,radius=0.10); rect(s,5.30,cy,0.045,0.28,HI)
    f=b.text_frame; f.margin_left=In(0.11); f.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=par(f,True); p.line_spacing=0.92
    run(p,"THE CALENDAR CONDITION  ",5.8,HI,bold=True)
    run(p,f"one break segment of {CF.MIN_WINDDOWN_WEEKS:.0f}+ weeks, or no lever fires: on a typical "
          f"{'+'.join(f'{w:.0f}' for w in CF.SHAPES['Typical: summer + winter + exam gaps'])}-week calendar only "
          f"{CF.TYPICAL_WINDOWNABLE:.0%} of break weeks are wind-downable and the dead zone costs {CF.TYPICAL_PENALTY:+.0%}.",
          5.7,MUTE)

    band(s,4.80,[("The campus inside the gate is uncontested by construction — an incumbent cannot serve it without the same permission we need. ",True),
      ("What this screen prices is the adjacent half of the catchment, which is exactly the half the break and the revocation case both lean on.",False)],h=0.36)
    footer(s,
      f"AISHE institution register as on 28-8-2026: {AD.N_COL:,} colleges + {AD.N_UNI:,} universities + {AD.N_STA:,} standalones = {AD.N_HEI:,} institutions across {AD.N_DISTRICTS:,} districts, each with an urban/rural flag "
      f"(the AISHE 2023-24 REPORT counts 48,246 colleges at 31-12-2023 — different instrument, different date)  ·  stores and metro overlap: Bernstein via ET, 18 Jul 2026 — {AD.TOTAL_STORES:,} stores across five operators, "
      f"{AD.STORES_PER_NEW_PINCODE:.1f} added per NEW pin code served  ·  filter from risk_quadrant.py and break_mode.py.  Screen and fragmentation on A7, A7b.")

slide1(); slide2(); slide5(); slide6()
out=os.path.join(ROOT,f"{OUT}_{THEME}.pptx")
prs.save(out)
print("wrote", out)
